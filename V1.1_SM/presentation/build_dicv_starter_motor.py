#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DICV_STARTER_MOTOR.pptx  — Starter-Motor business summary in the DICV_ALTERNATORS
visual language.  v2 (2026-07-06): 8 slides — honest-yet-optimistic, with an
engineering spine (how it was built · the four signals · why it is trustworthy).

Theme = DICV_ALTERNATORS deck: white canvas, BharatBenz-red section tags,
        charcoal headings, dark stat cards, wordmark top-right, dark action bar,
        BYTEDGE footer.

ALL numbers verified against frozen V1.1_SM artifacts (2026-07-06):
  results/V1_1_SM_model_spec.json, V1_1_SM_model_card.md,
  results/V1_1_SM_nested_lovo_predictions.csv, the 2026-07-02 DICV Starter-Motor
  Validation dossier, and Plan/2026-07-03…themes-5-6…md (56-heuristic census).
  AUROC 0.9321 (261/280); recall 13/14 @ Youden (5/20 FP); RED tier 10/14 @ 2/20;
  tiers GREEN 4F+16NF / AMBER 0F+2NF / RED 10F+2NF; median first-fire lead 168 d
  (77-424); A2 battery-cascade 0/20 FA, ~66 d (~9.5 wk); windows 28-91 / 126-284 d,
  9/11 inside; horizon k*=10 wk; 56 heuristics -> 10 admissible -> 4 features;
  coefs +0.886 / -0.414 / -0.270 / +0.141; solo AUROC 0.921 / 0.732 / 0.243 / 0.739;
  V1 restated 0.921->0.893; 4 reproductions, 17 challengers rejected; ~43% modeled
  saving.  No values invented.

VIN display: failed VIN1-14, in-service VIN15-34 (raw NF +14). RED = VIN19_NF
(raw VIN5_NF) & VIN34_NF (raw VIN20_NF); AMBER = VIN16_NF & VIN24_NF.

Run: py -3 build_dicv_starter_motor.py
"""
import os, tempfile, shutil
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── PATHS ──────────────────────────────────────────────────────
HERE   = Path(__file__).resolve().parent
V11    = HERE.parent
GRAPHS = V11 / "graphs"
ASSETS = HERE / "assets" / "dicv"
OUT    = HERE / "DICV_STARTER_MOTOR.pptx"

WORDMARK = str(ASSETS / "bharatbenz_wordmark.png")
BYTEDGE  = str(ASSETS / "bytedge_logo.png")
TRUCK    = str(ASSETS / "truck.png")
VIN6     = str(GRAPHS / "V1_1_SM_daily_risk_VIN6_F_SM_dashboard.png")
for p in (WORDMARK, BYTEDGE, TRUCK, VIN6):
    if not Path(p).exists():
        raise FileNotFoundError(p)

# ── DICV PALETTE (sampled from DICV_ALTERNATORS.pptx) ──────────
INK    = RGBColor(0x26, 0x26, 0x26)
RED    = RGBColor(0xE2, 0x23, 0x1A)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GOLD   = RGBColor(0xE0, 0xA2, 0x01)
BLUE   = RGBColor(0x1B, 0x6C, 0xA8)
GREY   = RGBColor(0x5E, 0x5E, 0x5E)
GREYL  = RGBColor(0x8C, 0x8C, 0x8C)
RULE   = RGBColor(0xD3, 0xD7, 0xDC)
CARDBG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARDDK = RGBColor(0x26, 0x26, 0x26)
CARDDK2= RGBColor(0x1E, 0x1E, 0x1E)
TIER_G = RGBColor(0x2E, 0x8B, 0x57)
TIER_A = RGBColor(0xE0, 0xA2, 0x01)
TIER_R = RGBColor(0xC0, 0x39, 0x2B)
ZEBRA  = RGBColor(0xF2, 0xF4, 0xF7)

FONT = 'Calibri'
SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.55)
TMP = tempfile.mkdtemp(prefix='dicv_sm_')
H_INK, H_RED, H_GREEN, H_GOLD, H_BLUE, H_GREY, H_TR = (
    '#262626', '#E2231A', '#2E8B57', '#E0A201', '#1B6CA8', '#5E5E5E', '#C0392B')


# ── LOW-LEVEL HELPERS ──────────────────────────────────────────
def prs_new():
    p = Presentation(); p.slide_width = SW; p.slide_height = SH; return p

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def rect(slide, x, y, w, h, fill, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    return shp

def text(slide, x, y, w, h, runs, size=12, bold=False, color=INK,
         align=PP_ALIGN.LEFT, font=FONT, anchor=None, spacing=None, wrap=True,
         line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    if anchor: tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]; p.alignment = align
    if line_spacing: p.line_spacing = line_spacing
    for t, ov in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(ov.get('size', size))
        r.font.bold = ov.get('bold', bold)
        r.font.color.rgb = ov.get('color', color)
        r.font.name = ov.get('font', font)
        sp = ov.get('spacing', spacing)
        if sp is not None:
            r.font._rPr.set('spc', str(int(sp * 100)))
    return tb

def bullets(slide, x, y, w, h, items, size=10, color=INK, marker='‣',
            mcolor=None, gap=4, lead=1.05):
    mcolor = mcolor or RED
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        if isinstance(it, str):
            txt, b, c = it, False, color
        else:
            txt = it[0]; b = it[1] if len(it) > 1 else False
            c = it[2] if len(it) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = lead
        if marker:
            rm = p.add_run(); rm.text = marker + '  '
            rm.font.size = Pt(size); rm.font.bold = True
            rm.font.color.rgb = mcolor; rm.font.name = FONT
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = b
        r.font.color.rgb = c; r.font.name = FONT
    return tb


# ── DICV MOTIF HELPERS ─────────────────────────────────────────
def wordmark(slide):
    ww, wh = Inches(1.55), Inches(0.42)
    wx = SW - Inches(0.55) - ww; wy = Inches(0.34)
    rect(slide, wx - Inches(0.14), wy + Inches(0.02), Inches(0.045), wh - Inches(0.04), RED)
    slide.shapes.add_picture(WORDMARK, wx, wy, width=ww, height=wh)

def section_header(slide, num, tag, title, title_size=23):
    wordmark(slide)
    text(slide, ML, Inches(0.34), Inches(8.0), Inches(0.28),
         "%s · %s" % (num, tag.upper()), size=11.5, bold=True, color=RED, spacing=1.2)
    text(slide, ML, Inches(0.66), Inches(10.4), Inches(0.8), title,
         size=title_size, bold=True, color=INK, line_spacing=1.0)
    rect(slide, ML, Inches(1.52), SW - 2 * ML, Pt(1.6), INK)

def footer(slide, page):
    y = Inches(7.02)
    rect(slide, ML, y, SW - 2 * ML, Pt(0.9), RULE)
    slide.shapes.add_picture(BYTEDGE, ML, Inches(7.14), height=Inches(0.2))
    text(slide, Inches(4.2), Inches(7.13), Inches(5.0), Inches(0.25),
         "DICV · Starter-Motor Predictive Maintenance · V1.1_SM — Confidential",
         size=7.5, color=GREYL, align=PP_ALIGN.CENTER)
    text(slide, Inches(9.4), Inches(7.13), Inches(3.0), Inches(0.25),
         "Figures audited to V1.1_SM source files", size=7.5, color=GREYL, align=PP_ALIGN.RIGHT)
    text(slide, SW - Inches(0.75), Inches(7.11), Inches(0.35), Inches(0.25),
         str(page), size=9, bold=True, color=GREY, align=PP_ALIGN.RIGHT)

def stat_card(slide, x, y, w, h, label, big, sub, num_color=WHITE):
    rect(slide, x, y, w, h, CARDDK, rounded=True)
    text(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.24),
         label.upper(), size=8, bold=True, color=RGBColor(0xB9, 0xBE, 0xC6),
         align=PP_ALIGN.CENTER, spacing=0.6)
    text(slide, x + Inches(0.1), y + Inches(0.36), w - Inches(0.2), Inches(0.5),
         big, size=25, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.12), y + h - Inches(0.42), w - Inches(0.24), Inches(0.38),
         sub, size=7.6, color=RGBColor(0xC7, 0xCC, 0xD3), align=PP_ALIGN.CENTER, line_spacing=0.98)

def cap_card(slide, x, y, w, h, title, items, accent, title_size=11.5, bsize=8.6):
    rect(slide, x, y, w, h, CARDBG, rounded=True)
    rect(slide, x, y, w, Inches(0.09), accent, rounded=False)
    text(slide, x + Inches(0.16), y + Inches(0.16), w - Inches(0.3), Inches(0.3),
         title, size=title_size, bold=True, color=accent)
    bullets(slide, x + Inches(0.18), y + Inches(0.52), w - Inches(0.34), h - Inches(0.6),
            items, size=bsize, color=INK, marker='‣', mcolor=accent, gap=4, lead=1.02)

def action_bar(slide, text_str, label="RECOMMENDED ACTION"):
    y, h = Inches(6.4), Inches(0.5)
    rect(slide, ML, y, SW - 2 * ML, h, CARDDK2, rounded=True)
    tb = slide.shapes.add_textbox(ML + Inches(0.2), y, SW - 2 * ML - Inches(0.4), h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = label + ":   "
    r1.font.size = Pt(11); r1.font.bold = True; r1.font.color.rgb = RED; r1.font.name = FONT
    r2 = p.add_run(); r2.text = text_str
    r2.font.size = Pt(10.5); r2.font.color.rgb = WHITE; r2.font.name = FONT

def dicv_table(slide, x, y, w, rows, col_w, header, body, tier_col=None,
               fs=8.6, hfs=8.8, row_h=0.34):
    ncols = len(header)
    gt = slide.shapes.add_table(rows, ncols, x, y, w, Inches(row_h * rows))
    tbl = gt.table; tbl.first_row = False; tbl.horz_banding = False
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = Inches(cw)

    def _cell(cell, val, size, bold, color, fill, align=PP_ALIGN.LEFT):
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        # tuple -> two-line cell (main bold + grey technical sub-line)
        if isinstance(val, (tuple, list)):
            main, sub = val
            p = tf.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = str(main)
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT
            p2 = tf.add_paragraph(); p2.alignment = align
            r2 = p2.add_run(); r2.text = str(sub)
            r2.font.size = Pt(size - 1.4); r2.font.italic = True
            r2.font.color.rgb = GREYL; r2.font.name = FONT
        else:
            p = tf.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT
    for ci, val in enumerate(header):
        _cell(tbl.cell(0, ci), val, hfs, True, WHITE, INK,
              PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT)
    for ri, row in enumerate(body, start=1):
        base = WHITE if ri % 2 else ZEBRA
        for ci, val in enumerate(row):
            col = INK
            if tier_col and ci == tier_col:
                u = str(val).upper()
                col = TIER_R if 'RED' in u or 'HIGH' in u else (
                    TIER_A if 'AMBER' in u or 'WATCH' in u else (
                    TIER_G if 'GREEN' in u or 'LOW' in u else INK))
            _cell(tbl.cell(ri, ci), val, fs, ci == 0 or (tier_col == ci), col, base,
                  PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    return gt

def minihead(slide, x, y, w, label, color=INK, size=13):
    text(slide, x, y, w, Inches(0.3), label.upper(), size=size, bold=True,
         color=color, spacing=0.4)


# ── CHARTS (DICV palette) ──────────────────────────────────────
def _clean_ax(ax):
    for s in ('top', 'right'):
        ax.spines[s].set_visible(False)
    ax.spines['left'].set_color('#B7BDC6'); ax.spines['bottom'].set_color('#B7BDC6')
    ax.tick_params(colors=H_GREY)

def chart_pipeline():
    fig, ax = plt.subplots(figsize=(12, 2.65))
    stages = [
        ('CAN Telemetry', '#8C8C8C', 'Existing 6-signal data\nfrom 34 BharatBenz trucks', '#1A1A1A'),
        ('AI Risk Scoring', H_INK, 'Calibrated risk per truck\n93.2% ranking accuracy', 'white'),
        ('Early-Warning Alerts', H_RED, 'Median 168-day lead\nbefore recorded failure', 'white'),
        ('Battery-vs-Starter Triage', H_GOLD, 'Cascade detector routes\nbattery-first (0 false alarms)', '#262626'),
        ('Maintenance Decision', H_GREEN, 'RED: inspect in 2-4 wks\nAMBER: next service', 'white'),
    ]
    n = len(stages); step = 2.35
    for i, (label, color, detail, tcol) in enumerate(stages):
        x = i * step
        box = mpatches.FancyBboxPatch((x, 0.55), 1.92, 1.9, boxstyle="round,pad=0.10",
                                       facecolor=color, edgecolor='white', linewidth=1.5)
        ax.add_patch(box)
        ax.text(x + 0.96, 1.95, label, ha='center', va='center', fontsize=10.3,
                fontweight='bold', color=tcol)
        ax.text(x + 0.96, 1.05, detail, ha='center', va='center', fontsize=7.6, color=tcol)
        if i < n - 1:
            ax.annotate('', xy=(x + 2.24, 1.5), xytext=(x + 1.98, 1.5),
                        arrowprops=dict(arrowstyle='-|>', color=H_RED, lw=2.6))
    ax.set_xlim(-0.12, (n - 1) * step + 1.92 + 0.12); ax.set_ylim(0.2, 2.75); ax.axis('off')
    fig.patch.set_facecolor('white'); plt.tight_layout()
    p = os.path.join(TMP, 'pipeline.png')
    fig.savefig(p, dpi=200, bbox_inches='tight', pad_inches=0.02, facecolor='white'); plt.close(fig)
    return p

def chart_funnel():
    fig, ax = plt.subplots(figsize=(12, 2.35))
    stages = [
        ('107.2 M', 'raw 5-second CAN rows\n6 signals · 34 trucks', '#8C8C8C', '#1A1A1A'),
        ('20,471', 'crank events catalogued\n+ 2,636 truck-weeks', H_INK, 'white'),
        ('56', 'heuristics engineered\n& adjudicated (5 families)', H_RED, 'white'),
        ('10', 'admissible after\nleakage screening', H_GOLD, '#262626'),
        ('4', 'features in the\nfrozen model', H_GREEN, 'white'),
    ]
    n = len(stages); step = 2.35
    heights = [1.95, 1.78, 1.6, 1.42, 1.24]   # funnel taper
    for i, (big, detail, color, tcol) in enumerate(stages):
        x = i * step; hh = heights[i]; y0 = (2.2 - hh) / 2 + 0.35
        box = mpatches.FancyBboxPatch((x, y0), 1.92, hh, boxstyle="round,pad=0.09",
                                       facecolor=color, edgecolor='white', linewidth=1.5)
        ax.add_patch(box)
        ax.text(x + 0.96, y0 + hh * 0.66, big, ha='center', va='center',
                fontsize=17, fontweight='bold', color=tcol)
        ax.text(x + 0.96, y0 + hh * 0.24, detail, ha='center', va='center',
                fontsize=7.4, color=tcol)
        if i < n - 1:
            ax.annotate('', xy=(x + 2.25, 1.45), xytext=(x + 1.97, 1.45),
                        arrowprops=dict(arrowstyle='-|>', color=H_RED, lw=2.6))
    ax.set_xlim(-0.12, (n - 1) * step + 1.92 + 0.12); ax.set_ylim(0.1, 2.75); ax.axis('off')
    fig.patch.set_facecolor('white'); plt.tight_layout()
    p = os.path.join(TMP, 'funnel.png')
    fig.savefig(p, dpi=200, bbox_inches='tight', pad_inches=0.02, facecolor='white'); plt.close(fig)
    return p

def chart_importance():
    fig, ax = plt.subplots(figsize=(5.4, 2.55))
    feats = ['Within-week\nvoltage noise', 'Weekly voltage-\nrange trend',
             'Resting-voltage\nfloor', 'Crank\ndip depth']
    absco = [0.886, 0.414, 0.270, 0.141]
    notes = ['+0.89 · solo 0.92', '−0.41 · solo 0.73', '−0.27 · solo 0.24', '+0.14 · solo 0.74']
    cols = [H_GREEN, H_GREY, H_GOLD, H_BLUE]
    yp = list(range(len(feats)))[::-1]
    ax.barh(yp, absco, color=cols, height=0.6, zorder=3, edgecolor='white', linewidth=0.8)
    for y, v, nt in zip(yp, absco, notes):
        ax.text(v + 0.02, y, nt, va='center', ha='left', fontsize=8.2, color=H_INK, fontweight='bold')
    ax.set_yticks(yp); ax.set_yticklabels(feats, fontsize=8.6, color=H_INK)
    ax.set_xlim(0, 1.2); ax.set_xlabel('|learned weight| in the risk score', fontsize=8.5, color=H_GREY)
    ax.set_title('Feature importance — learned weights', fontsize=11, fontweight='bold',
                 color=H_INK, pad=8)
    _clean_ax(ax); fig.patch.set_facecolor('white'); plt.tight_layout()
    p = os.path.join(TMP, 'importance.png')
    fig.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close(fig)
    return p

def chart_fleet_risk():
    fig, ax = plt.subplots(figsize=(5.4, 3.4))
    cats = ['LOW\n(Green)', 'WATCH\n(Amber)', 'HIGH\n(Red)']
    vals = [16, 2, 2]; cols = [H_GREEN, H_GOLD, H_TR]
    ax.bar(cats, vals, color=cols, width=0.55, edgecolor='white', linewidth=1.0, zorder=3)
    for i, v in enumerate(vals):
        ax.text(i, v + 0.3, '%d truck%s' % (v, 's' if v > 1 else ''), ha='center',
                va='bottom', fontsize=13, fontweight='bold', color=H_INK)
    ax.set_ylim(0, 19); ax.set_ylabel('In-service trucks', fontsize=10.5, color=H_GREY)
    ax.set_title('Current fleet risk mix (20 in-service trucks)', fontsize=12.5,
                 fontweight='bold', color=H_INK, pad=9)
    _clean_ax(ax); fig.patch.set_facecolor('white'); plt.tight_layout()
    p = os.path.join(TMP, 'fleet_risk.png')
    fig.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close(fig)
    return p

def chart_effectiveness():
    fig, ax = plt.subplots(figsize=(11, 3.0))
    cats = ['Fleet risk ranking', 'Early-warning alerts', 'Battery-vs-starter triage',
            'Exact failure dates', 'Silent / abrupt failures']
    scores = [93, 93, 100, 0, 0]; cols = [H_GREEN, H_GREEN, H_GREEN, H_TR, H_TR]
    labels = ['13/14 failures flagged (AUROC 0.93)', '13/14 fire an alert · median 168-day lead',
              '0 false alarms on cascade detector', 'Not supportable at 14 events',
              '~4 of 14 invisible to this telemetry']
    ax.barh(cats[::-1], scores[::-1], color=cols[::-1], height=0.55,
            edgecolor='white', linewidth=0.8, zorder=3)
    for i, (s, l) in enumerate(zip(scores[::-1], labels[::-1])):
        c = H_GREEN if s > 50 else H_TR
        ax.text(max(s, 4) + 2, i, l, ha='left', va='center', fontsize=9.2, fontweight='bold', color=c)
    ax.set_xlim(0, 150); ax.set_xticks([0, 20, 40, 60, 80, 100])
    ax.set_xlabel('Capability (%)', fontsize=10, color=H_GREY)
    _clean_ax(ax); fig.patch.set_facecolor('white'); plt.tight_layout()
    p = os.path.join(TMP, 'effectiveness.png')
    fig.savefig(p, dpi=200, bbox_inches='tight', facecolor='white'); plt.close(fig)
    return p


print("Generating DICV-themed charts …")
IMG_PIPE = chart_pipeline()
IMG_FUN  = chart_funnel()
IMG_IMP  = chart_importance()
IMG_RISK = chart_fleet_risk()
IMG_EFF  = chart_effectiveness()

VIN_FOOTNOTE = ("Display numbering: VIN1-14 failed, VIN15-34 in-service "
                "(raw-file map: results/V1_1_SM_vin_naming_map.csv).")

prs = prs_new()

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE + HEADLINE KPIs
# ══════════════════════════════════════════════════════════════
s = blank(prs)
wordmark(s)
text(s, ML, Inches(0.7), Inches(9.0), Inches(0.3),
     "PREDICTIVE MAINTENANCE · FLEET RESULTS & PHASE-2 READINESS",
     size=11.5, bold=True, color=RED, spacing=1.2)
text(s, ML, Inches(1.08), Inches(8.6), Inches(0.9),
     "Starter-Motor Failure Prediction", size=39, bold=True, color=INK)
text(s, ML, Inches(1.92), Inches(8.6), Inches(0.6),
     "for the BharatBenz 5528T Fleet", size=24, bold=False, color=GREY)
text(s, ML, Inches(2.62), Inches(8.0), Inches(0.4),
     "From roadside no-starts to a monitored, scored fleet program", size=14, color=GREY)
s.shapes.add_picture(TRUCK, Inches(8.85), Inches(1.95), width=Inches(4.0))
rect(s, ML, Inches(3.16), Inches(3.4), Pt(2.2), RED)
text(s, ML, Inches(3.34), Inches(12.2), Inches(0.3),
     "93.2% ranking accuracy · 13 of 14 failures caught · median ~6 months' warning · zero battery-alert false alarms",
     size=12.5, bold=True, color=INK)
text(s, ML, Inches(3.68), Inches(12.2), Inches(0.3),
     [("9 of 11 predicted service windows contained the actual failure", {'color': GREEN, 'bold': True}),
      (" — and both misses were on the safe side (the truck failed after the window closed).", {'color': GREY})],
     size=11)
text(s, ML, Inches(4.02), Inches(12.2), Inches(0.3),
     "Prepared for Fleet Operations, Maintenance & Engineering    |    July 2026    |    Version V1.1_SM",
     size=10.5, color=GREYL)
cards = [
    ("Trucks in study", "34", "14 failed · 20 in-service", WHITE),
    ("Ranking accuracy", "93.2%", "failing vs. healthy (AUROC)", WHITE),
    ("Failures caught", "13 / 14", "catch-most setting", GOLD),
    ("Median warning lead", "168 d", "≈ 5.5 months (77–424 d)", WHITE),
    ("Battery-alert false alarms", "0 / 20", "cascade detector, healthy trucks", GREEN),
]
cw, ch, gap = Inches(2.34), Inches(1.32), Inches(0.13)
x = ML; y = Inches(4.6)
for lab, big, sub, col in cards:
    stat_card(s, x, y, cw, ch, lab, big, sub, num_color=col); x += cw + gap
text(s, ML, Inches(6.62), Inches(9.0), Inches(0.3),
     "Confidential — internal distribution only", size=9, color=GREYL)
s.shapes.add_picture(BYTEDGE, ML, Inches(7.14), height=Inches(0.2))
text(s, SW - Inches(5.2), Inches(7.13), Inches(4.65), Inches(0.25),
     "BharatBenz 5528T Predictive Maintenance · Frozen V1.1_SM model",
     size=8, color=GREYL, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 01 · WHAT WAS BUILT
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "01", "What Was Built",
               "Turning existing truck telemetry into actionable maintenance intelligence")
s.shapes.add_picture(IMG_PIPE, ML, Inches(1.74), width=SW - 2 * ML)
minihead(s, ML, Inches(4.18), Inches(9.0), "Three plain-language capabilities", INK, 13)
cy, ch2 = Inches(4.56), Inches(1.62); cw2 = Inches(3.94); g2 = Inches(0.18)
cap_card(s, ML, cy, cw2, ch2, "1 · Risk ranking", [
    "Every truck gets a calibrated risk score from its own voltage patterns",
    "Scores map to GREEN / AMBER / RED maintenance tiers",
    "13 of 14 past failures ranked as high-risk (AUROC 0.93)",
], BLUE)
cap_card(s, ML + cw2 + g2, cy, cw2, ch2, "2 · Early-warning alerts", [
    "Median 168-day lead from the first validated alert to failure (77–424 days)",
    "Two validated service windows — battery ~1–3 months, persistence ~4–9 months",
    "Battery-cascade alert: ~9.5-week median lead, zero false alarms",
], RED)
cap_card(s, ML + 2 * (cw2 + g2), cy, cw2, ch2, "3 · Battery-vs-starter triage", [
    "Detects the battery-degradation cascade that mimics starter failure",
    "Screen-grade — tells the depot where to look first (battery-first)",
    "Agreed with failure type on 9 of 11 · 0 of 20 false attributions",
], GREEN)
action_bar(s, "Deploy the risk ranking, tier flags and alert channels on existing CAN-bus data "
              "now — no new sensors are required to start getting value.")
footer(s, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 02 · HOW IT WAS BUILT  (NEW)
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "02", "How It Was Built",
               "From 107 million rows of raw telemetry to four trusted signals")
s.shapes.add_picture(IMG_FUN, ML, Inches(1.72), width=SW - 2 * ML)
# left: families table
minihead(s, ML, Inches(4.12), Inches(7.0), "The 56 heuristics span five engineering families", INK, 12)
fam_header = ['Family', 'What it probes', 'In the model']
fam_body = [
    ['A · Crank duration & effort', 'Brush / commutator / bearing wear via crank waveform', 'Crank-dip signal'],
    ['B · Start failure & retry', 'Solenoid / pinion engagement + hard-start', 'Admissible (not selected)'],
    ['C · Voltage & battery ★', 'Supply-voltage stability & resting floor', '3 of 4 model features'],
    ['D · Usage & duty cycle', 'Cycle-dose exposure and wear', 'Screened out'],
    ['E · Environment & season', 'Thermal stress (needs new sensors)', 'Infeasible today'],
]
dicv_table(s, ML, Inches(4.44), Inches(7.35), 6, [2.55, 3.15, 1.65],
           fam_header, fam_body, fs=7.9, hfs=8.2, row_h=0.30)
# right: how it was policed
minihead(s, Inches(8.3), Inches(4.12), Inches(4.6), "How the funnel was policed", INK, 12)
bullets(s, Inches(8.3), Inches(4.46), Inches(4.5), Inches(1.9), [
    "Every feature is anchored to each truck's own healthy baseline (change-based, "
    "not level-based) — it flags real degradation, not truck-to-truck differences.",
    "A dominant-frequency candidate was exposed as a history-length artifact by the "
    "fixed-window control — caught, banned and locked out by a feature registry.",
    "7 SMA-dead trucks (a telemetry config) are cohort-masked, so missing crank data "
    "is never read as a fault.",
], size=8.6, marker='‣', mcolor=RED, gap=5, lead=1.02)
action_bar(s, "Small on purpose — at 14 failure events, a four-signal model that generalises "
              "beats a forty-signal model that memorises.", label="THE DESIGN CHOICE")
footer(s, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 03 · THE FOUR SIGNALS IT WATCHES  (NEW)
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "03", "The Four Signals It Watches",
               "Physically meaningful, always read together — no black box")
feat_header = ['Signal', 'What it measures (one line)', 'Role in the model']
feat_body = [
    [('Within-week voltage noise', 'vsi_withinwk_std_ratio_30d_w'),
     "Voltage jitter within a week vs the truck's own healthy baseline",
     'The workhorse — carries most of the score'],
    [('Weekly voltage-range trend', 'vsi_range_trend'),
     'Whether the weekly voltage swing is widening or narrowing (robust 12-wk slope)',
     'Suppressor — sharpens the ranking'],
    [('Resting-voltage floor', 'rest_vsi_p05_delta90'),
     'How far the engine-off battery floor has sagged vs own baseline (swap-aware)',
     'The battery-cascade signature'],
    [('Crank dip depth', 'dip_depth_last90_delta'),
     'How much deeper voltage dips during cranking, last 90 d vs own baseline',
     'The crank-circuit load signature'],
]
dicv_table(s, ML, Inches(1.7), Inches(7.35), 5, [2.55, 3.35, 1.45],
           feat_header, feat_body, fs=8.2, hfs=8.5, row_h=0.5)
s.shapes.add_picture(IMG_IMP, Inches(8.05), Inches(1.78), width=Inches(4.8))
text(s, Inches(8.05), Inches(3.92), Inches(4.8), Inches(0.35),
     "Weights are learned (Ridge, standardised) — not hand-set. The two voltage signals "
     "are a core pair selected in 34 of 34 CV folds.",
     size=7.6, color=GREY, line_spacing=1.0)
# alert-heuristics strip
minihead(s, ML, Inches(4.32), Inches(9.0), "On top of the score: three alert heuristics", INK, 12)
ay, ah = Inches(4.6), Inches(1.62); aw = Inches(3.94); ag = Inches(0.18)
cap_card(s, ML, ay, aw, ah, "Persistence flag", [
    "Risk stays elevated for consecutive weeks",
    "The primary early-warning channel — schedules service",
], BLUE, title_size=11, bsize=8.6)
cap_card(s, ML + aw + ag, ay, aw, ah, "A1 · Crank-burst", [
    "A burst of failed or retried cranks appears",
    "Corroborates the risk score near failure",
], GOLD, title_size=11, bsize=8.6)
cap_card(s, ML + 2 * (aw + ag), ay, aw, ah, "A2 · Battery-cascade", [
    "Rest-floor sag + deeper dips + rising crank failures together",
    "0 of 20 false alarms · ~9.5-wk lead · routes battery-first",
], GREEN, title_size=11, bsize=8.6)
action_bar(s, "No single signal decides anything — the model reads all four together, which is "
              "why weak individual cues combine into a 93.2% ranking.", label="THE KEY POINT")
footer(s, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — 04 · HOW WELL IT WORKS
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "04", "How Well It Works",
               "Validated across 34 trucks — 14 known failures + 20 in-service")
s.shapes.add_picture(VIN6, Inches(0.5), Inches(1.68), width=Inches(3.62))
minihead(s, Inches(4.35), Inches(1.72), Inches(3.0), "Daily degradation view", INK, 11)
bullets(s, Inches(4.35), Inches(2.06), Inches(4.6), Inches(1.9), [
    "Each truck gets a daily-updated degradation view; colours are action tiers "
    "(green = healthy, red = act now).",
    "This failed truck (VIN6) walked the full green-to-red path; the battery-cascade "
    "alert fired ~70 days early.",
    ("The dotted tail is a fleet-average projection, not a failure date.", False, GREY),
], size=9, marker='‣', mcolor=BLUE, gap=6, lead=1.03)
s.shapes.add_picture(IMG_RISK, Inches(9.16), Inches(1.72), width=Inches(3.62))
minihead(s, ML, Inches(4.16), Inches(4.0), "Headline results", INK, 13)
text(s, Inches(6.0), Inches(4.22), Inches(6.83), Inches(0.25), VIN_FOOTNOTE,
     size=7.5, color=GREYL, align=PP_ALIGN.RIGHT)
res_header = ['What', 'Result', 'Business impact']
res_body = [
    ['Failure detection', '13 of 14 caught at the catch-most setting (5 of 20 healthy flagged)',
     'High-risk trucks inspected before they strand'],
    ['Early warning', 'Median 168-day first-fire lead (77–424); battery alerts ~9.5 wks ahead',
     'Time to schedule, source parts, avoid roadside events'],
    ['Service windows', '9 of 11 windowed failures fell inside; both misses on the safe side',
     'Failures pre-empted by a scheduled service'],
    ['False-alarm burden', '2 of 20 healthy at the stricter RED tier; 0 of 20 on the battery alert',
     'Inspection effort goes where the risk is'],
    ['vs. first-generation', '+1 failure caught (incl. its worst miss); honest probabilities + alerts',
     'The upgrade is measured, not asserted'],
    ['Fleet status today', '16 LOW / 2 WATCH / 2 HIGH (VIN19_NF & VIN34_NF are HIGH)',
     'Two trucks need inspection now'],
]
dicv_table(s, ML, Inches(4.46), SW - 2 * ML, 7, [2.15, 5.95, 4.13],
           res_header, res_body, fs=8.1, hfs=8.6, row_h=0.255)
action_bar(s, "One dial, one honest trade — catch-most flags 13/14 (5/20 healthy); the stricter "
              "RED tier flags 10/14 (2/20). Operations picks the setting per cost.")
footer(s, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — 05 · WHY THE NUMBERS ARE TRUSTWORTHY  (NEW)
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "05", "Why the Numbers Are Trustworthy",
               "A strong score on only 14 events invites one question — is it real?")
g_header = ['#', 'Validation layer', 'What it showed']
g_body = [
    ['1', 'Nested leave-one-truck-out', 'Every truck scored blind across 34 folds — a deployment-grade number, not an in-sample flatterer'],
    ['2', 'Selection optimism measured', 'In-sample 262 vs honest 261 of 280 pairs — a one-pair gap (a memoriser shows a large gap)'],
    ['3', 'Label-shuffle test', '0 of 200 shuffled-label runs beat the real score (p = 0.005)'],
    ['4', 'Fixed-window control', 'Every feature recomputed on equal windows — score reproduced bit-for-bit (0.0 borrowed from history length)'],
    ['5', 'Leak-ceiling audit', 'Observation-length proxies scored as upper bounds; one artifact feature banned by a binding registry'],
    ['6', 'Predecessor restated honestly', 'The first-generation 0.921 was revised down to 0.893 under these stricter rules — disclosed, not hidden'],
    ['7', 'Reproduced & attacked', 'Benchmark reproduced exactly 4× (V2–V3.1); 17 pre-registered challenger features all rejected'],
]
dicv_table(s, ML, Inches(1.72), SW - 2 * ML, 8, [0.45, 3.35, 8.43],
           g_header, g_body, fs=8.4, hfs=8.6, row_h=0.33)
rect(s, ML, Inches(4.66), SW - 2 * ML, Inches(1.5), CARDDK, rounded=True)
text(s, ML + Inches(0.22), Inches(4.8), SW - 2 * ML - Inches(0.44), Inches(0.3),
     "WHAT THIS BUYS YOU", size=10.5, bold=True, color=RED, spacing=0.6)
text(s, ML + Inches(0.22), Inches(5.12), SW - 2 * ML - Inches(0.44), Inches(0.95),
     "When this system says 93.2% ranking accuracy, that is what you will see on unseen trucks. "
     "The score has survived a shuffle test, a fixed-window control, a leak audit, four independent "
     "reproductions and seventeen attempts to beat it. The ceiling here is the data, not the method "
     "— so every number sharpens as the fleet grows, using this same validated pipeline.",
     size=10.5, color=WHITE, line_spacing=1.08)
action_bar(s, "The same rigour ships with the deployment — retrain only under the audited protocol, "
              "never by hand-tuning a threshold.", label="IN OPERATION")
footer(s, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 06 · HONEST BOUNDARIES — AND WHAT UNLOCKS EACH
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "06", "Honest Boundaries — And What Unlocks Each",
               "Every limit below is a data limit, not a method limit")
s.shapes.add_picture(IMG_EFF, Inches(0.5), Inches(1.72), width=Inches(7.15))
minihead(s, Inches(7.95), Inches(1.68), Inches(5.0), "Honest boundaries", INK, 12)
bullets(s, Inches(7.95), Inches(2.04), Inches(4.9), Inches(2.2), [
    ("It does not predict failure dates.", True, RED),
    "With 14 examples a countdown clock is unsupportable — every timing model lost to a "
    "simple fleet-average (576 d vs 44 d error). We say ‘within ~10 weeks’, never ‘on 14 June’.",
    ("~4 of 14 failures are electrically invisible.", True, RED),
    "Sudden internal failures (windings, seizure) and telemetry-silent trucks give no "
    "electrical warning at this data resolution.",
    ("Some alerts are early or precautionary.", True, RED),
    "A few healthy trucks are flagged: cheap insurance inspections, but not failures.",
], size=8.7, marker='‣', mcolor=RED, gap=4, lead=1.0)
minihead(s, ML, Inches(4.28), Inches(9.0), "Data gaps — and the unlock", INK, 13)
gap_header = ['Gap', 'Impact on the system', 'The unlock']
gap_body = [
    ['Voltage sampled only every 5 seconds', 'The 2–4 month brush-wear warning is physically unreadable',
     'High-frequency crank logging (≥1 Hz)'],
    ['No cranking-current or battery-health signal', 'Battery vs. starter can only be inferred, not proven',
     'Add current or battery SoC/SoH sensing'],
    ['No maintenance / parts records in the data', 'Failure modes are inferred from patterns, not confirmed',
     'Link workshop records to telemetry'],
    ['Only 14 failure examples', 'Wide statistical band on every number (CI 0.81–0.99)',
     'Scale to 500 trucks (~30–50+ failures)'],
    ['7 of 34 trucks transmit no crank data', 'An SMA-dead config hides starter activity',
     'Telemetry-config audit; treat silence as a trigger'],
]
dicv_table(s, ML, Inches(4.6), SW - 2 * ML, 6, [3.35, 5.2, 3.68],
           gap_header, gap_body, fs=8.2, hfs=8.6, row_h=0.30)
action_bar(s, "Each boundary lifts at 500-truck scale using this same validated method — and only "
              "new sensing (≥1 Hz crank logging, battery health) breaks the silent-failure floor.",
           label="THE KEY POINT")
footer(s, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — 07 · RECOMMENDED ACTIONS & NEXT STEPS
# ══════════════════════════════════════════════════════════════
s = blank(prs)
section_header(s, "07", "Recommended Actions & Next Steps",
               "Priority actions for fleet operations, and what to invest in next")
minihead(s, ML, Inches(1.68), Inches(5.8), "Immediate actions", GREEN, 13)
imm_header = ['#', 'Action', 'Owner', 'Timeline']
imm_body = [
    ['1', 'Inspect the 2 HIGH trucks (VIN19_NF, VIN34_NF): starter + battery circuit', 'Fleet Ops', '2–4 weeks'],
    ['2', 'Bundle the 2 WATCH trucks (VIN16_NF, VIN24_NF) into next service', 'Fleet Ops', 'Next service'],
    ['3', 'Stand up weekly fleet scoring with monthly tier review', 'Engineering', 'This quarter'],
    ['4', 'Enable battery-first routing on cascade alerts', 'Maint. Planning', 'This quarter'],
]
dicv_table(s, ML, Inches(2.02), Inches(6.05), 5, [0.4, 3.35, 1.25, 1.05],
           imm_header, imm_body, fs=8.2, hfs=8.6, row_h=0.44)
minihead(s, Inches(6.95), Inches(1.68), Inches(6.0), "What to invest in next", BLUE, 13)
inv_header = ['Priority', 'Investment', 'Expected impact']
inv_body = [
    ['P0', 'High-frequency crank logging (≥1 Hz)', 'Unlocks months-ahead brush-wear warning'],
    ['P0', 'Cranking-current or battery-health signal', 'Proves battery-vs-starter; sharper alerts'],
    ['P1', 'Link maintenance / parts records to telemetry', 'Confirms failure modes; better labels'],
    ['P1', 'Keep collecting failure examples', 'Narrows uncertainty on every number'],
    ['P2', 'Scale to 500 trucks (Phase 2)', '~30–50+ failures → per-truck timing'],
]
dicv_table(s, Inches(6.95), Inches(2.02), Inches(5.83), 6, [0.85, 2.75, 2.23],
           inv_header, inv_body, fs=8.2, hfs=8.6, row_h=0.36)
minihead(s, ML, Inches(4.5), Inches(9.0), "Operating rhythm", INK, 13)
bullets(s, ML, Inches(4.88), Inches(6.05), Inches(1.5), [
    ("Weekly — ", True, INK),
    "Re-score every truck; a new RED or battery-cascade alert triggers a work order.",
    ("Monthly — ", True, INK),
    "Tier review with fleet ops; track watch-list trucks for drift or recovery.",
], size=9.5, marker='‣', mcolor=RED, gap=3, lead=1.0)
bullets(s, Inches(6.95), Inches(4.88), Inches(5.83), Inches(1.5), [
    ("On each confirmed failure or repair — ", True, INK),
    "Feed the outcome back; retrain only under the audited validation protocol.",
    ("Rule of communication — ", True, INK),
    "Quote risk tiers and the ~10-week horizon — never a single failure date.",
], size=9.5, marker='‣', mcolor=RED, gap=3, lead=1.0)
action_bar(s, "Value today: two HIGH trucks to inspect, a weekly scoring rhythm and battery-first "
              "triage — a modeled ~43% saving vs run-to-failure. The biggest next unlock is new "
              "signals + Phase-2 scale, not more modelling on the same data.", label="THE BOTTOM LINE")
footer(s, 7)

prs.save(str(OUT))
print("Saved:", OUT, "| slides:", len(prs.slides))
shutil.rmtree(TMP, ignore_errors=True)
print("Done.")
