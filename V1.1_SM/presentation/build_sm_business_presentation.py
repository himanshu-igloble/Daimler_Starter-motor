#!/usr/bin/env python3
"""
Build 5-slide Business-Oriented Starter Motor Predictive Maintenance presentation.
Results-focused for executive/fleet-ops audience -- no model internals.
Same visual language as the technical deck (and the Alternator V10.6.2 decks).

All numbers from STARTER MOTOR/V1.1/reports/ (comparison report, model card,
executive recommendation, alerts+horizon) -- no invention, no cost figures.
"""

import os
import sys
import tempfile
import shutil
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── PATHS ──────────────────────────────────────────────────────
ROOT   = Path(__file__).resolve().parent.parent.parent.parent
V11    = ROOT / "STARTER MOTOR" / "V1.1"
GRAPHS = V11 / "graphs"
OUT_DIR = V11 / "presentation"
OUT_DIR.mkdir(exist_ok=True)

# ── VIN DISPLAY RENUMBERING (2026-06-11) ───────────────────────
# Failed VIN1-14 unchanged; NF shifted +14 (old VIN1_NF -> VIN15_NF). All
# slide text/tables pass through map_nf_text at render time; the underlying
# results artifacts retain the ORIGINAL labels (results/V1_1_SM_vin_naming_map.csv).
sys.path.insert(0, str(V11 / "src"))
from V1_1_SM_vin_display_map import map_nf_text, display_label

VIN_FOOTNOTE = ("Sequential fleet numbering: VIN1-14 failed, VIN15-34 in-service "
                "(raw-file mapping: results/V1_1_SM_vin_naming_map.csv).")

# ── DIMENSIONS ─────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)

# ── COLOURS (identical to ALT decks) ───────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
GOLD       = RGBColor(0xC5, 0x8B, 0x1F)
GREY_MED   = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GREY  = RGBColor(0xB0, 0xB8, 0xC4)
KT_HEADER  = RGBColor(0x2E, 0x50, 0x90)
KT_BODY    = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_PASS = RGBColor(0x1E, 0x8C, 0x45)
RED_FAIL   = RGBColor(0xC0, 0x39, 0x2B)
ORANGE_W   = RGBColor(0xE3, 0x6C, 0x09)
BG_LIGHT   = RGBColor(0xF5, 0xF6, 0xFA)
BG_TABLE_H = RGBColor(0x0D, 0x1B, 0x2A)
BG_TABLE_A = RGBColor(0xF0, 0xF2, 0xF5)
BG_TABLE_B = RGBColor(0xFF, 0xFF, 0xFF)
KPI_BLUE   = RGBColor(0x1A, 0x5C, 0xB0)
KPI_GREEN  = RGBColor(0x15, 0x7F, 0x3D)
KPI_AMBER  = RGBColor(0xB8, 0x7A, 0x0F)
KPI_RED    = RGBColor(0xB0, 0x2A, 0x2A)

FONT = 'Calibri'
TMPDIR = tempfile.mkdtemp(prefix='sm_biz_charts_')


# ── HELPER FUNCTIONS (identical to ALT business deck) ──────────
def make_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    return shape

def add_rounded_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1.5)
    else:
        ln.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, font_size=12, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = map_nf_text(str(text))      # render-time VIN display mapping
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(2)
    return txBox

def add_multiline(slide, left, top, w, h, lines, font_size=11, color=DARK_TEXT,
                  bold=False, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, bold, color
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if bullet else "") + map_nf_text(str(txt))
        p.font.size = Pt(font_size)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = FONT
        p.space_after = Pt(4)
    return txBox

def add_header_bar(slide, title_text, subtitle_text=""):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             title_text, font_size=22, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.35),
                 subtitle_text, font_size=13, color=LIGHT_GREY)

def add_footer(slide, dark_bg=False):
    y = Inches(7.05)
    if not dark_bg:
        add_rect(slide, Inches(0), y, SW, Inches(0.45), BG_LIGHT)
    clr = LIGHT_GREY if dark_bg else GREY_MED
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "V1.1_SM  |  2026-06-10  |  BharatBenz 5528T Starter Motor Predictive Maintenance  |  BytEdge CONFIDENTIAL",
             font_size=8, color=clr, bold=False)

def add_kpi_tile(slide, left, top, w, h, label, value, color=KPI_BLUE, status=""):
    add_rounded_rect(slide, left, top, w, h, WHITE, border=color)
    add_text(slide, left + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.22),
             label, font_size=9, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.3), w - Inches(0.2), Inches(0.4),
             str(value), font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    if status:
        sc = GREEN_PASS if "PASS" in status.upper() or "YES" in status.upper() else (RED_FAIL if "FAIL" in status.upper() or "NO" in status.upper() else KPI_AMBER)
        add_text(slide, left + Inches(0.1), top + Inches(0.68), w - Inches(0.2), Inches(0.2),
                 status, font_size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)

def add_table_shape(slide, left, top, w, rows, cols, data, col_widths=None):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ri, row_data in enumerate(data):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = map_nf_text(str(val))   # render-time VIN display mapping
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            if ri == 0:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_H
            else:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_A if ri % 2 == 0 else BG_TABLE_B
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return tbl_shape

def add_key_takeaways(slide, bullets, left=Inches(0.4), top=Inches(5.6),
                      width=Inches(12.5), height=None):
    if height is None:
        height = Inches(0.25 + 0.24 * len(bullets))
    add_rounded_rect(slide, left, top, width, height,
                     RGBColor(0xE8, 0xEE, 0xF7), border=KT_HEADER)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), Inches(3), Inches(0.25),
             "KEY TAKEAWAYS", font_size=10, bold=True, color=KT_HEADER)
    y = top + Inches(0.28)
    for b in bullets:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.22),
                 "•  " + b, font_size=10, color=KT_BODY)
        y += Inches(0.22)


# ── CHART: BUSINESS VALUE PIPELINE ─────────────────────────────
def chart_value_pipeline():
    fig, ax = plt.subplots(figsize=(12, 3.0))
    stages = [
        ('CAN Bus\nTelemetry', '#B0B8C4', 'Existing data from\n34 BharatBenz trucks'),
        ('AI Risk\nScoring', '#1A5CB0', 'Calibrated probability,\n13 of 14 failures caught'),
        ('Early-Warning\nAlerts', '#2E5090', 'Flagged trucks typically\nwithin ~10 weeks of failure'),
        ('Battery-vs-Starter\nTriage', '#C58B1F', 'Cascade detector routes\nbattery-first (0 false alarms)'),
        ('Maintenance\nDecision', '#1E8C45', 'RED: inspect in 2-4 wks\nAMBER: next service'),
    ]
    for i, (label, color, detail) in enumerate(stages):
        x = i * 2.35
        rect = mpatches.FancyBboxPatch((x, 0.6), 1.9, 1.8, boxstyle="round,pad=0.12",
                                        facecolor=color, edgecolor='#0D1B2A', linewidth=1.2)
        ax.add_patch(rect)
        ax.text(x + 0.95, 1.75, label, ha='center', va='center', fontsize=10.5,
                fontweight='bold', color='white')
        ax.text(x + 0.95, 1.05, detail, ha='center', va='center', fontsize=7.5,
                color='#F5F6FA', style='italic')
        if i < len(stages) - 1:
            ax.annotate('', xy=(x + 2.15, 1.5), xytext=(x + 2.0, 1.5),
                       arrowprops=dict(arrowstyle='->', color='#C58B1F', lw=2.5))
    ax.set_xlim(-0.3, len(stages) * 2.35)
    ax.set_ylim(0, 3.0)
    ax.axis('off')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'value_pipeline.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_fleet_risk():
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    categories = ['LOW RISK\n(Green)', 'WATCH\n(Amber)', 'HIGH RISK\n(Red)']
    counts_nf = [16, 2, 2]
    colors = ['#1E8C45', '#E36C09', '#C0392B']
    ax.bar(categories, counts_nf, color=colors, width=0.5,
           edgecolor='#0D1B2A', linewidth=0.8, zorder=3)
    for i, v in enumerate(counts_nf):
        label = f'{v} truck{"s" if v > 1 else ""}'
        ax.text(i, v + 0.3, label, ha='center', va='bottom',
                fontsize=13, fontweight='bold', color='#0D1B2A')
    ax.set_ylim(0, 19)
    ax.set_ylabel('Number of In-Service Trucks', fontsize=11, color='#606060')
    ax.set_title('Current Fleet Risk Distribution (20 Trucks)', fontsize=13,
                 fontweight='bold', color='#0D1B2A', pad=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'fleet_risk.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_what_works_what_doesnt():
    fig, ax = plt.subplots(figsize=(11, 3.0))
    cats = ['Fleet risk\nranking', 'Early-warning\nalerts', 'Battery-vs-starter\ntriage', 'Exact failure\ndates', 'Silent/abrupt\nfailures']
    scores = [93, 93, 100, 0, 0]
    colors = ['#1E8C45', '#1E8C45', '#1E8C45', '#C0392B', '#C0392B']
    labels = ['13/14 failures flagged', '13/14 fire an alert, ~10-wk window', '0 false alarms on cascade detector', 'Not possible at 14 events', '~4 of 14 invisible to telemetry']
    ax.barh(cats[::-1], scores[::-1], color=colors[::-1], height=0.55,
            edgecolor='#0D1B2A', linewidth=0.5)
    for i, (s, l) in enumerate(zip(scores[::-1], labels[::-1])):
        c = '#1E8C45' if s > 50 else '#C0392B'
        ax.text(max(s, 5) + 2, i, l, ha='left', va='center',
                fontsize=10, fontweight='bold', color=c)
    ax.set_xlim(0, 165)
    ax.set_xlabel('Capability (%)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'effectiveness.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ── GENERATE CHARTS ────────────────────────────────────────────
print("Generating business charts...")
img_pipeline = chart_value_pipeline()
img_risk = chart_fleet_risk()
img_eff = chart_what_works_what_doesnt()
print("Charts generated.")

# Existing stable visualizations (daily-risk dashboard verified by the
# 2026-06-10 data audit; forecast-endpoint graphs re-rendered)
VIZ_FLEET_RISK = str(GRAPHS / "V1_1_SM_fleet_risk.png")
VIZ_DAILY_VIN6_F = str(GRAPHS / "V1_1_SM_daily_risk_VIN6_F_SM_dashboard.png")
for p in (VIZ_FLEET_RISK, VIZ_DAILY_VIN6_F):
    if not Path(p).exists():
        raise FileNotFoundError(f"Required image missing: {p}")


# ══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION -- 5 SLIDES
# ══════════════════════════════════════════════════════════════
prs = make_prs()


# ─── SLIDE 1: TITLE + HEADLINE RESULTS ──────────────────────
print("Building slide 1: Title & Headline")
s = blank_slide(prs)
add_rect(s, Inches(0), Inches(0), SW, SH, NAVY)
add_rect(s, Inches(0), Inches(4.8), SW, Inches(0.05), GOLD)

add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.7),
         "STARTER MOTOR PREDICTIVE MAINTENANCE",
         font_size=36, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
         "From Roadside No-Starts to a Monitored Fleet Program",
         font_size=20, color=GOLD)
add_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(0.4),
         "BharatBenz 5528T Heavy-Duty Trucks  |  34-Vehicle Study  |  June 2026  |  V1.1_SM",
         font_size=14, color=LIGHT_GREY)

# Headline KPI row below the gold line
y = Inches(5.1)
tw, th, gap, x0 = Inches(2.95), Inches(1.1), Inches(0.22), Inches(0.45)

add_rounded_rect(s, x0, y, tw, th, RGBColor(0x15, 0x25, 0x3A), border=GOLD)
add_text(s, x0 + Inches(0.1), y + Inches(0.08), tw - Inches(0.2), Inches(0.22),
         "FAILURES CAUGHT", font_size=9, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, x0 + Inches(0.1), y + Inches(0.35), tw - Inches(0.2), Inches(0.45),
         "13 of 14", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, x0 + Inches(0.1), y + Inches(0.82), tw - Inches(0.2), Inches(0.2),
         "Catch-most setting -- also flags 5 of 20 healthy trucks", font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

x1 = x0 + tw + gap
add_rounded_rect(s, x1, y, tw, th, RGBColor(0x15, 0x25, 0x3A), border=GOLD)
add_text(s, x1 + Inches(0.1), y + Inches(0.08), tw - Inches(0.2), Inches(0.22),
         "EARLY-WARNING WINDOW", font_size=9, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, x1 + Inches(0.1), y + Inches(0.35), tw - Inches(0.2), Inches(0.45),
         "~10 weeks", font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, x1 + Inches(0.1), y + Inches(0.82), tw - Inches(0.2), Inches(0.2),
         "Flagged trucks are typically within ~2.3 months of failure", font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

x2 = x1 + tw + gap
add_rounded_rect(s, x2, y, tw, th, RGBColor(0x15, 0x25, 0x3A), border=GOLD)
add_text(s, x2 + Inches(0.1), y + Inches(0.08), tw - Inches(0.2), Inches(0.22),
         "FALSE ALARMS (RED TIER)", font_size=9, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, x2 + Inches(0.1), y + Inches(0.35), tw - Inches(0.2), Inches(0.45),
         "2 of 20", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, x2 + Inches(0.1), y + Inches(0.82), tw - Inches(0.2), Inches(0.2),
         "Stricter RED-tier setting -- catches 10 of 14 failures", font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

x3 = x2 + tw + gap
add_rounded_rect(s, x3, y, tw, th, RGBColor(0x15, 0x25, 0x3A), border=GOLD)
add_text(s, x3 + Inches(0.1), y + Inches(0.08), tw - Inches(0.2), Inches(0.22),
         "FLEET VISIBILITY", font_size=9, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, x3 + Inches(0.1), y + Inches(0.35), tw - Inches(0.2), Inches(0.45),
         "34 Trucks", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, x3 + Inches(0.1), y + Inches(0.82), tw - Inches(0.2), Inches(0.2),
         "Every truck scored, tiered, explained", font_size=8, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_footer(s, dark_bg=True)


# ─── SLIDE 2: WHAT WAS BUILT ────────────────────────────────
print("Building slide 2: What Was Built")
s = blank_slide(prs)
add_header_bar(s, "WHAT WAS BUILT",
               "Turning existing truck telemetry into actionable starter motor maintenance intelligence")

s.shapes.add_picture(img_pipeline, Inches(0.4), Inches(1.15), Inches(12.5), Inches(2.8))

add_text(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.3),
         "THREE PLAIN-LANGUAGE CAPABILITIES", font_size=14, bold=True, color=NAVY)

col_w = Inches(3.8)
col_h = Inches(1.6)
col_y = Inches(4.55)

add_rounded_rect(s, Inches(0.5), col_y, col_w, col_h, BG_LIGHT, border=KPI_BLUE)
add_text(s, Inches(0.65), col_y + Inches(0.1), Inches(3.5), Inches(0.25),
         "1. RISK RANKING", font_size=12, bold=True, color=KPI_BLUE)
add_multiline(s, Inches(0.65), col_y + Inches(0.4), Inches(3.5), Inches(1.0), [
    "Every truck gets a calibrated risk score from its own voltage patterns",
    "Scores map to GREEN / AMBER / RED maintenance tiers",
    "13 of 14 past failures flagged before they failed (catch-most setting)",
], font_size=9, bullet=True, color=DARK_TEXT)

add_rounded_rect(s, Inches(4.75), col_y, col_w, col_h, BG_LIGHT, border=KPI_AMBER)
add_text(s, Inches(4.9), col_y + Inches(0.1), Inches(3.5), Inches(0.25),
         "2. EARLY-WARNING ALERTS", font_size=12, bold=True, color=KPI_AMBER)
add_multiline(s, Inches(4.9), col_y + Inches(0.4), Inches(3.5), Inches(1.0), [
    "A flagged truck is typically within ~10 weeks of failure",
    "Battery-cascade alert: ~9.5-week median lead, zero false alarms",
    "13 of 14 failed trucks fired at least one alert channel",
], font_size=9, bullet=True, color=DARK_TEXT)

add_rounded_rect(s, Inches(9.0), col_y, col_w, col_h, BG_LIGHT, border=KPI_GREEN)
add_text(s, Inches(9.15), col_y + Inches(0.1), Inches(3.5), Inches(0.25),
         "3. BATTERY-VS-STARTER TRIAGE", font_size=12, bold=True, color=KPI_GREEN)
add_multiline(s, Inches(9.15), col_y + Inches(0.4), Inches(3.5), Inches(1.0), [
    "Detects the battery-degradation cascade that mimics starter failure",
    "Routes the work order battery-first -- the cheapest correct fix",
    "Provably ignores battery replacements (no false trigger)",
], font_size=9, bullet=True, color=DARK_TEXT)

add_key_takeaways(s, [
    "No new sensors needed -- the system runs on existing CAN bus data already collected from the fleet",
    "Every alert comes with its reasons: which signal moved, by how much, and what would clear the flag",
], top=Inches(6.35), height=Inches(0.65))
add_footer(s)


# ─── SLIDE 3: HOW WELL IT WORKS ─────────────────────────────
print("Building slide 3: How Well It Works")
s = blank_slide(prs)
add_header_bar(s, "HOW WELL IT WORKS",
               "Validated results across 34 trucks (14 known failures + 20 in-service)")

# Daily-risk degradation view (3198x2065 px; width-only embed preserves aspect, h = 2.85")
s.shapes.add_picture(VIZ_DAILY_VIN6_F, Inches(0.4), Inches(1.15), width=Inches(4.42))
add_text(s, Inches(5.0), Inches(1.2), Inches(2.9), Inches(0.3),
         "DAILY DEGRADATION VIEW", font_size=11, bold=True, color=NAVY)
add_multiline(s, Inches(5.0), Inches(1.55), Inches(2.9), Inches(2.4), [
    ("Each truck gets a daily-updated degradation view; colors = action tiers (green = healthy, red = act now).", False, DARK_TEXT),
    ("This failed truck (VIN6) walked the full green-to-red path before its starter failed -- the battery alert fired ~70 days early.", False, DARK_TEXT),
    ("The dotted tail is a fleet-average projection, not a failure date.", False, GREY_MED),
], font_size=9, bullet=True)
# Fleet risk mini-chart (figsize 5.5x3.5; width-only embed preserves aspect, h = 2.86")
s.shapes.add_picture(img_risk, Inches(8.2), Inches(1.15), width=Inches(4.5))

add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.3),
         "HEADLINE RESULTS", font_size=14, bold=True, color=NAVY)
add_text(s, Inches(3.4), Inches(4.16), Inches(9.4), Inches(0.25),
         VIN_FOOTNOTE, font_size=7.5, color=GREY_MED, align=PP_ALIGN.RIGHT)

res_data = [
    ['What', 'Result', 'Business Impact'],
    ['Failure detection', '13 of 14 caught at the catch-most setting (5 of 20 healthy flagged)', 'High-risk trucks inspected before they strand'],
    ['False-alarm burden', '2 of 20 healthy at RED tier -- the stricter setting, which catches 10 of 14', 'Inspection effort goes where the risk is'],
    ['Early warning', 'Flagged trucks typically within ~10 weeks; battery alerts ~9.5 weeks ahead', 'Time to schedule, source parts, avoid roadside events'],
    ['Triage routing', 'Battery-cascade detector: 0 false alarms', 'Cheapest correct intervention happens first'],
    ['Fleet status today', '16 LOW / 2 WATCH / 2 HIGH risk', 'Two trucks (VIN5_NF, VIN20_NF) need inspection now'],
]
add_table_shape(s, Inches(0.4), Inches(4.45), Inches(12.5), 6, 3, res_data,
                col_widths=[2.5, 4.5, 5.5])

add_key_takeaways(s, [
    "One dial, one honest trade: the catch-most setting flags 13 of 14 failures but 5 of 20 healthy trucks; the stricter RED tier flags 10 of 14 at 2 of 20 -- operations picks the setting",
    "Versus the first-generation system: one more failure caught (including its worst miss), trustworthy probabilities, and an early-warning capability it did not have",
], top=Inches(6.55), height=Inches(0.45))
add_footer(s)


# ─── SLIDE 4: WHAT IT CANNOT DO ─────────────────────────────
print("Building slide 4: What It Cannot Do")
s = blank_slide(prs)
add_header_bar(s, "WHAT IT CANNOT DO -- AND WHY",
               "Transparent assessment of current system boundaries")

s.shapes.add_picture(img_eff, Inches(0.4), Inches(1.15), Inches(7.0), Inches(2.7))

add_text(s, Inches(7.8), Inches(1.15), Inches(5.2), Inches(0.3),
         "HONEST BOUNDARIES", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(7.8), Inches(1.55), Inches(5.2), Inches(2.5), [
    ("It does not predict failure dates", True, RED_FAIL),
    "With only 14 failure examples, a countdown clock is mathematically unsupportable -- every attempted timing model lost to a simple constant. The system says 'within ~10 weeks', never 'on June 14th'.",
    ("~4 of 14 failures are invisible", True, RED_FAIL),
    "Sudden internal failures (windings, seizure) and trucks whose telemetry goes silent give no electrical warning at this data resolution.",
    ("Some alerts will be early or wrong", True, RED_FAIL),
    "A few healthy trucks get flagged: one shows genuinely stressed electrical signals; another is flagged on its risk score alone. Those inspections are cheap insurance, but they are not failures.",
], font_size=9, bullet=True, color=DARK_TEXT)

add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.3),
         "DATA GAPS LIMITING ACCURACY", font_size=14, bold=True, color=NAVY)

gap_data = [
    ['Gap', 'Impact on System', 'How to Close'],
    ['Voltage sampled only every 5 seconds', 'The 2-4 month brush-wear warning signal is physically unreadable', 'High-frequency crank logging (post-2026 vehicle architecture)'],
    ['No cranking current or battery health signal', 'Battery vs starter cannot be proven, only inferred', 'Add current or battery SoC/SoH instrumentation'],
    ['No maintenance/parts records in the data', 'Failure modes are inferred from patterns, not confirmed', 'Link workshop records to telemetry'],
    ['Only 14 failure examples', 'Statistical uncertainty on every number', 'Keep collecting failures as the fleet ages'],
    ['Some trucks transmit no crank data', '7 of 34 trucks have a config that hides starter activity', 'Telemetry-config audit; treat silence as a maintenance trigger'],
]
add_table_shape(s, Inches(0.4), Inches(4.45), Inches(12.5), 6, 3, gap_data,
                col_widths=[3.4, 4.6, 4.5])

add_key_takeaways(s, [
    "These are physics and data-size constraints, not model deficiencies -- the system performs at the honest ceiling of what this data allows",
    "The boundaries are measured and stated, which is what makes the rest of the numbers trustworthy",
], top=Inches(6.55), height=Inches(0.45))
add_footer(s)


# ─── SLIDE 5: RECOMMENDED ACTIONS & NEXT STEPS ──────────────
print("Building slide 5: Recommended Actions & Next Steps")
s = blank_slide(prs)
add_header_bar(s, "RECOMMENDED ACTIONS & NEXT STEPS",
               "Priority actions for fleet operations and what to invest in next")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.3),
         "IMMEDIATE ACTIONS", font_size=16, bold=True, color=GREEN_PASS)

imm_data = [
    ['#', 'Action', 'Owner', 'Timeline'],
    ['1', 'Inspect the 2 RED trucks (VIN5_NF, VIN20_NF): starter + battery circuit', 'Fleet Ops', '2-4 weeks'],
    ['2', 'Bundle the 2 AMBER trucks (VIN2_NF, VIN10_NF) into next service', 'Fleet Ops', 'Next service'],
    ['3', 'Stand up weekly fleet scoring with monthly tier review', 'Engineering', 'This quarter'],
    ['4', 'Enable battery-first routing on cascade alerts', 'Maintenance Planning', 'This quarter'],
]
add_table_shape(s, Inches(0.5), Inches(1.55), Inches(5.8), 5, 4, imm_data,
                col_widths=[0.4, 2.9, 1.3, 1.2])

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.3),
         "WHAT TO INVEST IN NEXT", font_size=16, bold=True, color=KPI_BLUE)

road_data = [
    ['Priority', 'Investment', 'Expected Impact'],
    ['P0', 'High-frequency crank logging (>=1 Hz)', 'Unlocks months-ahead brush-wear warning'],
    ['P0', 'Cranking current or battery health signal', 'Proves battery-vs-starter; sharper alerts'],
    ['P1', 'Link maintenance/parts records to telemetry', 'Confirms failure modes; better labels'],
    ['P1', 'Keep collecting failure examples', 'Narrows uncertainty on every number'],
    ['P2', 'Telemetry-config and transmission-health audit', 'Closes the silent-truck blind spot'],
]
add_table_shape(s, Inches(7.0), Inches(1.55), Inches(5.8), 6, 3, road_data,
                col_widths=[0.9, 2.6, 2.3])

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.3),
         "OPERATING RHYTHM", font_size=16, bold=True, color=NAVY)

add_multiline(s, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.2), [
    ("Weekly", True, NAVY),
    "Re-score every truck; new RED or battery-cascade alert triggers a work order",
    ("Monthly", True, NAVY),
    "Tier review with fleet ops; track the watch-list trucks for drift or recovery",
], font_size=10, bullet=True)

add_multiline(s, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.2), [
    ("On each confirmed failure or repair", True, NAVY),
    "Feed the outcome back; retrain only under the audited validation protocol",
    ("Rule of communication", True, NAVY),
    "Quote risk tiers and the ~10-week window -- never a failure date",
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "The system delivers value today: two RED trucks to inspect, a weekly scoring rhythm, and battery-first triage on alerts",
    "Highest-leverage investments are new signals (crank logging, battery health), not more modeling on the same data",
], top=Inches(6.55), height=Inches(0.45))
add_footer(s)


# ── SAVE ───────────────────────────────────────────────────────
out_path = str(OUT_DIR / "SM_Business_Summary_V1.1.pptx")
prs.save(out_path)
print(f"\nPresentation saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")

shutil.rmtree(TMPDIR, ignore_errors=True)
print("Done.")
