"""
build_dicv_executive_deck.py - DICV executive business presentation for the
Starter Motor Failure Prediction Program V1.1.

External-facing deck: no model names, algorithms, features, or internal
methodology anywhere. The engine is referred to only as the "Predictive
Intelligence Engine". Every number is traced to validated V1.1 artifacts:

  results/V1_1_SM_model_spec.json         AUROC 0.9321 [0.811, 0.986], 13/14 @ 5/20,
                                          RED-tier 10/14 @ 2/20, tier counts 20/2/12
  results/V1_1_SM_alert_policy.csv        per-VIN first alert, channel, lead days
  reports/V1_1_SM_alerts_horizon.md       median lead 168 d, A2 0/20 FP & 66.5 d,
                                          10-week validated horizon (k*=10)
  reports/V1_1_SM_RESULTS_MASTER.md       fleet/data volumes, VIN stories
  reports/V1_1_SM_daily_graph_verification.md  VIN6_F zone dates

Output: presentation/DICV_SM_Executive_Business_Review_V1_1.pptx
Run:    py -3 "STARTER MOTOR/V1.1/presentation/build_dicv_executive_deck.py"
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(r"D:\Daimler-starter_motor_alternator_battery\STARTER MOTOR\V1.1")
ASSETS = ROOT / "presentation" / "assets"
OUT = ROOT / "presentation" / "DICV_SM_Executive_Business_Review_V1_1.pptx"

# ---------------------------------------------------------------------------
# Design system
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x14, 0x26, 0x3A)
NAVY2 = RGBColor(0x1F, 0x3A, 0x5F)
INK = RGBColor(0x21, 0x2B, 0x36)
SLATE = RGBColor(0x5B, 0x6B, 0x7C)
LINE = RGBColor(0xD7, 0xDE, 0xE5)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
REDACC = RGBColor(0xC8, 0x10, 0x2E)
GREEN_D = RGBColor(0x1B, 0x7A, 0x3D)
GREEN_L = RGBColor(0xDD, 0xEF, 0xE2)
AMBER_D = RGBColor(0xB4, 0x53, 0x09)
AMBER_L = RGBColor(0xFB, 0xEA, 0xD2)
RED_D = RGBColor(0x8B, 0x00, 0x00)
RED_L = RGBColor(0xFA, 0xE3, 0xE3)
BLUE_D = RGBColor(0x15, 0x65, 0xC0)
PURP_D = RGBColor(0x6C, 0x33, 0x83)

FONT = "Segoe UI"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _noline(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        _solid(sp, fill)
    if line_color is None:
        _noline(sp)
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=2):
    """paras: list of (text, size, color, bold) or (text, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        text, size, color, bold = para[:4]
        italic = para[4] if len(para) > 4 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = FONT
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return tb


def add_multirun(slide, x, y, w, h, runs, size=10.5, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, space_after=4):
    """One textbox, many paragraphs; each paragraph = list of (text, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runlist in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for text, color, bold in runlist:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.color.rgb = color
            f.bold = bold
    return tb


def add_pic(slide, path, x, y, w=None, h=None, border=True):
    img = Image.open(path)
    ar = img.width / img.height
    if w is not None and h is None:
        h = Emu(int(w / ar))
    elif h is not None and w is None:
        w = Emu(int(h * ar))
    pic = slide.shapes.add_picture(str(path), x, y, w, h)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def header(slide, kicker, title, num, total=9):
    add_rect(slide, 0, 0, Inches(0.12), SH, NAVY)
    add_rect(slide, Inches(0.12), 0, Inches(0.10), Inches(0.62), REDACC)
    add_text(slide, Inches(0.45), Inches(0.22), Inches(10.6), Inches(0.28),
             [(kicker.upper(), 10.5, SLATE, True)])
    add_text(slide, Inches(0.45), Inches(0.47), Inches(11.4), Inches(0.62),
             [(title, 25, NAVY, True)])
    add_text(slide, Inches(12.15), Inches(0.24), Inches(0.95), Inches(0.3),
             [(f"{num:02d} / {total:02d}", 10, SLATE, False)], align=PP_ALIGN.RIGHT)


def footer(slide, takeaway, action):
    y = Inches(6.52)
    h = Inches(0.68)
    add_rect(slide, Inches(0.45), y, Inches(7.35), h, LIGHT)
    add_rect(slide, Inches(0.45), y, Inches(0.07), h, NAVY)
    add_text(slide, Inches(0.68), y + Inches(0.08), Inches(7.0), Inches(0.2),
             [("KEY TAKEAWAY", 8.5, SLATE, True)])
    add_text(slide, Inches(0.68), y + Inches(0.27), Inches(7.0), Inches(0.4),
             [(takeaway, 10.5, NAVY, True)])
    add_rect(slide, Inches(7.95), y, Inches(4.93), h, LIGHT)
    add_rect(slide, Inches(7.95), y, Inches(0.07), h, REDACC)
    add_text(slide, Inches(8.18), y + Inches(0.08), Inches(4.6), Inches(0.2),
             [("RECOMMENDED ACTION", 8.5, SLATE, True)])
    add_text(slide, Inches(8.18), y + Inches(0.27), Inches(4.6), Inches(0.4),
             [(action, 10.5, REDACC, True)])
    add_text(slide, Inches(0.45), Inches(7.24), Inches(12.45), Inches(0.2),
             [("Daimler India Commercial Vehicles  ·  Starter Motor Failure "
               "Prediction Program V1.1  ·  Confidential — for internal DICV use",
               8, SLATE, False)])


def kpi_tile(slide, x, y, w, h, big, label, sub, accent=NAVY):
    add_rect(slide, x, y, w, h, WHITE, line_color=LINE, line_w=1.0)
    add_rect(slide, x, y, w, Inches(0.07), accent)
    add_text(slide, x + Inches(0.16), y + Inches(0.16), w - Inches(0.32), Inches(0.62),
             [(big, 30, accent, True)])
    add_text(slide, x + Inches(0.16), y + Inches(0.80), w - Inches(0.32), Inches(0.34),
             [(label, 11, INK, True)])
    add_text(slide, x + Inches(0.16), y + Inches(1.12), w - Inches(0.32),
             h - Inches(1.2), [(sub, 8.7, SLATE, False)], space_after=0)


def chip(slide, x, y, w, text, fill, txtcolor, size=9.5, h=Inches(0.30)):
    add_rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y + Inches(0.015), w, h, [(text, size, txtcolor, True)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# SLIDE 1 - Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.14), REDACC)
add_rect(s, Inches(0.9), Inches(2.02), Inches(0.14), Inches(1.62), REDACC)

add_text(s, Inches(0.9), Inches(1.28), Inches(11.5), Inches(0.35),
         [("DAIMLER INDIA COMMERCIAL VEHICLES  ·  FLEET RELIABILITY", 13,
           RGBColor(0x9F, 0xB3, 0xC8), True)])
add_text(s, Inches(1.28), Inches(2.02), Inches(11.2), Inches(1.7),
         [("Starter Motor Failure", 44, WHITE, True),
          ("Prediction Program", 44, WHITE, True)], space_after=0)
add_text(s, Inches(1.28), Inches(3.82), Inches(11.0), Inches(0.5),
         [("Executive Business Review  ·  Predictive Intelligence for the "
           "BharatBenz Fleet  ·  Version 1.1", 15, RGBColor(0xC9, 0xD6, 0xE3),
           False)])

stats = [("34", "vehicles analysed"), ("106M+", "telemetry records"),
         ("13 of 14", "failures identified in advance"),
         ("168 days", "median early warning")]
x = Inches(1.28)
for big, lab in stats:
    add_text(s, x, Inches(5.05), Inches(2.65), Inches(0.55),
             [(big, 26, WHITE, True)])
    add_text(s, x, Inches(5.62), Inches(2.65), Inches(0.35),
             [(lab, 10.5, RGBColor(0x9F, 0xB3, 0xC8), False)])
    x += Inches(2.85)

add_text(s, Inches(1.28), Inches(6.72), Inches(10.8), Inches(0.3),
         [("July 2026  ·  Confidential — prepared for DICV leadership, service "
           "engineering, warranty and fleet operations", 10,
           RGBColor(0x9F, 0xB3, 0xC8), False)])

# ===========================================================================
# SLIDE 2 - Executive Summary
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Executive summary", "Validated early warning for starter motor "
       "failures — months before they happen", 2)

tiles = [
    ("13 of 14", "Failures identified in advance",
     "Across 14 real starter motor failures, 13 vehicles raised a validated "
     "early warning before the failure date.", GREEN_D),
    ("168 days", "Median early-warning lead time",
     "First warning arrived a median 168 days before failure "
     "(range 77–424 days) — ample time to plan.", NAVY),
    ("93%", "Risk-ranking strength",
     "AUROC 0.93 (95% CI 0.81–0.99): independently validated ability to rank "
     "failing vehicles above healthy ones.", NAVY),
    ("10 weeks", "Validated action window",
     "A flagged vehicle is typically within ~10 weeks of failure — the "
     "confirmed window for maintenance action.", AMBER_D),
    ("0", "Battery-cascade false alarms",
     "The highest-precision alert fired on 4 failing vehicles and on none of "
     "the 20 healthy vehicles.", GREEN_D),
    ("106M", "Telemetry records analysed",
     "34 vehicles · 2,636 vehicle-weeks · 20,471 engine-start events — "
     "evidence base for every claim in this review.", SLATE),
]
tx, ty = Inches(0.45), Inches(1.38)
tw, th = Inches(2.72), Inches(2.28)
for i, (big, lab, sub, acc) in enumerate(tiles):
    kpi_tile(s, tx + (i % 3) * (tw + Inches(0.18)),
             ty + (i // 3) * (th + Inches(0.18)), tw, th, big, lab, sub, acc)

rx = Inches(9.28)
add_rect(s, rx, ty, Inches(3.6), Inches(4.74), LIGHT)
add_text(s, rx + Inches(0.22), ty + Inches(0.16), Inches(3.16), Inches(0.3),
         [("FLEET HEALTH AT A GLANCE", 10.5, NAVY, True)])
add_pic(s, ASSETS / "BIZ_tier_donut.png", rx + Inches(0.62), ty + Inches(0.52),
        w=Inches(2.36), border=False)
add_multirun(s, rx + Inches(0.22), ty + Inches(3.06), Inches(3.16), Inches(1.6), [
    [("Every vehicle, every week: ", NAVY, True),
     ("one clear health status and one clear action.", INK, False)],
    [("12 RED", RED_D, True), (" — 10 later failed; 2 flagged for inspection.",
      INK, False)],
    [("2 AMBER", AMBER_D, True), (" — bundled into next service visit.", INK, False)],
    [("20 GREEN", GREEN_D, True), (" — 16 healthy on routine ops; 4 subtle "
      "failures, 3 of them still caught by the alert layer (slide 8).",
      INK, False)],
], size=9.5, space_after=3)

footer(s, "V1.1 reliably separates degrading starter motors from healthy ones, "
          "with months of advance notice on real DICV fleet data.",
       "Endorse pilot deployment to convert validated warnings into avoided "
       "breakdowns.")

# ===========================================================================
# SLIDE 3 - Business solution overview
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "How it works — business view", "From vehicle operating data to "
       "maintenance decisions, every week", 3)

steps = [
    ("1", "Vehicle Operational Data", "Standard on-board electrical and "
     "usage signals — no new sensors required"),
    ("2", "Continuous Health Assessment", "Predictive Intelligence Engine "
     "evaluates each vehicle's electrical behaviour weekly"),
    ("3", "Early Risk Identification", "Departures from the vehicle's own "
     "healthy pattern are detected months ahead"),
    ("4", "Vehicle Risk Classification", "Each vehicle is classified GREEN / "
     "AMBER / RED with a validated risk score"),
    ("5", "Future Failure Outlook", "Flagged vehicles are typically within "
     "~10 weeks of failure — a confirmed planning window"),
    ("6", "Maintenance Prioritisation", "Workshops see a ranked list: which "
     "vehicles, how urgent, what to check first"),
    ("7", "Operational Decision Support", "Clear actions: inspect in 2–4 "
     "weeks, bundle at next service, or continue routine ops"),
]
bx, by = Inches(0.45), Inches(1.42)
bw, bh = Inches(3.03), Inches(1.28)
gapx, gapy = Inches(0.17), Inches(0.22)
positions = [(0, 0), (1, 0), (2, 0), (3, 0), (0, 1), (1, 1), (2, 1)]
for (col, row), (n, t, d) in zip(positions, steps):
    x = bx + col * (bw + gapx)
    y = by + row * (bh + gapy)
    acc = NAVY if row == 0 else NAVY2
    add_rect(s, x, y, bw, bh, WHITE, line_color=LINE, line_w=1.0)
    add_rect(s, x, y, Inches(0.07), bh, acc)
    add_rect(s, x + Inches(0.14), y + Inches(0.13), Inches(0.34), Inches(0.34),
             acc, shape=MSO_SHAPE.OVAL)
    add_text(s, x + Inches(0.14), y + Inches(0.155), Inches(0.34), Inches(0.3),
             [(n, 12, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.58), y + Inches(0.14), bw - Inches(0.72),
             Inches(0.35), [(t, 11, NAVY, True)])
    add_text(s, x + Inches(0.16), y + Inches(0.52), bw - Inches(0.30),
             Inches(0.72), [(d, 9, SLATE, False)], space_after=0)
    if (col, row) != (2, 1):
        ax = x + bw + Inches(0.008) if col < 3 else None
        if col < 3 and not (col == 3 and row == 0):
            add_rect(s, x + bw - Inches(0.005), y + bh / 2 - Inches(0.09),
                     Inches(0.18), Inches(0.18), NAVY,
                     shape=MSO_SHAPE.CHEVRON)

box_y = by + 2 * (bh + gapy) + Inches(0.12)
add_rect(s, Inches(0.45), box_y, Inches(9.4) + Inches(3.03), Inches(1.55), LIGHT)
add_text(s, Inches(0.72), box_y + Inches(0.15), Inches(11.9), Inches(0.3),
         [("WHAT THIS MEANS OPERATIONALLY", 10.5, NAVY, True)])
add_multirun(s, Inches(0.72), box_y + Inches(0.48), Inches(11.9), Inches(1.0), [
    [("The engine turns raw operating behaviour into a weekly fleet health "
      "report. Maintenance teams no longer wait for a breakdown call: "
      "degrading vehicles surface on a ranked watch-list months earlier, "
      "spare parts and workshop slots are arranged in advance, and healthy "
      "vehicles stay on the road undisturbed. The underlying analytics remain "
      "fully managed within the platform — DICV teams interact only with "
      "clear statuses and actions.", INK, False)],
], size=10.5, space_after=0)

footer(s, "A weekly, fully automated pipeline turns existing vehicle data "
          "into ranked, actionable maintenance intelligence.",
       "Integrate the weekly health report into existing workshop planning "
       "as the single fleet watch-list.")

# ===========================================================================
# SLIDE 4 - Case study: VIN6_F (flagship)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Vehicle case study — failure predicted", "VIN6: a starter motor "
       "failure signalled 168 days in advance", 4)

add_pic(s, ASSETS / "V1_1_SM_VIN6_F_SM_dashboard.png", Inches(0.45),
        Inches(1.30), w=Inches(7.85))

rx = Inches(8.52)
add_rect(s, rx, Inches(1.30), Inches(4.36), Inches(5.05), LIGHT)
add_text(s, rx + Inches(0.22), Inches(1.46), Inches(3.9), Inches(0.3),
         [("ANATOMY OF A PREDICTED FAILURE", 10.5, NAVY, True)])

events = [
    ("19 May 2025", "Early warning: sustained electrical volatility alert — "
     "vehicle enters the watch-list", PURP_D),
    ("25 May 2025", "Risk score leaves the healthy band — monitoring "
     "intensified", RGBColor(0xB7, 0x95, 0x0B)),
    ("25 Aug 2025", "Battery-cascade alert confirms the pathway — 70 days "
     "before failure", RED_D),
    ("31 Aug — 19 Oct", "Risk climbs steadily into RED: inspection and parts "
     "planning window", AMBER_D),
    ("04 Nov 2025", "Failure occurs — after a 168-day warning runway", RED_D),
]
ey = Inches(1.86)
for date, txt, col in events:
    add_rect(s, rx + Inches(0.22), ey + Inches(0.03), Inches(0.10),
             Inches(0.60), col)
    add_text(s, rx + Inches(0.44), ey, Inches(3.75), Inches(0.24),
             [(date, 9.5, col, True)])
    add_text(s, rx + Inches(0.44), ey + Inches(0.22), Inches(3.75),
             Inches(0.44), [(txt, 9, INK, False)], space_after=0)
    ey += Inches(0.72)

add_multirun(s, rx + Inches(0.22), ey + Inches(0.08), Inches(3.95),
             Inches(0.62), [
    [("With V1.1 live, this vehicle books a workshop slot in September — "
      "not a roadside recovery in November.", NAVY, True)],
], size=10, space_after=0)

footer(s, "Progressive degradation was visible, staged and confirmed by two "
          "independent alerts long before failure.",
       "Use this case as the pilot's operating template: alert → inspect → "
       "planned replacement.")

# ===========================================================================
# SLIDE 5 - Case studies: VIN10_F + VIN14_F + zone journey strip
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Vehicle case studies — failure predicted", "Different failure "
       "signatures, same result: months of advance notice", 5)

add_pic(s, ASSETS / "V1_1_SM_VIN10_F_SM_dashboard.png", Inches(0.45),
        Inches(1.32), w=Inches(6.12))
add_pic(s, ASSETS / "V1_1_SM_VIN14_F_SM_dashboard.png", Inches(6.76),
        Inches(1.32), w=Inches(6.12))

cap_y = Inches(5.42)
add_rect(s, Inches(0.45), cap_y, Inches(6.12), Inches(0.92), LIGHT)
add_multirun(s, Inches(0.62), cap_y + Inches(0.09), Inches(5.85), Inches(0.8), [
    [("VIN10 — starter-mechanism wear signature.  ", NAVY, True),
     ("Crank-anomaly alert fired 160 days before failure; risk climbed "
      "GREEN → RED five months ahead. Final risk 100%.", INK, False)],
], size=9.5, space_after=0)
add_rect(s, Inches(6.76), cap_y, Inches(6.12), Inches(0.92), LIGHT)
add_multirun(s, Inches(6.93), cap_y + Inches(0.09), Inches(5.85), Inches(0.8), [
    [("VIN14 — mixed electrical signature.  ", NAVY, True),
     ("First warning 245 days ahead; battery-cascade confirmation 28 days "
      "before failure. Joint-highest risk score in the fleet (1.00).",
      INK, False)],
], size=9.5, space_after=0)

footer(s, "The engine catches distinct degradation pathways — battery-driven, "
          "starter-mechanism and electrical-noise signatures alike.",
       "Route each alert type to its matching inspection checklist "
       "(battery circuit vs. starter system).")

# ===========================================================================
# SLIDE 6 - Healthy contrast: VIN17_NF
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Vehicle case study — healthy fleet", "What healthy looks like: "
       "stable scores, self-resolving deviations, zero alerts", 6)

add_pic(s, ASSETS / "V1_1_SM_VIN17_NF_SM_dashboard.png", Inches(0.45),
        Inches(1.30), w=Inches(7.85))

rx = Inches(8.52)
add_rect(s, rx, Inches(1.30), Inches(4.36), Inches(5.05), LIGHT)
add_text(s, rx + Inches(0.22), Inches(1.46), Inches(3.9), Inches(0.3),
         [("THE HEALTHY-VEHICLE PATTERN", 10.5, NAVY, True)])

hstats = [
    ("6%", "final risk score after 702 days of monitoring — firmly GREEN",
     GREEN_D),
    ("0", "alerts raised across the full observation period", GREEN_D),
    ("10 of 20", "healthy vehicles never triggered any alert channel",
     NAVY),
    ("16 of 20", "healthy vehicles finished in GREEN; 2 AMBER, 2 RED "
     "flagged for precautionary inspection", NAVY),
]
ey = Inches(1.88)
for big, txt, col in hstats:
    add_text(s, rx + Inches(0.22), ey, Inches(1.25), Inches(0.5),
             [(big, 21, col, True)])
    add_text(s, rx + Inches(1.5), ey + Inches(0.04), Inches(2.68),
             Inches(0.62), [(txt, 9, INK, False)], space_after=0)
    ey += Inches(0.72)

add_multirun(s, rx + Inches(0.22), ey + Inches(0.05), Inches(3.95),
             Inches(1.35), [
    [("Reading the contrast:  ", NAVY, True),
     ("short-lived score elevations on healthy vehicles settle back to GREEN "
      "on their own. Failing vehicles climb and stay high. That difference — "
      "sustained versus transient — is what the engine measures every week.",
      INK, False)],
], size=9.5, space_after=0)

footer(s, "Healthy vehicles stay green and quiet — the system earns trust by "
          "leaving well-running trucks alone.",
       "Track alert-free healthy-fleet share as a standing KPI during the "
       "pilot (baseline: 50% fully quiet, 80% green).")

# ===========================================================================
# SLIDE 7 - Vehicle-level intelligence (fleet ranking + health cards)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Vehicle-level intelligence", "One ranked view of all 34 vehicles "
       "— maintenance effort flows where risk is", 7)

add_pic(s, ASSETS / "BIZ_fleet_risk_bars.png", Inches(0.45), Inches(1.28),
        w=Inches(12.43), border=False)

cards = [
    ("VIN6_F", "RED · risk 100%", "Failed as predicted — 168-day warning "
     "validated", RED_D, RED_L),
    ("VIN19_NF", "RED · risk 96%", "In service — priority inspection within "
     "2–4 weeks", RED_D, RED_L),
    ("VIN16_NF", "AMBER · risk 45%", "In service — bundle checks into next "
     "scheduled service", AMBER_D, AMBER_L),
    ("VIN17_NF", "GREEN · risk 6%", "In service — routine operation, no "
     "action required", GREEN_D, GREEN_L),
]
cx = Inches(0.45)
cw = Inches(3.02)
cy = Inches(5.32)
for vin, status, txt, col, fill in cards:
    add_rect(s, cx, cy, cw, Inches(1.02), fill)
    add_rect(s, cx, cy, Inches(0.07), Inches(1.02), col)
    add_text(s, cx + Inches(0.18), cy + Inches(0.08), cw - Inches(0.3),
             Inches(0.26), [(f"{vin}   ·   {status}", 10, col, True)])
    add_text(s, cx + Inches(0.18), cy + Inches(0.36), cw - Inches(0.3),
             Inches(0.6), [(txt, 9, INK, False)], space_after=0)
    cx += cw + Inches(0.115)

footer(s, "Every vehicle carries a current status, a validated risk score "
          "and a clear next action — updated weekly.",
       "Adopt the tier rule in workshop planning: RED inspect in 2–4 weeks; "
       "AMBER at next service; GREEN routine.")

# ===========================================================================
# SLIDE 8 - Fleet risk outlook (runway + zone ladder)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Fleet risk outlook", "The early-warning runway: how much time "
       "maintenance teams actually get", 8)

add_pic(s, ASSETS / "BIZ_warning_runway.png", Inches(0.45), Inches(1.30),
        w=Inches(7.35), border=False)

rx = Inches(8.10)
add_text(s, rx, Inches(1.34), Inches(4.75), Inches(0.3),
         [("FROM RISK STAGE TO OPERATIONAL RESPONSE", 10.5, NAVY, True)])
ladder = [
    ("GREEN", "Normal operation", "Weekly monitoring continues in the "
     "background; no workshop action.", GREEN_D, GREEN_L),
    ("AMBER", "Early degradation", "Plan ahead: reserve a service slot, "
     "stage spares, add checks to the next visit.", AMBER_D, AMBER_L),
    ("RED", "High risk", "Inspect starter and battery circuit within 2–4 "
     "weeks; vehicle typically within ~10 weeks of failure.", RED_D, RED_L),
    ("CRITICAL", "Battery-cascade alert active", "Immediate intervention — "
     "this alert produced zero false alarms in validation.", WHITE, NAVY),
]
ly = Inches(1.74)
for name, sub, txt, col, fill in ladder:
    add_rect(s, rx, ly, Inches(4.78), Inches(1.02), fill)
    tcol = WHITE if fill == NAVY else col
    add_text(s, rx + Inches(0.18), ly + Inches(0.09), Inches(4.4),
             Inches(0.26), [(f"{name}  —  {sub}", 10.5, tcol, True)])
    add_text(s, rx + Inches(0.18), ly + Inches(0.38), Inches(4.45),
             Inches(0.58),
             [(txt, 9, WHITE if fill == NAVY else INK, False)], space_after=0)
    ly += Inches(1.14)

footer(s, "12 of 13 detected failures gave more than 3 months of notice — "
          "enough to plan parts, slots and vehicle rotation.",
       "Build pilot SOPs around the runway: alert → diagnose → schedule → "
       "replace before failure.")

# ===========================================================================
# SLIDE 9 - Business outcomes & strategic recommendations
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Business outcomes & strategic recommendations",
       "Deploy what is validated; grow what the data supports", 9)

col_y = Inches(1.34)
col_h = Inches(3.42)

# Column 1 - validated effectiveness
c1x, c1w = Inches(0.45), Inches(4.15)
add_rect(s, c1x, col_y, c1w, col_h, WHITE, line_color=LINE, line_w=1.0)
add_rect(s, c1x, col_y, c1w, Inches(0.07), NAVY)
add_text(s, c1x + Inches(0.2), col_y + Inches(0.14), c1w - Inches(0.4),
         Inches(0.3), [("PREDICTION EFFECTIVENESS — VALIDATED", 10, NAVY, True)])
rows = [
    ("Coverage-first mode", "13 of 14 failures caught; 5 of 20 healthy "
     "flagged for inspection"),
    ("Precision-first mode", "10 of 14 failures caught; only 2 of 20 healthy "
     "flagged"),
    ("Risk-ranking strength", "AUROC 0.93 (95% CI 0.81–0.99), independently "
     "validated"),
    ("Early warning", "Median 168 days; battery-cascade confirmation median "
     "~9 weeks, zero false alarms"),
    ("Known limit", "1 of 14 failures gave no signal — the vehicle stopped "
     "transmitting 142 days before failing"),
]
ry = col_y + Inches(0.5)
for t, d in rows:
    add_multirun(s, c1x + Inches(0.2), ry, c1w - Inches(0.4), Inches(0.56), [
        [(t + ":  ", NAVY, True), (d, INK, False)]], size=9, space_after=0)
    ry += Inches(0.57)

# Column 2 - operational benefits
c2x, c2w = Inches(4.78), Inches(3.95)
add_rect(s, c2x, col_y, c2w, col_h, WHITE, line_color=LINE, line_w=1.0)
add_rect(s, c2x, col_y, c2w, Inches(0.07), GREEN_D)
add_text(s, c2x + Inches(0.2), col_y + Inches(0.14), c2w - Inches(0.4),
         Inches(0.3), [("OPERATIONAL BENEFITS", 10, NAVY, True)])
bens = [
    "Fewer roadside failures — degrading vehicles reach the workshop first",
    "Better workshop scheduling — median 5+ months of planning runway",
    "Smarter spares — battery vs. starter signature guides what to stock",
    "Higher fleet availability — planned repairs beat breakdown recovery",
    "Focused effort — 80% of the healthy fleet stays untouched",
    "Warranty insight — degradation evidence per vehicle, per week",
]
ry = col_y + Inches(0.5)
for b in bens:
    add_rect(s, c2x + Inches(0.2), ry + Inches(0.045), Inches(0.09),
             Inches(0.09), GREEN_D, shape=MSO_SHAPE.OVAL)
    add_text(s, c2x + Inches(0.42), ry, c2w - Inches(0.62), Inches(0.5),
             [(b, 9, INK, False)], space_after=0)
    ry += Inches(0.47)

# Column 3 - roadmap
c3x, c3w = Inches(8.93), Inches(3.95)
add_rect(s, c3x, col_y, c3w, col_h, WHITE, line_color=LINE, line_w=1.0)
add_rect(s, c3x, col_y, c3w, Inches(0.07), REDACC)
add_text(s, c3x + Inches(0.2), col_y + Inches(0.14), c3w - Inches(0.4),
         Inches(0.3), [("STRATEGIC ROADMAP", 10, NAVY, True)])
phases = [
    ("NOW — Pilot", "Deploy on the active fleet; validate alerts against "
     "workshop findings; adopt tier-based SOPs", REDACC),
    ("NEXT — Integrate", "Service-centre integration, fleet dashboards, "
     "automated scheduling hooks", NAVY),
    ("THEN — Scale", "Enterprise rollout across the BharatBenz fleet; richer "
     "vehicle data feeds — identified during validation — to push detection "
     "coverage and precision further", NAVY2),
]
ry = col_y + Inches(0.5)
for t, d, col in phases:
    add_rect(s, c3x + Inches(0.2), ry + Inches(0.02), Inches(0.10),
             Inches(0.82), col)
    add_text(s, c3x + Inches(0.42), ry, c3w - Inches(0.62), Inches(0.26),
             [(t, 9.5, col, True)])
    add_text(s, c3x + Inches(0.42), ry + Inches(0.24), c3w - Inches(0.62),
             Inches(0.66), [(d, 8.7, INK, False)], space_after=0)
    ry += Inches(0.94)

# Closing message band
bz = Inches(4.98)
add_rect(s, Inches(0.45), bz, Inches(12.43), Inches(1.34), NAVY)
add_rect(s, Inches(0.45), bz, Inches(0.10), Inches(1.34), REDACC)
add_text(s, Inches(0.78), bz + Inches(0.14), Inches(11.9), Inches(0.28),
         [("EXECUTIVE CLOSING MESSAGE", 9.5, RGBColor(0x9F, 0xB3, 0xC8), True)])
add_text(s, Inches(0.78), bz + Inches(0.44), Inches(11.9), Inches(0.85),
         [("V1.1 demonstrates that predictive vehicle intelligence moves DICV "
           "from reactive repair to proactive, data-informed fleet reliability "
           "management — improving operational continuity, maintenance "
           "planning and customer satisfaction, while the underlying "
           "predictive technology remains fully protected.", 12.5, WHITE,
           True)], space_after=0)

footer(s, "The capability is validated, the operating model is defined, and "
          "the scaling path is clear.",
       "Approve the pilot and nominate the first operating fleet and service "
       "centres.")

prs.save(OUT)
print(f"Saved {OUT}")
print(f"Slides: {len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)}")
