"""
build_internal_technical_deck.py - INTERNAL technical deep-dive deck for the
Starter Motor V1.1 iteration, aimed at senior engineers for work approval.

Full technical disclosure (internal only): model names, validation protocol,
feature definitions, audit gates, alert rules, limitations. Every number is
traced to validated V1.1 artifacts (model_spec.json, gates.json, alert files,
horizon curve, reports).

Output: presentation/SM_V1_1_Technical_DeepDive_Internal.pptx
Run:    py -3 "STARTER MOTOR/V1.1/presentation/build_internal_technical_deck.py"
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(r"D:\Daimler-starter_motor_alternator_battery\STARTER MOTOR\V1.1")
ASSETS = ROOT / "presentation" / "assets"
GRAPHS = ROOT / "graphs"
OUT = ROOT / "presentation" / "SM_V1_1_Technical_DeepDive_Internal.pptx"

NAVY = RGBColor(0x14, 0x26, 0x3A)
NAVY2 = RGBColor(0x1F, 0x3A, 0x5F)
INK = RGBColor(0x21, 0x2B, 0x36)
SLATE = RGBColor(0x5B, 0x6B, 0x7C)
LINE = RGBColor(0xD7, 0xDE, 0xE5)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x0E, 0x6E, 0x6E)
REDACC = RGBColor(0xC8, 0x10, 0x2E)
GREEN_D = RGBColor(0x1B, 0x7A, 0x3D)
GREEN_L = RGBColor(0xDD, 0xEF, 0xE2)
AMBER_D = RGBColor(0xB4, 0x53, 0x09)
AMBER_L = RGBColor(0xFB, 0xEA, 0xD2)
RED_D = RGBColor(0x8B, 0x00, 0x00)
RED_L = RGBColor(0xFA, 0xE3, 0xE3)
BLUE_D = RGBColor(0x15, 0x65, 0xC0)
BLUE_L = RGBColor(0xE3, 0xEC, 0xF7)
PURP_D = RGBColor(0x6C, 0x33, 0x83)
PURP_L = RGBColor(0xEF, 0xE6, 0xF4)

FONT = "Segoe UI"
MONO = "Consolas"
SW, SH = Inches(13.333), Inches(7.5)
N_SLIDES = 25

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.75,
             shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=2, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        text, size, color, bold = para[:4]
        italic = para[4] if len(para) > 4 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return tb


def add_multirun(slide, x, y, w, h, runs, size=10.5, align=PP_ALIGN.LEFT,
                 space_after=4, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runlist in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        for text, color, bold in runlist:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = font
            f.size = Pt(size)
            f.color.rgb = color
            f.bold = bold
    return tb


def add_pic(slide, path, x, y, w=None, h=None, border=True):
    img = Image.open(path)
    ar = img.width / img.height
    if w is not None and h is None:
        h = Emu(int(w / ar))
    elif h is not None and w is None:
        w = Emu(int(h * ar))
    pic = slide.shapes.add_picture(str(path), x, y, w, h)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    return pic


def header(slide, kicker, title, num):
    add_rect(slide, 0, 0, Inches(0.12), SH, NAVY)
    add_rect(slide, Inches(0.12), 0, Inches(0.10), Inches(0.62), TEAL)
    add_text(slide, Inches(0.45), Inches(0.22), Inches(10.6), Inches(0.28),
             [(f"INTERNAL · TECHNICAL REVIEW  —  {kicker.upper()}", 10, SLATE,
               True)])
    add_text(slide, Inches(0.45), Inches(0.47), Inches(11.4), Inches(0.62),
             [(title, 23, NAVY, True)])
    add_text(slide, Inches(12.15), Inches(0.24), Inches(0.95), Inches(0.3),
             [(f"{num:02d} / {N_SLIDES:02d}", 10, SLATE, False)],
             align=PP_ALIGN.RIGHT)


def footer(slide, takeaway, review):
    y = Inches(6.62)
    h = Inches(0.58)
    add_rect(slide, Inches(0.45), y, Inches(7.35), h, LIGHT)
    add_rect(slide, Inches(0.45), y, Inches(0.07), h, NAVY)
    add_text(slide, Inches(0.68), y + Inches(0.06), Inches(7.0), Inches(0.2),
             [("ENGINEERING TAKEAWAY", 8, SLATE, True)])
    add_text(slide, Inches(0.68), y + Inches(0.24), Inches(7.0), Inches(0.32),
             [(takeaway, 9.5, NAVY, True)])
    add_rect(slide, Inches(7.95), y, Inches(4.93), h, LIGHT)
    add_rect(slide, Inches(7.95), y, Inches(0.07), h, TEAL)
    add_text(slide, Inches(8.18), y + Inches(0.06), Inches(4.6), Inches(0.2),
             [("REVIEW / DECISION POINT", 8, SLATE, True)])
    add_text(slide, Inches(8.18), y + Inches(0.24), Inches(4.6), Inches(0.32),
             [(review, 9.5, TEAL, True)])
    add_text(slide, Inches(0.45), Inches(7.26), Inches(12.45), Inches(0.2),
             [("Starter Motor Failure Prediction V1.1  ·  Internal engineering "
               "review — full technical disclosure  ·  Not for external "
               "distribution", 8, SLATE, False)])


def chipbox(slide, x, y, w, h, title, fill, edge, title_color=None):
    add_rect(slide, x, y, w, h, fill, line_color=edge, line_w=1.0)
    add_text(slide, x + Inches(0.14), y + Inches(0.08), w - Inches(0.28),
             Inches(0.26), [(title, 10, title_color or edge, True)])


# ===========================================================================
# SLIDE 1 - Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 0, SW, Inches(0.14), TEAL)
add_rect(s, Inches(0.9), Inches(1.95), Inches(0.14), Inches(1.62), TEAL)
add_text(s, Inches(0.9), Inches(1.22), Inches(11.5), Inches(0.35),
         [("INTERNAL — SENIOR ENGINEERING REVIEW FOR WORK APPROVAL", 13,
           RGBColor(0x9F, 0xB3, 0xC8), True)])
add_text(s, Inches(1.28), Inches(1.95), Inches(11.4), Inches(1.7),
         [("Starter Motor Failure Prediction V1.1", 40, WHITE, True),
          ("Technical Deep-Dive: Data, Methods, Validation, Results", 22,
           RGBColor(0xC9, 0xD6, 0xE3), False)], space_after=6)
add_text(s, Inches(1.28), Inches(3.75), Inches(11.2), Inches(0.4),
         [("Nested leave-one-vehicle-out RidgeClassifier pipeline · 34-truck "
           "fleet · 106.4M telemetry rows · frozen 4-feature model",
           13, RGBColor(0x9F, 0xB3, 0xC8), False)])
stats = [("0.9321", "nested-LOVO AUROC [0.811, 0.986]"),
         ("p = 0.005", "full-pipeline permutation test"),
         ("13 / 14", "recall @ per-fold Youden"),
         ("k* = 10 wk", "validated prequential horizon")]
x = Inches(1.28)
for big, lab in stats:
    add_text(s, x, Inches(4.95), Inches(2.75), Inches(0.5),
             [(big, 24, WHITE, True)])
    add_text(s, x, Inches(5.48), Inches(2.75), Inches(0.55),
             [(lab, 10, RGBColor(0x9F, 0xB3, 0xC8), False)])
    x += Inches(2.87)
add_text(s, Inches(1.28), Inches(6.65), Inches(10.8), Inches(0.3),
         [("July 2026  ·  V1.1 frozen 2026-06-10  ·  artifacts: STARTER "
           "MOTOR/V1.1/{results, reports, audit, discovery, graphs}", 10,
           RGBColor(0x9F, 0xB3, 0xC8), False)])

# ===========================================================================
# SLIDE 2 - Scope & approval asks
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Scope", "What V1.1 delivers, and what we are asking you to "
       "approve", 2)

chipbox(s, Inches(0.45), Inches(1.35), Inches(6.0), Inches(2.5),
        "SCOPE OF THE V1.1 ITERATION", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(1.75), Inches(5.6), Inches(2.0), [
    [("Objective:  ", NAVY, True),
     ("weekly per-vehicle starter-motor failure-risk classification on the "
      "34-truck SM fleet (14 failed / 20 non-failed).", INK, False)],
    [("Method:  ", NAVY, True),
     ("nested leave-one-vehicle-out (LOVO) RidgeClassifier(alpha=1.0), "
      "fold-internal imputation/scaling/selection, per-fold Platt "
      "recalibration, pre-registered thresholds.", INK, False)],
    [("Deliverables:  ", NAVY, True),
     ("recalibrated probabilities + GREEN/AMBER/RED tiers, 3-channel alert "
      "layer, 10-week validity horizon, 34 per-VIN dashboards, full audit "
      "trail.", INK, False)],
    [("Explicit non-deliverable:  ", NAVY, True),
     ("day-level RUL / failure-date forecasts (all survival variants lose "
      "to a constant — evidence on slide 17).", INK, False)],
], size=10, space_after=6)

chipbox(s, Inches(0.45), Inches(4.05), Inches(6.0), Inches(2.4),
        "WHY V1.1 EXISTS (V1 POST-MORTEM)", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(4.45), Inches(5.6), Inches(1.9), [
    [("V1 reported 0.9214 without nested selection; honest restatement is "
      "0.8929 (selection optimism +0.0285).", INK, False)],
    [("V1's vsi_dominant_freq was exposed as a 1/n_weeks observation-length "
      "artifact (fake TP on VIN2_F) - banned in V1.1.", INK, False)],
    [("V1 missed VIN8_F (P 0.303). V1.1 recovers it (P 0.716, RED) via the "
      "battery-step-aware rest-VSI baseline.", INK, False)],
], size=10, space_after=6)

chipbox(s, Inches(6.75), Inches(1.35), Inches(6.13), Inches(5.1),
        "APPROVAL ASKS (THIS REVIEW)", NAVY, NAVY, WHITE)
asks = [
    ("1", "Freeze the modal 4-feature nested model as the V1.1 production "
     "scorer (spec: V1_1_SM_model_spec.json)."),
    ("2", "Ship recalibrated probabilities + tier rule (GREEN < 0.35, AMBER "
     "0.35-0.55, RED >= 0.55) — calibration gate passed (slope 0.86)."),
    ("3", "Adopt the 3-channel alert policy: A2 battery-cascade as pager, "
     "persistence as terminal condition flag, A1 as corroborator."),
    ("4", "Approve weekly scoring cadence (monthly gives only 2 reads inside "
     "the 10-week horizon)."),
    ("5", "Endorse instrumentation requests: >=1 Hz crank logging, cranking "
     "current / battery SoC, maintenance records, more failure cases."),
]
ay = Inches(1.85)
for n, t in asks:
    add_rect(s, Inches(6.98), ay, Inches(0.30), Inches(0.30), TEAL,
             shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(6.98), ay + Inches(0.02), Inches(0.30), Inches(0.26),
             [(n, 11, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.42), ay - Inches(0.02), Inches(5.3), Inches(0.85),
             [(t, 10, WHITE, False)], space_after=0)
    ay += Inches(0.88)

footer(s, "V1.1 is a frozen, fully audited iteration; every claim in this "
          "deck is reproducible from committed artifacts.",
       "Confirm the five asks or flag which require deeper review.")

# ===========================================================================
# SLIDE 3 - Data foundation & data reality
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Data", "Data foundation — and the data reality the design had to "
       "survive", 3)

chipbox(s, Inches(0.45), Inches(1.35), Inches(6.0), Inches(2.32),
        "CORPUS", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(1.75), Inches(5.6), Inches(1.85), [
    [("34 trucks (14 F / 20 NF), _SM suffix — fully independent of the ALT "
      "fleet (no cross-VIN identity).", INK, False)],
    [("106,445,161 raw 5-second rows after null-timestamp drop; 17,522 "
      "active vin-days.", INK, False)],
    [("Reduced to 2,636 truck-weeks (weekly cache) + 20,471 crank events "
      "(event extraction, artifact-filtered).", INK, False)],
    [("Signals: VSI, SMA, RPM, CSP, ANR, GED. No current, no temperature, "
      "no battery SoC. VSI quantized at 0.2 V.", INK, False)],
], size=10, space_after=5)

chipbox(s, Inches(0.45), Inches(3.88), Inches(6.0), Inches(2.55),
        "DATA REALITY (DESIGN CONSTRAINTS)", WHITE, LINE, RED_D)
add_multirun(s, Inches(0.65), Inches(4.28), Inches(5.6), Inches(2.1), [
    [("SMA-dead cohort (7 trucks):  ", RED_D, True),
     ("VIN8/9_F + 5 NF have <1% SMA coverage - all crank features NaN, "
      "fold-internal median-imputed.", INK, False)],
    [("Terminal silent gaps (5 failed):  ", RED_D, True),
     ("VIN1 72d, VIN4 97d, VIN5 32d, VIN8 37d, VIN9 142d - no telemetry "
      "before failure.", INK, False)],
    [("Leakage ceilings:  ", RED_D, True),
     ("observation length ALONE classifies at AUROC 0.952 (n_weeks) / 0.893 "
      "(t_start). Every feature must beat the L40 fixed-window control to "
      "prove it is not a length proxy.", INK, False)],
], size=10, space_after=6)

chipbox(s, Inches(6.75), Inches(1.35), Inches(6.13), Inches(5.1),
        "SIGNAL DICTIONARY (AS USED)", WHITE, LINE, NAVY)
sig_rows = [
    ("VSI", "Power-supply voltage (V)", "0-36 valid; sentinels 0/255; "
     "x0.2 scaling; drive vs rest split"),
    ("SMA", "Starter Motor Active {0,1}", "crank-event extraction basis; "
     "dead on 7 trucks"),
    ("RPM", "Engine speed (rev/min)", "sentinel 65535; crank success "
     "criterion rpm_max_15s >= 550"),
    ("CSP", "Vehicle speed (km/h)", "sentinel 65535; drive-state gating"),
    ("ANR", "Engine torque (Nm)", "sentinels 65535 / -5000; load context"),
    ("GED", "Alternator excitation state {0..3}", "no GED=2 events in "
     "failed SM fleet (unlike ALT)"),
]
ry = Inches(1.82)
for name, meaning, use in sig_rows:
    add_rect(s, Inches(6.98), ry, Inches(0.78), Inches(0.62), BLUE_L)
    add_text(s, Inches(6.98), ry + Inches(0.16), Inches(0.78), Inches(0.3),
             [(name, 11, BLUE_D, True)], align=PP_ALIGN.CENTER, font=MONO)
    add_text(s, Inches(7.92), ry + Inches(0.015), Inches(4.85), Inches(0.24),
             [(meaning, 9.5, NAVY, True)])
    add_text(s, Inches(7.92), ry + Inches(0.26), Inches(4.85), Inches(0.38),
             [(use, 8.5, SLATE, False)], space_after=0)
    ry += Inches(0.75)

footer(s, "The pipeline is designed around missingness and length-leakage — "
          "not around an idealized dataset.",
       "Confirm the SMA-dead imputation and L40 anti-leak strategy are "
       "acceptable engineering practice.")

# ===========================================================================
# SLIDE 4 - Pipeline architecture flow
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Architecture", "End-to-end pipeline: raw telemetry to audited "
       "risk artifacts", 4)


def flow_card(slide, x, y, w, h, title, sub, fill, edge, tcol=None):
    add_rect(slide, x, y, w, h, fill, line_color=edge, line_w=1.0)
    add_text(slide, x + Inches(0.10), y + Inches(0.07), w - Inches(0.2),
             Inches(0.4), [(title, 9.5, tcol or edge, True)], space_after=0)
    add_text(slide, x + Inches(0.10), y + Inches(0.47), w - Inches(0.2),
             h - Inches(0.52), [(sub, 8, INK, False)], space_after=0)


def harrow(slide, x, y):
    add_rect(slide, x, y, Inches(0.20), Inches(0.16), NAVY,
             shape=MSO_SHAPE.CHEVRON)


row1 = [
    ("1 · RAW TELEMETRY", "2 SM-fleet CSVs, 106.4M rows @5 s\nVSI SMA RPM "
     "CSP ANR GED", BLUE_L, BLUE_D),
    ("2 · CLEANING", "sentinel handling (65535 / -5000 / 0 / 255)\nVSI x0.2 "
     "scaling, null-ts drop", BLUE_L, BLUE_D),
    ("3 · WEEKLY CACHE", "2,636 truck-weeks\nV1_SM_weekly_*.parquet",
     BLUE_L, BLUE_D),
    ("4 · CRANK EVENTS", "20,471 events, artifact-filtered\n"
     "V1_SM_crank_events.parquet", BLUE_L, BLUE_D),
    ("5 · FEATURE MATRIX", "34 x 10 candidates, L40-anchored\n+ admissibility "
     "audit (reject rule)", PURP_L, PURP_D),
]
row2 = [
    ("6 · NESTED LOVO MODEL", "RidgeClassifier(alpha=1.0)\n34 outer folds, "
     "inner select k=3-6", PURP_L, PURP_D),
    ("7 · RECALIBRATION", "per-fold Platt on inner-OOF\nper-fold Youden "
     "threshold", PURP_L, PURP_D),
    ("8 · OOF RISK + TIERS", "prob_recal per VIN\nGREEN <0.35 | AMBER | "
     "RED >=0.55", GREEN_L, GREEN_D),
    ("9 · ALERT LAYER", "persistence / A1 burst / A2 cascade\nalert_"
     "validation + policy CSVs", AMBER_L, AMBER_D),
    ("10 · HORIZON + OUTPUTS", "prequential k*=10 wk\n34 dashboards, reports, "
     "gates.json", GREEN_L, GREEN_D),
]
cw, ch = Inches(2.34), Inches(1.06)
gx = Inches(0.13)
y1, y2 = Inches(1.42), Inches(2.86)
for i, (t, sub, f, e) in enumerate(row1):
    x = Inches(0.45) + i * (cw + gx)
    flow_card(s, x, y1, cw, ch, t, sub, f, e)
    if i < 4:
        harrow(s, x + cw - Inches(0.03), y1 + ch / 2 - Inches(0.08))
for i, (t, sub, f, e) in enumerate(row2):
    x = Inches(0.45) + i * (cw + gx)
    flow_card(s, x, y2, cw, ch, t, sub, f, e)
    if i < 4:
        harrow(s, x + cw - Inches(0.03), y2 + ch / 2 - Inches(0.08))
# wrap indicator 5 -> 6 (text only; a drawn bar would touch the wrong cards)
add_text(s, Inches(11.35), y1 + ch + Inches(0.04), Inches(1.5), Inches(0.3),
         [("continues at stage 6 ↓", 8.5, SLATE, True)],
         align=PP_ALIGN.RIGHT)

chipbox(s, Inches(0.45), Inches(4.14), Inches(12.43), Inches(2.28),
        "AUDIT SPINE (RUNS ACROSS EVERY STAGE)", WHITE, LINE, NAVY)
audits = [
    ("A · Data-quality audit", "labels, null structure, sampling gaps, VSI "
     "sensor probes (probe1-7)", "audit/A_data_quality_audit.md"),
    ("B · Feature audit", "candidate screening, incremental LOVO, "
     "truncation control, model variants", "audit/B_feature_audit.md"),
    ("C · Model audit", "nested folds, jackknife, permutation, "
     "C1_audit_results.json", "audit/C_model_audit.md"),
    ("D-G · Physics & discovery", "failure-physics brief (D); clustering/"
     "archetypes, survival & fleet-clock, sequence probes (E-G)",
     "audit/D + discovery/E,F,G"),
    ("Gates G1-G6", "L40 control, proxy audit, calibration, stability, "
     "jackknife, token scan", "results/V1_1_SM_gates.json"),
]
ax_ = Inches(0.65)
aw = Inches(2.36)
for t, d, f in audits:
    add_rect(s, ax_, Inches(4.62), aw, Inches(1.62), LIGHT)
    add_text(s, ax_ + Inches(0.10), Inches(4.72), aw - Inches(0.2),
             Inches(0.3), [(t, 9, NAVY, True)])
    add_text(s, ax_ + Inches(0.10), Inches(5.04), aw - Inches(0.2),
             Inches(0.78), [(d, 8, INK, False)], space_after=0)
    add_text(s, ax_ + Inches(0.10), Inches(5.88), aw - Inches(0.2),
             Inches(0.3), [(f, 7.5, SLATE, False, True)], space_after=0,
             font=MONO)
    ax_ += aw + Inches(0.055)

footer(s, "Ten pipeline stages, five audit workstreams — every artifact "
          "named, versioned and re-runnable with py -3.",
       "Architecture sign-off: any stage requiring design changes before "
       "pilot integration?")

# ===========================================================================
# SLIDE 5 - Feature engineering
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Features", "Candidate pool: 10 admissible features, 3 physics "
       "families, 1 hard reject rule", 5)

fam_specs = [
    ("VOLTAGE-VOLATILITY FAMILY (4)", BLUE_D, BLUE_L, [
        ("vsi_std_ratio_30d_L40", "between-week drive-VSI volatility, last 4 "
         "wks vs fixed L40 window"),
        ("vsi_withinwk_std_ratio_30d_w", "within-week drive-VSI noise ratio, "
         "30 d vs L40 baseline  [WINNER]"),
        ("vsi_range_trend", "Theil-Sen slope of weekly p95-p05 envelope, "
         "last 12 masked wks  [WINNER]"),
        ("vsi_trend_persistence", "|mean sign| of rolling 4-wk OLS slopes of "
         "vsi_drive_mean, last 12 wks"),
    ]),
    ("CRANK-EVENT FAMILY (4) — NaN on SMA-dead cohort", AMBER_D, AMBER_L, [
        ("failed_crank_rate_last90", "share of failed cranks among "
         "success-known events, last 90 d"),
        ("retry_burst_rate_last90", "burst episodes (>=2 events/10 min, no "
         "intervening success) per active day"),
        ("extended_crank_tail_rate_last90", "P(n_rows>=2) last 90 d minus "
         "L40 baseline share"),
        ("first_crank_fail_rate_last90", "first-crank-after-6h-rest failure "
         "share, last 90 d"),
    ]),
    ("BATTERY-STATE FAMILY (2)", GREEN_D, GREEN_L, [
        ("rest_vsi_p05_delta90", "rest-VSI p05 shift vs baseline, "
         "RE-BASELINED after detected battery-replacement step (>= +0.5 V, "
         "SNR >= 2)  [WINNER]"),
        ("dip_depth_last90_delta", "crank dip-depth deepening, last 90 d vs "
         "L40 baseline  [WINNER]"),
    ]),
]
fx = Inches(0.45)
fw = Inches(4.02)
for title, edge, fill, items in fam_specs:
    add_rect(s, fx, Inches(1.35), fw, Inches(3.9), WHITE, line_color=LINE,
             line_w=1.0)
    add_rect(s, fx, Inches(1.35), fw, Inches(0.34), fill)
    add_text(s, fx + Inches(0.12), Inches(1.41), fw - Inches(0.24),
             Inches(0.26), [(title, 8.5, edge, True)])
    iy = Inches(1.82)
    for name, desc in items:
        add_text(s, fx + Inches(0.14), iy, fw - Inches(0.28), Inches(0.22),
                 [(name, 8.5, NAVY, True)], font=MONO, space_after=0)
        add_text(s, fx + Inches(0.14), iy + Inches(0.21), fw - Inches(0.28),
                 Inches(0.56), [(desc, 8, SLATE, False)], space_after=0)
        iy += Inches(0.82)
    fx += fw + Inches(0.185)

chipbox(s, Inches(0.45), Inches(5.42), Inches(12.43), Inches(1.05),
        "ADMISSIBILITY RULE + THE BANNED ARTIFACT", WHITE, LINE, RED_D)
add_multirun(s, Inches(0.65), Inches(5.80), Inches(12.0), Inches(0.62), [
    [("REJECT iff max |Spearman r| vs {n_weeks, t_start, span} > 0.5 AND "
      "L40 fixed-window AUROC drop > 0.05.  ", INK, False),
     ("Cautionary tale: V1's vsi_dominant_freq = 1/n_weeks in disguise "
      "(AUROC 0.748 raw -> 0.525 under L40) — banned from the V1.1 pool.",
      RED_D, True)],
], size=9.5, space_after=0)

footer(s, "All 10 candidates are window-anchored by construction; the L40 "
          "control is the anti-leak backbone.",
       "Approve the admissibility rule as the standing bar for all future "
       "feature proposals.")

# ===========================================================================
# SLIDE 6 - Frozen winner subset & stability
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Model inputs", "The frozen 4-feature subset: selection stability "
       "and one disclosed suppressor", 6)

add_pic(s, ASSETS / "TECH_feature_freq.png", Inches(0.45), Inches(1.40),
        w=Inches(7.5), border=False)

chipbox(s, Inches(0.45), Inches(4.20), Inches(7.5), Inches(2.22),
        "G4 WINNER-STABILITY GATE — STRICT FAIL, DISCLOSED", WHITE, LINE,
        AMBER_D)
add_multirun(s, Inches(0.65), Inches(4.60), Inches(7.1), Inches(1.75), [
    [("Only 3 distinct subsets ever chosen across 34 folds: the modal "
      "4-feature set (14), its 3-feature sibling without dip_depth (14), "
      "and a 3-feature sibling without rest_p05 (6).", INK, False)],
    [("Strict criterion needs the modal subset in >=17/34 folds; 14/34 "
      "fails it. Substantive stability is strong (core pair in 34/34), but "
      "the gate is reported as FAIL — not smoothed over.", INK, False)],
], size=9.5, space_after=5)

chipbox(s, Inches(8.25), Inches(1.40), Inches(4.63), Inches(5.02),
        "PHYSICS READING & THE SUPPRESSOR", WHITE, LINE, NAVY)
add_multirun(s, Inches(8.45), Inches(1.82), Inches(4.25), Inches(4.5), [
    [("withinwk_std_ratio  ", BLUE_D, True),
     ("- the workhorse: sustained electrical noise rises before failure.",
      INK, False)],
    [("rest_vsi_p05_delta90  ", GREEN_D, True),
     ("- battery floor sag; re-baselining after replacement steps is what "
      "recovered VIN8_F.", INK, False)],
    [("dip_depth_last90_delta  ", GREEN_D, True),
     ("- crank load signature deepening (battery-cascade component).",
      INK, False)],
    [("vsi_range_trend  ", AMBER_D, True),
     ("- SUPPRESSOR: physics-positive (widening envelope = risk) but "
      "enters with a NEGATIVE multivariate weight, driven by r = +0.82 "
      "collinearity with the noise ratio. Per-VIN explanation cards state "
      "the physical direction and the model usage separately - do not read "
      "its coefficient as physics.", INK, False)],
], size=9.5, space_after=7)

footer(s, "Two features carry every fold; the suppressor is a documented "
          "modeling artifact, not a physics claim.",
       "Accept G4 strict-fail with substantive-stability evidence, or "
       "request a fixed-subset re-run.")

# ===========================================================================
# SLIDE 7 - Validation protocol diagram
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Validation", "Nested LOVO: selection, threshold and calibration "
       "all inside the fold", 7)

flow_card(s, Inches(0.45), Inches(1.45), Inches(2.6), Inches(1.30),
          "OUTER LOOP (x34)", "hold out truck i entirely\n(features, "
          "selection, scaling,\nimputation never see it)", BLUE_L, BLUE_D)
flow_card(s, Inches(3.35), Inches(1.45), Inches(3.45), Inches(1.30),
          "INNER SCREENING (33 trucks)", "Mann-Whitney p < 0.1 · AUROC >= "
          "0.6\nSpearman dedup |rho| <= 0.85\nstability >= 0.8 · pool cap 10",
          PURP_L, PURP_D)
flow_card(s, Inches(7.10), Inches(1.45), Inches(3.15), Inches(1.30),
          "INNER SUBSET SEARCH", "exhaustive k = 3..6 subsets\ninner-LOVO "
          "AUROC picks winner\n(modal 4 in 14/34 folds; 3 subsets ever)",
          PURP_L, PURP_D)
flow_card(s, Inches(10.55), Inches(1.45), Inches(2.33), Inches(1.30),
          "REFIT + SCORE", "RidgeClassifier(alpha=1.0)\nclosed-form replica "
          "(<1e-9 vs sklearn)\nscore held-out truck i", GREEN_L, GREEN_D)
harrow(s, Inches(3.07), Inches(2.02))
harrow(s, Inches(6.82), Inches(2.02))
harrow(s, Inches(10.27), Inches(2.02))

flow_card(s, Inches(0.45), Inches(3.05), Inches(4.0), Inches(1.22),
          "PER-FOLD POST-PROCESSING", "threshold = inner-OOF Youden "
          "(pre-registered rule)\nrecalibration = Platt on inner-OOF "
          "decision values", AMBER_L, AMBER_D)
flow_card(s, Inches(4.75), Inches(3.05), Inches(4.0), Inches(1.22),
          "AGGREGATE OOF", "34 out-of-fold recalibrated probabilities\n"
          "nested AUROC = 0.9321 (optimism vs non-nested: +0.0036)",
          GREEN_L, GREEN_D)
flow_card(s, Inches(9.05), Inches(3.05), Inches(3.83), Inches(1.22),
          "INFERENCE", "bootstrap N=200 (seed 42): CI [0.8107, 0.9861]\n"
          "full-pipeline permutation N=200 (seed 43): p = 0.005",
          GREEN_L, GREEN_D)
harrow(s, Inches(4.47), Inches(3.56))
harrow(s, Inches(8.77), Inches(3.56))

chipbox(s, Inches(0.45), Inches(4.55), Inches(12.43), Inches(1.85),
        "WHY THIS PROTOCOL (WHAT IT CLOSES OFF)", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(4.95), Inches(12.0), Inches(1.4), [
    [("Selection leakage:  ", NAVY, True),
     ("feature screening and subset choice re-run inside every fold — the "
      "held-out truck never influences its own scorer (this alone is worth "
      "-0.0285 on V1's reported number).", INK, False)],
    [("Threshold leakage:  ", NAVY, True),
     ("Youden point comes from inner-OOF only; no post-hoc threshold "
      "tuning on test outcomes.", INK, False)],
    [("Calibration leakage:  ", NAVY, True),
     ("Platt fitted per fold on inner-OOF decision values; recalibrated "
      "probabilities are honestly out-of-fold.", INK, False)],
], size=9.5, space_after=5)

footer(s, "Everything that could overfit is re-derived 34 times inside the "
          "folds; the permutation test validates the WHOLE pipeline.",
       "Protocol sign-off: is nested-LOVO accepted as the standing "
       "evaluation bar for SM iterations?")

# ===========================================================================
# SLIDE 8 - Headline results (bridge)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Results", "Headline: nested AUROC 0.9321 — and where the gain "
       "actually comes from", 8)

add_pic(s, ASSETS / "TECH_v1_bridge.png", Inches(0.45), Inches(1.42),
        w=Inches(8.1), border=False)

chipbox(s, Inches(8.85), Inches(1.42), Inches(4.03), Inches(5.0),
        "HEADLINE NUMBERS", WHITE, LINE, NAVY)
hl = [
    ("0.9321", "nested-LOVO AUROC, 34 OOF probabilities", NAVY),
    ("[0.811, 0.986]", "bootstrap 95% CI (N=200, seed 42) - small-n width, "
     "read honestly", NAVY),
    ("p = 0.005", "permutation test: 0/200 full-pipeline label shuffles "
     "reach 0.9321 (N=200 floor)", NAVY),
    ("+0.0036", "nesting optimism vs non-nested modal 0.9357 - selection "
     "bias is tiny in V1.1", GREEN_D),
    ("+0.089", "V1.1 features vs V1-era 22 features under the SAME nested "
     "protocol (0.9321 vs 0.8429): the gain is features, not protocol",
     GREEN_D),
]
ry = Inches(1.88)
for big, sub, col in hl:
    add_text(s, Inches(9.05), ry, Inches(3.7), Inches(0.35),
             [(big, 17, col, True)])
    add_text(s, Inches(9.05), ry + Inches(0.34), Inches(3.65), Inches(0.5),
             [(sub, 8.5, SLATE, False)], space_after=0)
    ry += Inches(0.90)

footer(s, "V1.1 beats the honestly-restated V1 by +0.039 AUROC and includes "
          "VIN8_F, V1's worst miss.",
       "Accept 0.9321 [0.811, 0.986] as the frozen reference number for all "
       "downstream claims.")

# ===========================================================================
# SLIDE 9 - Operating points & confusion matrices
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Operating points", "Two pre-registered operating points — choose "
       "per maintenance economics", 9)


def confusion(slide, x, y, title, tp, fn, fp, tn, sub):
    cw_, ch_ = Inches(1.55), Inches(0.92)
    add_text(slide, x, y, Inches(4.3), Inches(0.3), [(title, 12, NAVY, True)])
    add_text(slide, x + Inches(1.15), y + Inches(0.40), cw_, Inches(0.22),
             [("Flagged", 9, SLATE, True)], align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(1.15) + cw_ + Inches(0.06), y + Inches(0.40),
             cw_, Inches(0.22), [("Clear", 9, SLATE, True)],
             align=PP_ALIGN.CENTER)
    rows = [("Failed (14)", (tp, GREEN_L, GREEN_D, "caught"),
             (fn, RED_L, RED_D, "missed")),
            ("Healthy (20)", (fp, AMBER_L, AMBER_D, "false alarm"),
             (tn, GREEN_L, GREEN_D, "correct clear"))]
    yy = y + Inches(0.64)
    for lab, c1, c2 in rows:
        add_text(slide, x, yy + Inches(0.32), Inches(1.10), Inches(0.3),
                 [(lab, 9, SLATE, True)], align=PP_ALIGN.RIGHT)
        for j, (v, fill, edge, note) in enumerate((c1, c2)):
            cx = x + Inches(1.15) + j * (cw_ + Inches(0.06))
            add_rect(slide, cx, yy, cw_, ch_, fill, line_color=edge,
                     line_w=1.2)
            add_text(slide, cx, yy + Inches(0.12), cw_, Inches(0.42),
                     [(str(v), 22, edge, True)], align=PP_ALIGN.CENTER)
            add_text(slide, cx, yy + Inches(0.58), cw_, Inches(0.26),
                     [(note, 8, edge, False)], align=PP_ALIGN.CENTER)
        yy += ch_ + Inches(0.06)
    add_text(slide, x, yy + Inches(0.06), Inches(4.35), Inches(0.55),
             [(sub, 8.5, SLATE, False)], space_after=0)


confusion(s, Inches(0.75), Inches(1.50),
          "A · Per-fold Youden (coverage-first)", 13, 1, 5, 15,
          "recall 0.929 · specificity 0.75 · F1 0.8125 · MCC 0.6691. "
          "Threshold from inner-OOF Youden, per fold (pre-registered).")
confusion(s, Inches(5.55), Inches(1.50),
          "B · RED tier >= 0.55 (precision-first)", 10, 4, 2, 18,
          "recall 0.714 · specificity 0.90. RED false alarms: VIN19_NF "
          "(P 0.96), VIN34_NF (P 0.62) - plausibly right-censored degraders.")

chipbox(s, Inches(10.30), Inches(1.50), Inches(2.58), Inches(3.55),
        "NON-DOMINANCE", WHITE, LINE, AMBER_D)
add_multirun(s, Inches(10.48), Inches(1.92), Inches(2.25), Inches(3.05), [
    [("No point beats V1-restated (12/14 @ 18/20) on every cell:", INK,
      False)],
    [("Youden: +1 recall, -3 specificity", INK, False)],
    [("RED tier: -2 recall, equal specificity (18/20)", INK, False)],
    [("Both are honest, pre-registered; the choice is an economics "
      "decision, not a statistics one.", NAVY, True)],
], size=8.5, space_after=5)

chipbox(s, Inches(0.45), Inches(5.30), Inches(12.43), Inches(1.10),
        "THE ONE MISS AT YOUDEN", WHITE, LINE, RED_D)
add_multirun(s, Inches(0.65), Inches(5.68), Inches(12.0), Inches(0.65), [
    [("VIN9_F: raw score 0.401 vs fold threshold 0.406 — loses by 0.005 "
      "(recalibrated P 0.224, the GREEN bar on the fleet chart). SMA-dead "
      "(0.4% coverage), 142-day terminal silent gap, mid-fleet on every "
      "live feature. This truck is structurally unobservable in the current "
      "telemetry; no admissible model sees it (full anatomy on slide 16).",
      INK, False)],
], size=9.5, space_after=0)

footer(s, "13/14 @ 5/20 or 10/14 @ 2/20 — both validated; the trade is "
          "explicit and priced.",
       "Pick the default operating point for the pilot (recommend: RED tier "
       "pager + Youden watch-list).")

# ===========================================================================
# SLIDE 10 - Fleet OOF view
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Results", "Per-VIN out-of-fold risk: the full 34-truck picture", 10)

add_pic(s, GRAPHS / "V1_1_SM_fleet_risk.png", Inches(0.72), Inches(1.35),
        w=Inches(9.2))

chipbox(s, Inches(10.20), Inches(1.35), Inches(2.68), Inches(5.05),
        "READING NOTES", WHITE, LINE, NAVY)
add_multirun(s, Inches(10.38), Inches(1.77), Inches(2.35), Inches(4.55), [
    [("Tiers at end of history: RED 12 (10 F + 2 NF), AMBER 2 (NF), GREEN "
      "20 (4 F + 16 NF).", INK, False)],
    [("4 failed in GREEN:  ", RED_D, True),
     ("VIN1/3/4_F caught at Youden or by alert channels; VIN9_F is the "
      "structural miss.", INK, False)],
    [("2 NF in RED:  ", AMBER_D, True),
     ("VIN19_NF (0.96) and VIN34_NF (0.62) fire multiple channels - "
      "plausibly genuinely degrading (right-censored).", INK, False)],
    [("Display note: graphs renumber NF VINs +14 (VIN1_NF -> VIN15_NF); "
      "results CSVs keep original labels.", SLATE, False)],
], size=8.5, space_after=6)

footer(s, "The ranking is clean: 10 of the top 12 scores are true failures; "
          "the borderline zone is where alerts take over.",
       "Confirm the two RED healthy trucks are queued for physical "
       "inspection feedback.")

# ===========================================================================
# SLIDE 11 - Gates G1-G6
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Audit gates", "Six gates, reported as measured — including the "
       "one that fails", 11)

gates = [
    ("G1", "Fixed L40-window control", "recompute every feature with data "
     "clipped to last 40 masked weeks; AUROC drop = 0.0000", "PASS",
     GREEN_D, GREEN_L),
    ("G2", "OOF-vs-proxy audit", "Spearman: n_weeks -0.640 · t_start +0.507 "
     "· span -0.653 (all >|0.5| tripwire). Defense: label-mediated - G1 "
     "drop 0.0000 + prequential collapse at k=11 prove failure-locking, "
     "not length-locking", "REPORT", SLATE, LIGHT),
    ("G3", "Calibration", "Platt slope 0.86 (in [0.5, 2]) · Brier 0.124 vs "
     "constant 0.242 · CITL -0.062 -> probabilities shippable", "REPORT",
     SLATE, LIGHT),
    ("G4", "Winner stability (strict)", "modal subset 14/34 < 17 required; "
     "3 distinct subsets; core pair 34/34 - substantive stability strong",
     "FAIL", RED_D, RED_L),
    ("G5", "Jackknife (drop-1-truck)", "AUROC range 0.9269-0.9511, std "
     "0.007 - no single truck carries the result", "REPORT", SLATE,
     LIGHT),
    ("G6", "Banned-token scan", "0 hits for banned leakage tokens "
     "(vsi_dominant_freq family, epoch fields) in the frozen pipeline",
     "PASS", GREEN_D, GREEN_L),
]
gy = Inches(1.38)
for gid, name, detail, verdict, vc, vf in gates:
    add_rect(s, Inches(0.45), gy, Inches(12.43), Inches(0.80), WHITE,
             line_color=LINE, line_w=0.9)
    add_rect(s, Inches(0.45), gy, Inches(0.55), Inches(0.80), NAVY)
    add_text(s, Inches(0.45), gy + Inches(0.26), Inches(0.55), Inches(0.3),
             [(gid, 12, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.18), gy + Inches(0.08), Inches(2.55), Inches(0.6),
             [(name, 9.5, NAVY, True)], space_after=0)
    add_text(s, Inches(3.85), gy + Inches(0.08), Inches(7.55), Inches(0.66),
             [(detail, 8.5, INK, False)], space_after=0)
    add_rect(s, Inches(11.62), gy + Inches(0.18), Inches(1.05), Inches(0.44),
             vf, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(11.62), gy + Inches(0.27), Inches(1.05), Inches(0.28),
             [(verdict, 10, vc, True)], align=PP_ALIGN.CENTER)
    gy += Inches(0.86)

footer(s, "The audit is adversarial by design: G2's correlations and G4's "
          "strict fail are printed, defended, and left visible.",
       "Countersign the gate table (results/V1_1_SM_gates.json) as the "
       "audit record for V1.1.")

# ===========================================================================
# SLIDE 12 - Calibration & tiers
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Calibration", "Recalibrated probabilities are shippable — and "
       "drive a three-tier action rule", 12)

chipbox(s, Inches(0.45), Inches(1.38), Inches(6.0), Inches(2.6),
        "G3 CALIBRATION EVIDENCE (PER-FOLD PLATT)", WHITE, LINE, NAVY)
cal = [
    ("slope 0.86", "recalibration slope inside the [0.5, 2] ship-range "
     "defined in the V1.1 spec; V1 had slope 4.72 = rank-only"),
    ("Brier 0.124", "vs constant-predictor reference 0.242 - probabilities "
     "carry real information"),
    ("CITL -0.062", "calibration-in-the-large logit gap ~ 0 - no systematic "
     "over/under-confidence"),
]
ry = Inches(1.84)
for big, sub in cal:
    add_text(s, Inches(0.68), ry, Inches(1.75), Inches(0.35),
             [(big, 15, TEAL, True)])
    add_text(s, Inches(2.50), ry + Inches(0.02), Inches(3.8), Inches(0.62),
             [(sub, 8.5, INK, False)], space_after=0)
    ry += Inches(0.70)

chipbox(s, Inches(0.45), Inches(4.18), Inches(6.0), Inches(2.25),
        "WEEKLY OPERATING CADENCE", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(4.58), Inches(5.6), Inches(1.75), [
    [("Score every truck weekly from fresh 30-90 d feature windows. "
      "Monthly cadence gives only ~2 reads inside the 10-week validity "
      "horizon - too coarse to act on.", INK, False)],
    [("Tier is re-evaluated each week; alerts (slide 13) run alongside "
      "as independent channels.", INK, False)],
], size=9.5, space_after=5)

tiers = [
    ("GREEN", "P < 0.35", "20 trucks (4 F + 16 NF)", "routine operation; "
     "weekly re-score only", GREEN_D, GREEN_L),
    ("AMBER", "0.35 <= P < 0.55", "2 trucks (both NF)", "bundle inspection "
     "into next scheduled service", AMBER_D, AMBER_L),
    ("RED", "P >= 0.55", "12 trucks (10 F + 2 NF)", "inspect starter + "
     "battery circuit within 2-4 weeks", RED_D, RED_L),
]
ty = Inches(1.38)
for name, rule, count, action, tc, tf_ in tiers:
    add_rect(s, Inches(6.75), ty, Inches(6.13), Inches(1.52), tf_,
             line_color=tc, line_w=1.0)
    add_text(s, Inches(6.95), ty + Inches(0.10), Inches(2.4), Inches(0.35),
             [(name, 15, tc, True)])
    add_text(s, Inches(6.95), ty + Inches(0.52), Inches(2.6), Inches(0.3),
             [(rule, 10.5, tc, True)], font=MONO)
    add_text(s, Inches(6.95), ty + Inches(0.88), Inches(2.6), Inches(0.3),
             [(count, 9, SLATE, True)])
    add_text(s, Inches(9.65), ty + Inches(0.16), Inches(3.05), Inches(1.2),
             [(action, 10, INK, False)], space_after=0)
    ty += Inches(1.68)

footer(s, "Slope 0.86 means a shown 70% is a real ~70% — tiers are "
          "probability statements, not just ranks.",
       "Approve shipping recalibrated probabilities (not tier-only) to the "
       "pilot dashboards.")

# ===========================================================================
# SLIDE 13 - Alert layer specification
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Alert layer", "Three channels, three verdicts: pager, condition "
       "flag, corroborator", 13)

cards = [
    ("A2 · BATTERY-CASCADE TRIPLE", GREEN_D, GREEN_L, "SHIP AS PAGER",
     [("Rule (all three, causal at scoring week):", NAVY, True),
      ("rest-VSI step <= -0.5 V (SNR >= 2)\n+ drive-VSI step >= +0.3 V "
       "within +/-8 wk\n+ dip-depth widening > +1 V (>= 10 events)", INK,
       False),
      ("Validation: 4/5 battery-archetype failures caught (miss: VIN2_F, "
       "no qualifying drive-step); 0/20 NF false alarms; median lead 66.5 d "
       "(28-91).", INK, False),
      ("Battery-replacement separation: 5 NF trucks with rest-VSI steps UP "
       "do NOT fire.", GREEN_D, True)]),
    ("PERSISTENCE · VOLATILITY FLAG", PURP_D, PURP_L,
     "TERMINAL CONDITION FLAG ONLY",
     [("Rule:", NAVY, True),
      ("trailing-4-wk within-week VSI-std ratio above training-fold NF p90 "
       "envelope in >= 4 of last 12 wks", INK, False),
      ("Validation: fires on 13/14 failed (median terminal lead 168 d). "
       "BUT: all 20/20 NF enter fire-state at least once (avg 31.4% of "
       "weeks in-fire) - walking-alarm pathology.", INK, False),
      ("Never a standalone first-crossing pager; gate behind AMBER/RED "
       "tier.", PURP_D, True)]),
    ("A1 · CRANK-BURST", BLUE_D, BLUE_L, "AMBER/RED CORROBORATOR",
     [("Rule:", NAVY, True),
      ("daily failed-cranks + retries within 120 s; 7-day rolling sum S7 > "
       "own-first-half mean + 3 SD (floor S7 >= 3) for >= 2 consecutive "
       "days; SMA-dead excluded", INK, False),
      ("Validation: 4/12 applicable failed fire (leads 128-179 d vs last "
       "telemetry; VIN1_F 232 d vs recorded failure date); 8/15 applicable "
       "NF fire = 1.52 episodes/truck-year - too noisy alone.", INK, False),
      ("Value: recovered VIN1_F, a GREEN-tier miss - 232 d before recorded "
       "failure.", BLUE_D, True)]),
]
cx = Inches(0.45)
cw_ = Inches(4.02)
for title, edge, fill, verdict, paras in cards:
    add_rect(s, cx, Inches(1.38), cw_, Inches(4.55), WHITE, line_color=LINE,
             line_w=1.0)
    add_rect(s, cx, Inches(1.38), cw_, Inches(0.36), fill)
    add_text(s, cx + Inches(0.12), Inches(1.44), cw_ - Inches(0.24),
             Inches(0.28), [(title, 9.5, edge, True)])
    add_multirun(s, cx + Inches(0.14), Inches(1.88), cw_ - Inches(0.28),
                 Inches(3.3), [[p] for p in paras], size=8.5, space_after=5)
    add_rect(s, cx + Inches(0.14), Inches(5.42), cw_ - Inches(0.28),
             Inches(0.36), fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, cx + Inches(0.14), Inches(5.49), cw_ - Inches(0.28),
             Inches(0.26), [(verdict, 9, edge, True)], align=PP_ALIGN.CENTER)
    cx += cw_ + Inches(0.185)

add_multirun(s, Inches(0.45), Inches(6.02), Inches(12.4), Inches(0.5), [
    [("Combined policy: 13/14 failed trucks fire >= 1 channel (median "
      "first-fire lead 168 d - a distinct metric, equal to the persistence "
      "terminal median by coincidence); 10/20 NF completely clean; "
      "VIN16/19_NF fire all three - candidates for physical inspection.",
      NAVY, True)],
], size=10, space_after=0)

footer(s, "Only A2 is precise enough to page on its own; the other channels "
          "add recall and corroboration behind the tier gate.",
       "Approve the channel verdicts (pager / condition flag / corroborator) "
       "as the pilot alert policy.")

# ===========================================================================
# SLIDE 14 - Combined alert policy & runway
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Alert outcomes", "Validated per-VIN alert leads: the runway the "
       "policy actually delivers", 14)

add_pic(s, ASSETS / "BIZ_warning_runway.png", Inches(0.45), Inches(1.35),
        w=Inches(7.6), border=False)

chipbox(s, Inches(8.35), Inches(1.35), Inches(4.53), Inches(5.05),
        "POLICY OUTCOMES (VALIDATED)", WHITE, LINE, NAVY)
add_multirun(s, Inches(8.55), Inches(1.77), Inches(4.15), Inches(4.55), [
    [("Failed fleet:  ", NAVY, True),
     ("13/14 fire >= 1 channel; median first lead 168 d (range 77-424 vs "
      "recorded failure date); persistence first x10, A1 first x3.", INK,
      False)],
    [("A2 confirmations:  ", GREEN_D, True),
     ("VIN3 (91 d), VIN6 (70 d), VIN13 (63 d), VIN14 (28 d) - short-fuse, "
      "zero-false-alarm escalation on top of early channels.", INK, False)],
    [("Healthy fleet burden:  ", AMBER_D, True),
     ("10/20 zero channels; 6 show one; 2 show two; VIN16/19_NF show all "
      "three (treat as inspection candidates, possibly true degraders).",
      INK, False)],
    [("The gap case:  ", RED_D, True),
     ("VIN9_F fires nothing - silent 142 d before failure. Telemetry "
      "continuity itself must become a monitored health signal.", INK,
      False)],
], size=9, space_after=7)

footer(s, "The alert layer converts a ranking model into an operating "
          "procedure with measured leads and measured burden.",
       "Confirm inspection follow-up for VIN2_NF / VIN5_NF to close the "
       "loop on all-three-channel healthy trucks.")

# ===========================================================================
# SLIDE 15 - Horizon validation
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Horizon", "Prequential validation: the score is trustworthy for "
       "10 weeks — and we can prove why", 15)

add_pic(s, ASSETS / "TECH_horizon_curve.png", Inches(0.45), Inches(1.38),
        w=Inches(8.3), border=False)

chipbox(s, Inches(9.05), Inches(1.38), Inches(3.83), Inches(5.0),
        "METHOD & READING", WHITE, LINE, NAVY)
add_multirun(s, Inches(9.25), Inches(1.80), Inches(3.45), Inches(4.5), [
    [("Method:  ", NAVY, True),
     ("cut each truck's history k weeks before its end, re-anchor all "
      "feature windows at the cut, re-score with the frozen 4-feature "
      "LOVO models, recompute AUROC.", INK, False)],
    [("k* = 10:  ", GREEN_D, True),
     ("largest k with AUROC >= 0.75 sustained from k=0 (isolated k=15 blip "
      "not counted).", INK, False)],
    [("Recall at 18/20 specificity decays 0.71 (k<=3) -> 0.36 (k=10): "
      "flag late, act fast.", INK, False)],
    [("Anti-leak dividend:  ", NAVY, True),
     ("an epoch/length artifact would NOT decay with distance-to-failure; "
      "the collapse at k=11 is the strongest single evidence the score is "
      "failure-locked.", INK, False)],
], size=9, space_after=7)

footer(s, "The 10-week horizon is a measured validity window — the same "
          "curve that certifies it also certifies non-leakage.",
       "Approve '<=10 weeks' as the only forward-looking claim allowed in "
       "external material.")

# ===========================================================================
# SLIDE 16 - Case studies: VIN8_F redemption & VIN9_F structural miss
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Case anatomy", "VIN8_F: the engineered catch · VIN9_F: the "
       "honest miss", 16)

add_pic(s, GRAPHS / "V1_1_SM_VIN8_F_SM_dashboard.png", Inches(0.45),
        Inches(1.35), w=Inches(6.12))
add_pic(s, GRAPHS / "V1_1_SM_VIN9_F_SM_dashboard.png", Inches(6.76),
        Inches(1.35), w=Inches(6.12))

cap_y = Inches(5.46)
add_rect(s, Inches(0.45), cap_y, Inches(6.12), Inches(1.02), GREEN_L)
add_multirun(s, Inches(0.62), cap_y + Inches(0.09), Inches(5.85),
             Inches(0.88), [
    [("VIN8_F - V1 missed it (P 0.303); V1.1 catches it (P 0.716, RED).  ",
      GREEN_D, True),
     ("SMA-dead + 37 d silent gap. The battery-step-aware rest-VSI "
      "re-baseline recovers the pre-silence signal; persistence lead "
      "98-135 d.", INK, False)],
], size=9, space_after=0)
add_rect(s, Inches(6.76), cap_y, Inches(6.12), Inches(1.02), RED_L)
add_multirun(s, Inches(6.93), cap_y + Inches(0.09), Inches(5.85),
             Inches(0.88), [
    [("VIN9_F - the structural miss (raw score 0.401 vs threshold 0.406; "
      "recalibrated P 0.22 shown on the dashboard).  ", RED_D, True),
     ("SMA-dead (0.4%), longest silent gap (142 d), mid-fleet on all live "
      "features, zero alerts. Fix is instrumentation + telemetry-continuity "
      "monitoring, not a better classifier.", INK, False)],
], size=9, space_after=0)

footer(s, "The same physics that explains the catch explains the miss — "
          "observability, not model capacity, is the boundary.",
       "Approve 'telemetry-continuity monitoring' as a pilot workstream "
       "(a truck going quiet is itself a trigger).")

# ===========================================================================
# SLIDE 17 - Failure physics, archetypes & limitations
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Limitations", "Failure archetypes, physics ceilings, and what "
       "V1.1 deliberately does not claim", 17)

arch = [
    ("A1 · solenoid intermittency", "x3", "failed-crank rate up to 4.3x, "
     "retries 7.5x; A1 channel territory", BLUE_D, BLUE_L),
    ("A2 · battery cascade", "x4-5", "rest-floor sag + drive rise + dip "
     "widening; A2 detector territory", GREEN_D, GREEN_L),
    ("A3 · VSI volatility", "x3", "sustained electrical noise without "
     "battery/crank signature", PURP_D, PURP_L),
    ("A4 · silent / abrupt", "x4", "no admissible precursor in 5 s / 0.2 V "
     "telemetry - the recall ceiling", RED_D, RED_L),
]
ax_ = Inches(0.45)
aw = Inches(3.02)
for name, n, desc, ec, fc in arch:
    add_rect(s, ax_, Inches(1.38), aw, Inches(1.45), fc, line_color=ec,
             line_w=1.0)
    add_text(s, ax_ + Inches(0.14), Inches(1.48), aw - Inches(0.28),
             Inches(0.3), [(f"{name}  {n}", 9.5, ec, True)])
    add_text(s, ax_ + Inches(0.14), Inches(1.82), aw - Inches(0.28),
             Inches(0.92), [(desc, 8.5, INK, False)], space_after=0)
    ax_ += aw + Inches(0.115)
add_text(s, Inches(0.45), Inches(2.92), Inches(12.4), Inches(0.3),
         [("Archetypes are in-sample descriptors (not out-of-fold "
           "validated); counts sum past 14 because VIN14_F carries a dual "
           "A1+A2 signature. Honest tier-RED recall ceiling ~= 10-11/14, "
           "set by the A1+A2+A3 count.", 9.5, SLATE, False, True)])

chipbox(s, Inches(0.45), Inches(3.42), Inches(6.0), Inches(2.98),
        "PHYSICS CEILINGS (WHY MORE MODELING WON'T HELP)", WHITE, LINE, RED_D)
add_multirun(s, Inches(0.65), Inches(3.82), Inches(5.6), Inches(2.5), [
    [("Brush wear:  ", RED_D, True),
     ("the real 60-120 d declining-crank-voltage precursor lives in "
      "sub-sample dip shape and sub-second duration growth - destroyed by "
      "5 s / 0.2 V sampling.", INK, False)],
    [("Bearing seizure:  ", RED_D, True),
     ("zero electrical precursor at any sampling rate.", INK, False)],
    [("No RUL:  ", RED_D, True),
     ("discrete-time hazard, Cox, Weibull and deep-survival variants all "
      "lose to a constant predictor (MAE 576 d vs 44 d). V1.1 ships a "
      "validity window, never a date.", INK, False)],
], size=9.5, space_after=7)

chipbox(s, Inches(6.75), Inches(3.42), Inches(6.13), Inches(2.98),
        "STATISTICAL HONESTY LEDGER", WHITE, LINE, NAVY)
add_multirun(s, Inches(6.95), Inches(3.82), Inches(5.75), Inches(2.5), [
    [("n = 34, 14 failures: CI floor 0.811; a handful of trucks decide "
      "every threshold - treat point estimates with CI-width humility.",
      INK, False)],
    [("OOF scores correlate with observation-length proxies (up to "
      "|0.653|) - defended as label-mediated via G1 + horizon collapse, "
      "and disclosed rather than hidden.", INK, False)],
    [("Suppressor coefficient (vsi_range_trend) documented; per-VIN "
      "explanations separate physics from model usage.", INK, False)],
    [("Permutation N=200 (runtime-bounded; target 1000) - p-value floor "
      "0.005 stated as such.", INK, False)],
], size=9.5, space_after=7)

footer(s, "V1.1's boundaries are measured, physical, and documented — the "
          "ceiling is data, not method.",
       "Accept the limitation ledger as the agreed basis for external "
       "claims and future-iteration bars.")

# ===========================================================================
# SLIDE 18 - Crank-event extraction (mermaid D1)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Event extraction", "Crank-event state machine: 20,471 events "
       "distilled from 106M raw rows", 18)

add_pic(s, ASSETS / "MMD_D1_crank_events.png", Inches(0.97), Inches(1.30),
        w=Inches(11.4), border=False)

consts = [
    ("<= 10 s", "intra-event gap: consecutive SMA==1 rows merge"),
    ("> 60 s", "max plausible duration -> artifact, excluded"),
    (">= 550 RPM", "success criterion on rpm_max_15s (+15 s window)"),
    ("[-90, -10] s", "dip baseline window (>= 3 valid VSI readings)"),
    ("120 s", "retry window: failed crank -> retry cluster / burst"),
    ("45 s", "recovery window for post-crank stabilisation"),
    ("+ 5 s", "sample width added to event duration"),
    ("1.2%", "reconciliation delta vs gap-naive reference (20,729)"),
]
cx0, cy0 = Inches(0.45), Inches(3.94)
cw_, chh = Inches(3.02), Inches(0.98)
for i, (big, sub) in enumerate(consts):
    x = cx0 + (i % 4) * (cw_ + Inches(0.115))
    y = cy0 + (i // 4) * (chh + Inches(0.14))
    add_rect(s, x, y, cw_, chh, LIGHT)
    add_text(s, x + Inches(0.14), y + Inches(0.09), cw_ - Inches(0.28),
             Inches(0.32), [(big, 14, NAVY, True)], font=MONO)
    add_text(s, x + Inches(0.14), y + Inches(0.44), cw_ - Inches(0.28),
             Inches(0.54), [(sub, 8.5, SLATE, False)], space_after=0)

add_multirun(s, Inches(0.45), Inches(6.14), Inches(12.4), Inches(0.4), [
    [("Cohort note: ", NAVY, True),
     ("the 7 SMA-dead trucks produce no usable events - their event-rate "
      "artifacts (up to 8.4/active-day) are excluded by cohort masking. The "
      "trucks themselves are still scored weekly via the voltage features "
      "(crank features imputed - see slide 19).", INK, False)],
], size=9.5, space_after=0)

footer(s, "Every crank feature rests on a deterministic, artifact-filtered "
          "event definition - reproducible from raw rows.",
       "Confirm the 10 s / 60 s / 550 RPM constants against powertrain "
       "engineering expectations.")

# ===========================================================================
# SLIDE 19 - Weekly scoring & alert decision logic (mermaid D2)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Runtime logic", "One week in production: scoring, tiering and "
       "alert gating as a decision flow", 19)

add_pic(s, ASSETS / "MMD_D2_weekly_scoring.png", Inches(0.45), Inches(1.30),
        w=Inches(12.43), border=False)

chipbox(s, Inches(0.45), Inches(4.80), Inches(12.43), Inches(1.66),
        "RUNTIME NOTES", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(5.18), Inches(12.0), Inches(1.25), [
    [("Cadence:  ", NAVY, True),
     ("weekly, on masked-week close. Feature windows re-anchor at the "
      "scoring cut - the exact machinery validated by the prequential "
      "horizon test (slide 15).", INK, False)],
    [("Honest scoring:  ", NAVY, True),
     ("each truck is scored by a model trained on the other 33 (LOVO "
      "parity with validation); recalibration is per-fold Platt.", INK,
      False)],
    [("Gating:  ", NAVY, True),
     ("A2 battery-cascade pages independently of tier (0/20 NF false "
      "alarms earns that right); persistence and A1 annotate only at "
      "AMBER/RED - never page on their own.", INK, False)],
], size=9.5, space_after=5)

footer(s, "The deployment flow is the validation flow - no new logic "
          "appears between the audit and the pilot.",
       "Sign off the decision flow as the pilot SOP skeleton (workshop "
       "order = tier + alert annotations).")

# ===========================================================================
# SLIDE 20 - Battery-step re-baselining (mermaid D3) + step ledger
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Maintenance separation", "Battery-step detection: replacements "
       "re-baseline features, cascades feed the A2 pager", 20)

add_pic(s, ASSETS / "MMD_D3_battery_step.png", Inches(2.97), Inches(1.30),
        w=Inches(7.4), border=False)

chipbox(s, Inches(0.45), Inches(4.42), Inches(6.0), Inches(2.0),
        "DETECTED UP-STEPS (REPLACEMENT CANDIDATES, E5)", WHITE, LINE,
        GREEN_D)
add_multirun(s, Inches(0.65), Inches(4.82), Inches(5.6), Inches(1.5), [
    [("VIN32_NF +1.40 V (2024-05) · VIN26_NF +0.70 V · VIN17_NF +0.61 V · "
      "VIN19_NF +0.61 V · VIN31_NF +0.59 V · VIN8_F +0.60 V (2024-06).",
      INK, False)],
    [("None of the 5 NF up-step trucks fires A2 - the detector separates "
      "maintenance from degradation by construction.", GREEN_D, True)],
], size=9.5, space_after=5)

chipbox(s, Inches(6.75), Inches(4.42), Inches(6.13), Inches(2.0),
        "DETECTED DOWN-STEPS (CASCADE EVIDENCE)", WHITE, LINE, RED_D)
add_multirun(s, Inches(6.95), Inches(4.82), Inches(5.75), Inches(1.5), [
    [("VIN6_F -2.71 V (2025-08) · VIN14_F -2.31 V (2025-09) · VIN3_F "
      "-1.70 V (2025-06) · VIN2_F -1.59 V (2025-09). All pass the SNR >= 2 "
      "detector gate (strongest 5.8).", INK, False)],
    [("The four largest rest-VSI collapses are exactly the A2 battery-"
      "cascade archetype trucks - independent corroboration of the "
      "archetype assignment.", RED_D, True)],
], size=9.5, space_after=5)

footer(s, "One detector, two uses: up-steps clean the features (recovered "
          "VIN8_F), down-steps power the zero-false-alarm pager.",
       "Validate the 6 replacement candidates against service records when "
       "DICV maintenance data becomes available.")

# ===========================================================================
# SLIDE 21 - Survival / RUL rejection (F-track)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Rejected alternative: RUL", "Survival analysis was built, "
       "measured, and rejected - a constant beats every RUL model", 21)

chipbox(s, Inches(0.45), Inches(1.38), Inches(4.02), Inches(3.3),
        "RUL MAE - FAILED VINS, LAST 26 WEEKS", WHITE, LINE, NAVY)
mae_rows = [
    ("44.4 d", "constant '91 days' predictor (structural floor; actuals "
     "~uniform 7-182 d)", GREEN_D),
    ("461.9 d", "Weibull fleet-clock conditional median (LOVO)", RED_D),
    ("576.1 d", "discrete-time hazard model (logistic, LOVO)", RED_D),
]
ry = Inches(1.86)
for big, sub, col in mae_rows:
    add_text(s, Inches(0.68), ry, Inches(1.6), Inches(0.4),
             [(big, 17, col, True)])
    add_text(s, Inches(2.30), ry + Inches(0.02), Inches(2.0), Inches(0.75),
             [(sub, 8.5, INK, False)], space_after=0)
    ry += Inches(0.85)
add_text(s, Inches(0.68), ry + Inches(0.02), Inches(3.6), Inches(0.5),
         [("The constant wins by 10.4x and 13.0x.", 10, NAVY, True)])

chipbox(s, Inches(4.66), Inches(1.38), Inches(4.02), Inches(3.3),
        "HAZARD-MODEL DIAGNOSTICS", WHITE, LINE, NAVY)
add_multirun(s, Inches(4.86), Inches(1.80), Inches(3.65), Inches(2.8), [
    [("Pooled weekly-hazard AUROC 0.747 (14 event-weeks vs 2,622 at-risk "
      "weeks).", INK, False)],
    [("Age-matched concordance 0.654 (332 pairs).", INK, False)],
    [("P(fail <= H) classification: 0.744 @30 d · 0.708 @60 d · 0.688 "
      "@90 d.", INK, False)],
    [("Truck-level ranking from held-out hazards: 0.586 - collapses vs "
      "the static model's 0.893.", RED_D, True)],
], size=9.5, space_after=7)

chipbox(s, Inches(8.87), Inches(1.38), Inches(4.02), Inches(3.3),
        "WHAT SURVIVES FROM THE F-TRACK", WHITE, LINE, TEAL)
add_multirun(s, Inches(9.07), Inches(1.80), Inches(3.65), Inches(2.8), [
    [("Fleet Weibull clock: ", NAVY, True),
     ("lambda 133.3 wk, rho 2.03 (wear-out), median 111.3 wk (779 d) - "
      "used ONLY as the dashboard RUL-curve illustration anchor.", INK,
      False)],
    [("Cox time-varying: ", NAVY, True),
     ("vsi_std_ratio HR 1.74/unit (p=0.002), coefficient stable +0.69 "
      "+/-0.04 across folds - independent confirmation of the volatility "
      "workhorse.", INK, False)],
    [("Everything else is documented and retired.", SLATE, False)],
], size=9.5, space_after=7)

add_rect(s, Inches(0.45), Inches(4.88), Inches(12.43), Inches(0.62), LIGHT)
add_text(s, Inches(0.68), Inches(5.02), Inches(12.0), Inches(0.35),
         [("RUL claims are withheld because we measured them failing - not "
           "because we did not try. Three model families, one conclusion.",
           11, NAVY, True)])

footer(s, "14 events cannot parameterize a per-truck failure clock; the "
          "validated deliverable is risk + a 10-week window.",
       "Ratify the 'risk + window, no dates' deliverable contract for the "
       "pilot and all external material.")

# ===========================================================================
# SLIDE 22 - Capacity budget & variant sweep (G1 + B4)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Rejected alternative: deep learning", "Parameter-budget math: "
       "why 14 failures cannot feed a sequence model", 22)

chipbox(s, Inches(0.45), Inches(1.38), Inches(6.35), Inches(3.62),
        "MINIMAL-CONFIG PARAMETER BUDGET (EPV >= 10 -> ~1.4 PARAMS)",
        WHITE, LINE, NAVY)
budget = [
    ("V1.1 Ridge (4 features)", "5", "4x", GREEN_D),
    ("LSTM, 1 layer, h=8", "329", "235x", RED_D),
    ("TCN, 2 blocks, 8 ch", "489", "349x", RED_D),
    ("Transformer enc, d=16", "2,273", "1,624x", RED_D),
    ("PatchTST, d=16", "2,385", "1,704x", RED_D),
    ("Informer, d=16", "3,057", "2,184x", RED_D),
    ("TFT, d=16", "8,785", "6,275x", RED_D),
]
ry = Inches(1.84)
add_text(s, Inches(0.68), ry, Inches(3.4), Inches(0.24),
         [("architecture (minimal config)", 8.5, SLATE, True)])
add_text(s, Inches(4.18), ry, Inches(1.1), Inches(0.24),
         [("params", 8.5, SLATE, True)], align=PP_ALIGN.RIGHT)
add_text(s, Inches(5.38), ry, Inches(1.2), Inches(0.24),
         [("over budget", 8.5, SLATE, True)], align=PP_ALIGN.RIGHT)
ry += Inches(0.3)
for name, params, over, col in budget:
    add_text(s, Inches(0.68), ry, Inches(3.5), Inches(0.26),
             [(name, 9.5, INK, False)])
    add_text(s, Inches(4.18), ry, Inches(1.1), Inches(0.26),
             [(params, 9.5, NAVY, True)], align=PP_ALIGN.RIGHT, font=MONO)
    add_text(s, Inches(5.38), ry, Inches(1.2), Inches(0.26),
             [(over, 9.5, col, True)], align=PP_ALIGN.RIGHT, font=MONO)
    ry += Inches(0.315)
add_text(s, Inches(0.68), ry + Inches(0.04), Inches(5.9), Inches(0.55),
         [("Even LSTM h=2 (43 params, 31x over) is seed-unstable: AUROC "
           "0.854 / 0.882 / 0.918 across three seeds.", 9, RED_D, True)],
         space_after=0)

chipbox(s, Inches(7.00), Inches(1.38), Inches(5.88), Inches(3.62),
        "WHAT WAS STILL TRIED (HONEST PROBES, LOVO)", WHITE, LINE, NAVY)
add_multirun(s, Inches(7.20), Inches(1.80), Inches(5.5), Inches(3.1), [
    [("Sequence probes (G2, screening):  ", NAVY, True),
     ("PCA3 0.868-0.900 (PC1 leak-flagged: r +0.50 vs n_weeks); trend-"
      "coefficient Ridge 0.918-0.925; 1-NN margin 0.914. None survives the "
      "leak audit with a genuine edge.", INK, False)],
    [("Variant sweep (B4, non-nested V1-era screening):  ", NAVY, True),
     ("baseline 0.9214; drop-artifact 0.8643; +within-week noise feature "
      "0.9679 - the screening high-water mark that motivated V1.1's "
      "winning feature.", INK, False)],
    [("Screening numbers are NOT comparable to the nested 0.9321 - they "
      "exist to pick candidates, and are labeled as such wherever quoted.",
      AMBER_D, True)],
], size=9.5, space_after=8)

footer(s, "At 14 failures, capacity is the enemy: anything bigger than ~5 "
          "parameters memorizes epoch, not physics.",
       "Adopt the EPV parameter-budget rule for all proposed model work "
       "until n_failed >= 30-50.")

# ===========================================================================
# SLIDE 23 - Data-quality probe atlas (A-track)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Data-quality audit", "Probe atlas: what the A-track found before "
       "a single feature was built", 23)

probes = [
    ("P1 · LABELS & LIFETIME", "Failure_type constant across all failed "
     "rows; 0 rows after JCOPENDATE; 9/14 failed VINs transmit up to the "
     "JCO day; sale-to-failure life 171-671 d.", BLUE_D),
    ("P2 · MISSINGNESS = 3 LAYERS", "config (SMA-dead: 99.7-99.9% null), "
     "regime (engine-off nulls 13-86% = duty cycle), cohort (VSI null "
     "13.3% F vs 18.6% NF). Features must be regime-aware.", BLUE_D),
    ("P3 · SAMPLING FAMILIES", "median dt 5 s; continuous (p99 6 s) vs "
     "rest-heartbeat (p99 ~900 s). Gap counts classify at 0.868-0.875 = "
     "volume proxies -> banned as features.", AMBER_D),
    ("P4 · VSI SENSOR REALITY", "0.2 V quantization (168 distinct values); "
     "per-truck setpoints 27.6-28.2 V (fleet std ~1 LSB) -> per-VIN "
     "baselining mandatory; stuck runs (<=58 h) not failure-correlated.",
     AMBER_D),
    ("P5 · SINGLE-SIGNAL SCREENS", "vsi_drive_std 0.732 · settled p05 "
     "0.734 · rest p05 0.729. Dead ends: RPM profile 0.500, crank duration "
     "0.507, dip-depth LEVEL 0.575, success rate 0.543.", GREEN_D),
    ("P6 · LATENT DEGRADATION", "within-VIN delta(vsi_drive_std), last 8 "
     "wks vs own baseline: AUROC 0.893 (calendar-matched control 0.889); "
     "9/14 failed respond; onset up to 294 d (VIN13_F) / 280 d (VIN6_F).",
     GREEN_D),
]
px0, py0 = Inches(0.45), Inches(1.38)
pw, ph = Inches(4.02), Inches(2.42)
for i, (t, d, col) in enumerate(probes):
    x = px0 + (i % 3) * (pw + Inches(0.185))
    y = py0 + (i // 3) * (ph + Inches(0.16))
    add_rect(s, x, y, pw, ph, WHITE, line_color=LINE, line_w=1.0)
    add_rect(s, x, y, Inches(0.07), ph, col)
    add_text(s, x + Inches(0.18), y + Inches(0.10), pw - Inches(0.32),
             Inches(0.28), [(t, 9.5, col, True)])
    add_text(s, x + Inches(0.18), y + Inches(0.44), pw - Inches(0.32),
             ph - Inches(0.55), [(d, 8.8, INK, False)], space_after=0)

footer(s, "The audit found the signal (within-VIN volatility deltas, P6) "
          "and the traps (volume, epoch, config) before modeling began.",
       "Carry P2-P4 forward as mandatory preprocessing requirements for "
       "the 500-truck scale-up.")

# ===========================================================================
# SLIDE 24 - Reproducibility & artifact map (mermaid D4)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Reproducibility", "Artifact dependency map: review this deck by "
       "re-running it", 24)

add_pic(s, ASSETS / "MMD_D4_artifact_map.png", Inches(0.45), Inches(1.40),
        w=Inches(9.0), border=False)

chipbox(s, Inches(9.70), Inches(1.40), Inches(3.18), Inches(5.0),
        "DETERMINISM LEDGER", WHITE, LINE, NAVY)
add_multirun(s, Inches(9.88), Inches(1.82), Inches(2.85), Inches(4.5), [
    [("Seeds: ", NAVY, True),
     ("bootstrap 42, permutation 43, global numpy 42.", INK, False)],
    [("Runtimes: ", NAVY, True),
     ("nested main 132.7 s; permutation 828.5 s; dashboards ~1 min/VIN.",
      INK, False)],
    [("Parity: ", NAVY, True),
     ("closed-form Ridge replica verified < 1e-9 vs sklearn.", INK, False)],
    [("Interpreter: ", NAVY, True),
     ("py -3 (repo convention); no GPU, no cloud dependency.", INK, False)],
    [("Traceability: ", NAVY, True),
     ("display renumbering in V1_1_SM_vin_naming_map.csv; gates.json is "
      "the machine-readable audit record.", INK, False)],
], size=9, space_after=7)

add_multirun(s, Inches(0.45), Inches(5.75), Inches(9.0), Inches(0.7), [
    [("Single-command re-runs:  ", NAVY, True),
     ("V1_1_SM_features.py -> V1_1_SM_nested_ridge.py -> V1_1_SM_alerts.py "
      "-> V1_1_SM_horizon.py -> V1_1_SM_production_graphs.py (each "
      "regenerates its artifacts in place).", INK, False)],
], size=9.5, space_after=0)

footer(s, "Fixed seeds, committed artifacts, named scripts - every number "
          "in this deck regenerates on demand.",
       "Archive the V1.1 artifact set as the frozen reference package for "
       "all future iteration comparisons.")

# ===========================================================================
# SLIDE 25 - Recommendations & sign-off
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Sign-off", "Recommendations: ship the validated core, "
       "instrument for the next ceiling", 25)

chipbox(s, Inches(0.45), Inches(1.38), Inches(6.0), Inches(2.85),
        "SHIP NOW (VALIDATED IN V1.1)", WHITE, LINE, GREEN_D)
ship = [
    "Frozen modal-4 nested RidgeClassifier scorer + per-fold Platt "
    "probabilities (spec, seeds, artifacts committed)",
    "GREEN/AMBER/RED tier rule on recalibrated P (0.35 / 0.55)",
    "3-channel alert policy with the verdicts of slide 13",
    "Weekly scoring cadence; <=10-week horizon as the only forward claim",
    "34 per-VIN dashboards + explanation cards for workshop use",
]
ry = Inches(1.82)
for t in ship:
    add_rect(s, Inches(0.68), ry + Inches(0.05), Inches(0.11), Inches(0.11),
             GREEN_D, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(0.92), ry - Inches(0.02), Inches(5.4), Inches(0.5),
             [(t, 9, INK, False)], space_after=0)
    ry += Inches(0.47)

chipbox(s, Inches(0.45), Inches(4.42), Inches(6.0), Inches(1.98),
        "PILOT WORKSTREAMS", WHITE, LINE, NAVY)
add_multirun(s, Inches(0.65), Inches(4.82), Inches(5.6), Inches(1.5), [
    [("1. Alert -> workshop-finding feedback loop (labels for the next "
      "iteration).", INK, False)],
    [("2. Telemetry-continuity monitoring (the VIN9_F fix).", INK, False)],
    [("3. Physical inspection of VIN19_NF / VIN16_NF (all-channel healthy "
      "trucks).", INK, False)],
], size=9.5, space_after=5)

chipbox(s, Inches(6.75), Inches(1.38), Inches(6.13), Inches(2.85),
        "INSTRUMENTATION ASKS (WHAT BREAKS THE CEILING)", WHITE, LINE, TEAL)
inst = [
    (">= 1 Hz crank-window logging (during SMA=1)", "revives the physically "
     "real brush-wear prognosis channel - biggest unlock"),
    ("Cranking current or battery SoC/SoH", "ends the battery-vs-starter "
     "ambiguity that caps alert precision"),
    ("Maintenance / parts-replacement records", "turns data-derived "
     "archetypes into supervised labels"),
    ("More failure cases (n_failed >= 30-50)", "unlocks survival modeling "
     "and representation learning credibly"),
]
ry = Inches(1.84)
for t, d in inst:
    add_text(s, Inches(6.95), ry, Inches(5.75), Inches(0.24),
             [(t, 9.5, TEAL, True)], space_after=0)
    add_text(s, Inches(6.95), ry + Inches(0.22), Inches(5.75), Inches(0.36),
             [(d, 8.5, SLATE, False)], space_after=0)
    ry += Inches(0.60)

add_rect(s, Inches(6.75), Inches(4.42), Inches(6.13), Inches(1.98), NAVY)
add_text(s, Inches(6.98), Inches(4.58), Inches(5.7), Inches(0.3),
         [("BOTTOM LINE FOR APPROVERS", 9.5,
           RGBColor(0x9F, 0xB3, 0xC8), True)])
add_text(s, Inches(6.98), Inches(4.90), Inches(5.7), Inches(1.4),
         [("V1.1 is the audited ceiling of the current telemetry: 0.9321 "
           "nested AUROC, calibrated probabilities, a zero-false-alarm "
           "pager and a proven 10-week window. Approving the five asks of "
           "slide 2 - model freeze, calibrated tiers, alert policy, weekly "
           "cadence, instrumentation - ships everything the data supports "
           "and funds exactly the data that lifts the ceiling.", 11, WHITE,
           True)], space_after=0)

footer(s, "Approve asks 1-5 (slide 2) to move V1.1 from frozen iteration to "
          "operating pilot.",
       "Sign-off: model freeze · tier+alert policy · weekly cadence · "
       "instrumentation requests.")

prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides._sldIdLst)} slides)")
