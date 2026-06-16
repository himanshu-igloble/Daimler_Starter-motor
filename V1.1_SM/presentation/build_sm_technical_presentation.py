#!/usr/bin/env python3
"""
Build Daimler Starter Motor Predictive Maintenance Technical Review presentation.
V1.1_SM -- validated metrics only. Same design language as the Alternator V10.6.2 deck.

All numbers sourced from:
  STARTER MOTOR/V1.1/reports/V1_1_SM_comparison_report.md
  STARTER MOTOR/V1.1/reports/V1_1_SM_model_card.md
  STARTER MOTOR/V1.1/reports/V1_1_SM_executive_recommendation.md
  STARTER MOTOR/V1.1/reports/V1_1_SM_alerts_horizon.md
  STARTER MOTOR/V1.1/reports/V1_1_SM_experiment_results.md
  STARTER MOTOR/V1.1/Plan/V1_1_SM_spec.md
  STARTER MOTOR/reports/V1_SM_final_report.md
  STARTER MOTOR/V1.1/results/V1_1_SM_horizon_curve.csv
  STARTER MOTOR/V1.1/results/V1_1_SM_nested_lovo_predictions.csv
"""

import os
import csv
import sys
import tempfile
import shutil
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── PATHS ──────────────────────────────────────────────────────
ROOT   = Path(__file__).resolve().parent.parent.parent.parent
V11    = ROOT / "STARTER MOTOR" / "V1.1"
GRAPHS = V11 / "graphs"
RESULTS = V11 / "results"

OUT_DIR = V11 / "presentation"
OUT_DIR.mkdir(exist_ok=True)

# ── VIN DISPLAY RENUMBERING (2026-06-11) ───────────────────────
# Failed VIN1-14 unchanged; NF shifted +14 (old VIN1_NF -> VIN15_NF). All
# slide text/tables pass through map_nf_text at render time; the underlying
# results artifacts retain the ORIGINAL labels (results/V1_1_SM_vin_naming_map.csv).
sys.path.insert(0, str(V11 / "src"))
from V1_1_SM_vin_display_map import map_nf_text, display_label

VIN_FOOTNOTE = ("Sequential fleet numbering: VIN1-14 failed, VIN15-34 in-service "
                "(raw-file mapping: results/V1_1_SM_vin_naming_map.csv).")

# ── DIMENSIONS ─────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)

# ── COLOURS (identical to ALT V10.6.2 deck) ────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
GOLD       = RGBColor(0xC5, 0x8B, 0x1F)
GREY_MED   = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GREY  = RGBColor(0xB0, 0xB8, 0xC4)
KT_HEADER  = RGBColor(0x2E, 0x50, 0x90)
KT_BODY    = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_PASS = RGBColor(0x1E, 0x8C, 0x45)
RED_FAIL   = RGBColor(0xC0, 0x39, 0x2B)
ORANGE_W   = RGBColor(0xE3, 0x6C, 0x09)
BG_LIGHT   = RGBColor(0xF5, 0xF6, 0xFA)
BG_TABLE_H = RGBColor(0x0D, 0x1B, 0x2A)
BG_TABLE_A = RGBColor(0xF0, 0xF2, 0xF5)
BG_TABLE_B = RGBColor(0xFF, 0xFF, 0xFF)
KPI_BLUE   = RGBColor(0x1A, 0x5C, 0xB0)
KPI_GREEN  = RGBColor(0x15, 0x7F, 0x3D)
KPI_AMBER  = RGBColor(0xB8, 0x7A, 0x0F)
KPI_RED    = RGBColor(0xB0, 0x2A, 0x2A)

FONT = 'Calibri'
TMPDIR = tempfile.mkdtemp(prefix='sm_charts_')

# ── VALIDATED METRICS ──────────────────────────────────────────
# Nested-LOVO Ridge (X2 -- V1_1_SM_experiment_results.md / model card)
AUROC_V11      = 0.9321
AUROC_V1_REP   = 0.9214   # V1 as originally reported (non-nested)
AUROC_V1_REST  = 0.8929   # V1 restated under nested protocol
AUROC_ABLATION = 0.8429   # nested protocol on V1-era de-artifacted features
BOOTSTRAP_CI   = (0.811, 0.986)
PERMUTATION_P  = 0.005
OPTIMISM_V11   = 0.0036
OPTIMISM_V1    = 0.0285
JACKKNIFE      = (0.927, 0.951)
TP, FP, FN, TN = 13, 5, 1, 15          # per-fold Youden operating point
RECALL_YOUDEN  = "13/14"
SPEC_YOUDEN    = "15/20"
RECALL_RED     = "10/14"
SPEC_RED       = "18/20"
F1_SCORE, MCC  = 0.812, 0.669
CAL_SLOPE, BRIER, CITL = 0.86, 0.124, -0.062
N_FEATURES     = 4

FEATURES = [
    ("vsi_withinwk_std_ratio_30d_w", +0.8862, "Within-week VSI noise, last 4 wk vs own 40-wk baseline"),
    ("vsi_range_trend",              -0.4139, "Weekly drive-VSI range Theil-Sen slope, last 12 wk (suppressor)"),
    ("rest_vsi_p05_delta90",         -0.2704, "Engine-off rest-voltage floor delta vs own baseline (battery-step aware)"),
    ("dip_depth_last90_delta",       +0.1409, "Crank dip-depth delta, last 90 d vs own baseline"),
]

# Fleet (V1_SM_final_report.md)
N_VINS, N_FAILED, N_NF = 34, 14, 20
N_ROWS = 106_445_161
N_TRUCK_WEEKS = 2_636
N_CRANK_EVENTS = 20_471

# Alerts (V1_1_SM_alerts_horizon.md)
A2_RECALL, A2_FP, A2_MEDIAN_LEAD = "4/5", "0/20", 66.5
PERS_RECALL, PERS_FP, PERS_MEDIAN_LEAD = "13/14", "4/20", 168
A1_FP_RATE = 1.52
COMBINED_RECALL = "13/14"
NF_CLEAN = "10/20"

# Horizon (X4)
K_STAR = 10
HORIZON_K0, HORIZON_K10, HORIZON_K11 = 0.9357, 0.768, 0.704
HORIZON_TAIL = 0.592

# Closed negatives (comparison report / exec recommendation)
SURV_RUL_MAE, CONST_RUL_MAE = 576, 44
SURV_RANKING = 0.586

# Fleet tiers (recalibrated, V1_1_SM_nested_lovo_predictions.csv)
NF_GREEN, NF_AMBER, NF_RED = 16, 2, 2
F_RED, F_GREEN = 10, 4


# ── HELPER FUNCTIONS ───────────────────────────────────────────
# Identical to reference Alternator presentation for visual consistency.

def make_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    return shape

def add_rounded_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1.5)
    else:
        ln.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, font_size=12, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = map_nf_text(str(text))      # render-time VIN display mapping
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(2)
    return txBox

def add_multiline(slide, left, top, w, h, lines, font_size=11, color=DARK_TEXT,
                  bold=False, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, bold, color
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if bullet else "") + map_nf_text(str(txt))
        p.font.size = Pt(font_size)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = FONT
        p.space_after = Pt(3)
    return txBox

def add_header_bar(slide, title_text, subtitle_text=""):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             title_text, font_size=22, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.35),
                 subtitle_text, font_size=13, color=LIGHT_GREY)

def add_section_slide(slide, number, title_text, subtitle_text=""):
    add_rect(slide, Inches(0), Inches(0), SW, SH, NAVY)
    add_text(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.8),
             f"{number:02d}  {title_text}", font_size=34, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, Inches(0.8), Inches(3.4), Inches(11), Inches(0.5),
                 subtitle_text, font_size=16, color=LIGHT_GREY)
    add_footer(slide, dark_bg=True)

def add_footer(slide, dark_bg=False):
    y = Inches(7.05)
    if not dark_bg:
        add_rect(slide, Inches(0), y, SW, Inches(0.45), BG_LIGHT)
    clr = LIGHT_GREY if dark_bg else GREY_MED
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "V1.1_SM  |  2026-06-10  |  BharatBenz 5528T Starter Motor Predictive Maintenance  |  BytEdge CONFIDENTIAL",
             font_size=8, color=clr, bold=False)

def add_key_takeaways(slide, bullets, left=Inches(0.4), top=Inches(5.45),
                      width=Inches(12.5), height=None):
    if height is None:
        height = Inches(0.25 + 0.22 * len(bullets))
    add_rounded_rect(slide, left, top, width, height,
                     RGBColor(0xE8, 0xEE, 0xF7), border=KT_HEADER)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), Inches(3), Inches(0.25),
             "KEY TAKEAWAYS", font_size=10, bold=True, color=KT_HEADER)
    y = top + Inches(0.28)
    for b in bullets:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.2),
                 "•  " + b, font_size=9, color=KT_BODY)
        y += Inches(0.19)

def add_kpi_tile(slide, left, top, w, h, label, value, color=KPI_BLUE, status=""):
    add_rounded_rect(slide, left, top, w, h, WHITE, border=color)
    add_text(slide, left + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.22),
             label, font_size=9, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.3), w - Inches(0.2), Inches(0.4),
             str(value), font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    if status:
        sc = GREEN_PASS if "PASS" in status else (RED_FAIL if "FAIL" in status else KPI_AMBER)
        add_text(slide, left + Inches(0.1), top + Inches(0.68), w - Inches(0.2), Inches(0.2),
                 status, font_size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)

def add_table_shape(slide, left, top, w, rows, cols, data, col_widths=None):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ri, row_data in enumerate(data):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = map_nf_text(str(val))   # render-time VIN display mapping
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            if ri == 0:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_H
            else:
                p.font.size = Pt(8.5)
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_A if ri % 2 == 0 else BG_TABLE_B
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return tbl_shape


# ── CHART GENERATION (matplotlib, ALT styling) ─────────────────

def chart_auroc_progression():
    """The honest-validation differentiator: V1 reported -> restated -> V1.1."""
    fig, ax = plt.subplots(figsize=(10, 3.2))
    cats = ['V1 as reported\n(non-nested)', 'V1 restated\n(nested, honest)',
            'V1-era features\nunder nested protocol', 'V1.1 nested\n(headline)']
    vals = [AUROC_V1_REP, AUROC_V1_REST, AUROC_ABLATION, AUROC_V11]
    colors = ['#B0B8C4', '#E36C09', '#B0B8C4', '#1A5CB0']
    bars = ax.bar(cats, vals, color=colors, width=0.55,
                  edgecolor='#0D1B2A', linewidth=0.5, zorder=3)
    for i, v in enumerate(vals):
        ax.text(i, v + 0.008, f'{v:.4f}', ha='center', va='bottom',
                fontsize=11, fontweight='bold', color='#0D1B2A')
    ax.axhline(y=AUROC_V1_REST, color='#C0392B', linewidth=1.2, linestyle='--',
               label=f'Honest V1 baseline {AUROC_V1_REST:.3f}')
    ax.set_ylim(0.75, 1.0)
    ax.set_ylabel('Nested-LOVO AUROC', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    ax.legend(fontsize=8, frameon=False, loc='upper left')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'auroc_progression.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_ridge_metrics():
    fig, ax = plt.subplots(figsize=(10, 3.5))
    metrics = ['AUC-ROC', 'Recall\n(Youden)', 'Specificity\n(Youden)', 'F1 Score', 'Calibration\nSlope']
    values  = [AUROC_V11 * 100, 13/14 * 100, 15/20 * 100, F1_SCORE * 100, CAL_SLOPE * 100]
    bars = ax.bar(metrics, values, color='#1A5CB0', width=0.55,
                  edgecolor='#0D1B2A', linewidth=0.5, zorder=3)
    labels = [f'{AUROC_V11*100:.1f}%', '13/14', '15/20', f'{F1_SCORE:.3f}', f'{CAL_SLOPE:.2f}']
    for i, (v, l) in enumerate(zip(values, labels)):
        ax.text(i, v + 1.2, l, ha='center', va='bottom',
                fontsize=11, fontweight='bold', color='#0D1B2A')
    ax.set_ylim(0, 115)
    ax.set_ylabel('Percentage / scaled value', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    ax.set_facecolor('white')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'ridge_metrics.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_confusion():
    fig, ax = plt.subplots(figsize=(4.5, 3.5))
    # Layout: rows = Actual (Failed, Healthy), cols = Predicted (Failed, Healthy)
    # -> [[TP, FN], [FP, TN]] (review fix W1: off-diagonals were transposed)
    matrix = np.array([[TP, FN], [FP, TN]])
    ax.imshow(matrix, cmap='Blues', vmin=0, vmax=15)
    labels = [[f'TP\n{TP}', f'FN\n{FN}'], [f'FP\n{FP}', f'TN\n{TN}']]
    for i in range(2):
        for j in range(2):
            c = 'white' if matrix[i, j] > 7 else '#0D1B2A'
            ax.text(j, i, labels[i][j], ha='center', va='center',
                    fontsize=14, fontweight='bold', color=c)
    ax.set_xticks([0, 1])
    ax.set_yticks([0, 1])
    ax.set_xticklabels(['Predicted\nFailed', 'Predicted\nHealthy'],
                       fontsize=10, color='#0D1B2A')
    ax.set_yticklabels(['Actual\nFailed', 'Actual\nHealthy'],
                       fontsize=10, color='#0D1B2A')
    ax.set_title(f'34-fold Nested LOVO, per-fold Youden ({N_VINS} VINs)',
                 fontsize=11, fontweight='bold', color='#0D1B2A', pad=10)
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'confusion.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_feature_coefs():
    fig, ax = plt.subplots(figsize=(10, 3.0))
    names = [f[0] for f in FEATURES]
    vals = [f[1] for f in FEATURES]
    colors = ['#1A5CB0' if v > 0 else '#C58B1F' for v in vals]
    ax.barh(names[::-1], vals[::-1], color=colors[::-1], height=0.55,
            edgecolor='#0D1B2A', linewidth=0.5)
    for i, v in enumerate(vals[::-1]):
        off = 0.02 if v >= 0 else -0.02
        ha = 'left' if v >= 0 else 'right'
        ax.text(v + off, i, f'{v:+.3f}', ha=ha, va='center',
                fontsize=10, fontweight='bold', color='#0D1B2A')
    ax.axvline(x=0, color='#B0B8C4', linewidth=1)
    ax.set_xlim(-0.7, 1.15)
    ax.set_xlabel('Standardized Ridge coefficient (production refit)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'feature_coefs.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_archetypes():
    """Failure archetype split for 14 failed trucks (VIN14_F is mixed A1+A2)."""
    fig, ax = plt.subplots(figsize=(10, 3.0))
    cats = ['A1 Solenoid\nIntermittency', 'A2 Battery\nCascade', 'A3 Volatility\nDrift', 'A4 Silent /\nAbrupt']
    counts = [3, 5, 3, 4]
    colors = ['#1A5CB0', '#C58B1F', '#2E5090', '#C0392B']
    ax.bar(cats, counts, color=colors, width=0.5, edgecolor='#0D1B2A', linewidth=0.5)
    vins = ['VIN10, 14, 1_F', 'VIN2, 3, 6, 13, 14_F', 'VIN7, 11, 12_F', 'VIN4, 5, 8, 9_F']
    for i, (v, lab) in enumerate(zip(counts, vins)):
        ax.text(i, v + 0.12, f'{v} trucks', ha='center', va='bottom',
                fontsize=13, fontweight='bold', color='#0D1B2A')
        ax.text(i, v / 2, lab, ha='center', va='center',
                fontsize=8, fontweight='bold', color='white')
    ax.set_ylim(0, 6.5)
    ax.set_ylabel('Number of Failures', fontsize=10, color='#606060')
    ax.set_title('14 failed trucks; VIN14_F is mixed A1+A2 (counted in both)',
                 fontsize=9, color='#606060', pad=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'archetypes.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_pipeline():
    fig, ax = plt.subplots(figsize=(12, 2.5))
    stages = [
        ('CAN\nIngestion', '#1A5CB0'),
        ('Sentinel\nCleaning', '#2E5090'),
        ('Crank Event\nCatalog', '#1A5CB0'),
        ('Weekly\nAggregation', '#2E5090'),
        ('Windowed\nFeatures', '#0D1B2A'),
        ('Nested-LOVO\nRidge', '#2E5090'),
        ('Alert\nChannels', '#0D1B2A'),
        ('Horizon\nStatement', '#C58B1F'),
        ('Risk\nTier', '#1E8C45'),
    ]
    for i, (label, color) in enumerate(stages):
        x = i * 1.3
        rect = mpatches.FancyBboxPatch((x, 0.3), 1.0, 1.4, boxstyle="round,pad=0.1",
                                        facecolor=color, edgecolor='#0D1B2A', linewidth=1)
        ax.add_patch(rect)
        ax.text(x + 0.5, 1.0, label, ha='center', va='center', fontsize=9,
                fontweight='bold', color='white')
        if i < len(stages) - 1:
            ax.annotate('', xy=(x + 1.15, 1.0), xytext=(x + 1.05, 1.0),
                       arrowprops=dict(arrowstyle='->', color='#C58B1F', lw=2))
    ax.set_xlim(-0.3, len(stages) * 1.3)
    ax.set_ylim(-0.1, 2.2)
    ax.axis('off')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'pipeline.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_horizon():
    """Prequential detection-horizon curve from V1_1_SM_horizon_curve.csv."""
    ks, aurocs, los, his = [], [], [], []
    with open(RESULTS / "V1_1_SM_horizon_curve.csv", newline='') as f:
        for row in csv.DictReader(f):
            ks.append(int(row['k_weeks']))
            aurocs.append(float(row['auroc']))
            los.append(float(row['ci95_lo']))
            his.append(float(row['ci95_hi']))
    ks, aurocs = np.array(ks), np.array(aurocs)
    los, his = np.array(los), np.array(his)

    fig, ax = plt.subplots(figsize=(11, 3.4))
    ax.fill_between(ks, los, his, color='#1A5CB0', alpha=0.15, label='Bootstrap 95% CI')
    ax.plot(ks, aurocs, color='#1A5CB0', linewidth=2.2, marker='o', markersize=4,
            markerfacecolor='#0D1B2A', label='LOVO AUROC at cut k')
    ax.axhline(y=0.75, color='#C0392B', linewidth=1.5, linestyle='--',
               label='Sustained-detection threshold 0.75')
    ax.axhline(y=0.5, color='#B0B8C4', linewidth=1.2, linestyle=':', label='Chance (0.5)')
    ax.axvline(x=K_STAR, color='#C58B1F', linewidth=2, linestyle='-')
    ax.text(K_STAR + 0.2, 0.97, f'k* = {K_STAR} weeks', fontsize=11, fontweight='bold',
            color='#C58B1F', va='top')
    ax.annotate(f'k=0: {HORIZON_K0:.3f}', xy=(0, HORIZON_K0), xytext=(1.2, 1.0),
                fontsize=9, fontweight='bold', color='#0D1B2A',
                arrowprops=dict(arrowstyle='->', color='#606060', lw=1))
    ax.text(23.5, 0.66, f'far-tail mean {HORIZON_TAIL:.3f}\n(CIs span 0.5)',
            fontsize=8.5, color='#606060', ha='center')
    ax.set_xlabel('k = weeks of data removed before each truck\'s end of history', fontsize=10, color='#606060')
    ax.set_ylabel('AUROC', fontsize=10, color='#606060')
    ax.set_xlim(-0.5, 26.5)
    ax.set_ylim(0.25, 1.05)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    ax.legend(fontsize=8, frameon=False, loc='lower left', ncol=2)
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'horizon.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_alert_leads():
    """Combined-policy first-fire lead per failed truck (vs t_end)."""
    data = [  # (vin, lead_days, first_channel)
        ('VIN5_F', 392, 'persistence'), ('VIN13_F', 301, 'persistence'),
        ('VIN11_F', 266, 'persistence'), ('VIN7_F', 266, 'persistence'),
        ('VIN14_F', 245, 'persistence'), ('VIN3_F', 168, 'persistence'),
        ('VIN6_F', 168, 'persistence'), ('VIN1_F', 160, 'A1 burst'),
        ('VIN10_F', 160, 'A1 burst'), ('VIN12_F', 128, 'A1 burst'),
        ('VIN8_F', 98, 'persistence'), ('VIN2_F', 77, 'persistence'),
        ('VIN4_F', 28, 'persistence'), ('VIN9_F', 0, 'NONE'),
    ]
    fig, ax = plt.subplots(figsize=(11, 3.4))
    vins = [d[0] for d in data]
    leads = [d[1] for d in data]
    colors = ['#C0392B' if d[2] == 'NONE' else ('#C58B1F' if d[2] == 'A1 burst' else '#1A5CB0')
              for d in data]
    ax.barh(vins[::-1], leads[::-1], color=colors[::-1], height=0.6,
            edgecolor='#0D1B2A', linewidth=0.5)
    for i, (v, l, ch) in enumerate(data[::-1]):
        label = 'NO CHANNEL FIRES (A4 + SMA-dead)' if ch == 'NONE' else f'{l}d ({ch})'
        c = '#C0392B' if ch == 'NONE' else '#0D1B2A'
        ax.text(max(l, 2) + 4, i, label, ha='left', va='center',
                fontsize=8, fontweight='bold', color=c)
    ax.axvline(x=PERS_MEDIAN_LEAD, color='#1E8C45', linewidth=1.5, linestyle='--',
               label=f'Median first-fire lead {PERS_MEDIAN_LEAD}d')
    ax.set_xlim(0, 480)
    ax.set_xlabel('First-fire lead before end of telemetry (days)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060', labelsize=8)
    ax.legend(fontsize=8, frameon=False, loc='lower right')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'alert_leads.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_nf_tiers():
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    cats = ['GREEN', 'AMBER', 'RED']
    counts = [NF_GREEN, NF_AMBER, NF_RED]
    colors = ['#1E8C45', '#E36C09', '#C0392B']
    ax.bar(cats, counts, color=colors, width=0.5, edgecolor='#0D1B2A', linewidth=0.8, zorder=3)
    notes = [map_nf_text(n) for n in ['', 'VIN2_NF, VIN10_NF', 'VIN5_NF, VIN20_NF']]
    for i, (v, n) in enumerate(zip(counts, notes)):
        ax.text(i, v + 0.3, f'{v} trucks', ha='center', va='bottom',
                fontsize=12, fontweight='bold', color='#0D1B2A')
        if n:
            ax.text(i, v + 1.8, n, ha='center', va='bottom', fontsize=8, color='#606060')
    ax.set_ylim(0, 19)
    ax.set_ylabel('In-service trucks (20)', fontsize=10, color='#606060')
    ax.set_title('Recalibrated tier distribution -- 20 NF trucks', fontsize=11,
                 fontweight='bold', color='#0D1B2A', pad=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'nf_tiers.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ── GENERATE CHARTS ────────────────────────────────────────────
print("Generating charts...")
img_auroc_prog = chart_auroc_progression()
img_ridge      = chart_ridge_metrics()
img_confusion  = chart_confusion()
img_coefs      = chart_feature_coefs()
img_archetypes = chart_archetypes()
img_pipeline   = chart_pipeline()
img_horizon    = chart_horizon()
img_leads      = chart_alert_leads()
img_nf_tiers   = chart_nf_tiers()
print("Charts generated.")

# Existing stable visualizations to embed.
# Daily-risk dashboards verified by the 2026-06-10 data audit (5 forecast-endpoint
# graphs re-rendered) -- now embedded in the Daily-Risk RUL Views slides.
VIZ_FLEET_RISK = str(GRAPHS / "V1_1_SM_fleet_risk.png")
VIZ_VIN8_F     = str(GRAPHS / "V1_1_SM_VIN8_F_SM_dashboard.png")
VIZ_VIN6_F     = str(GRAPHS / "V1_1_SM_VIN6_F_SM_dashboard.png")
VIZ_DAILY_VIN1_F  = str(GRAPHS / "V1_1_SM_daily_risk_VIN1_F_SM_dashboard.png")
VIZ_DAILY_VIN6_F  = str(GRAPHS / "V1_1_SM_daily_risk_VIN6_F_SM_dashboard.png")
VIZ_DAILY_VIN1_NF = str(GRAPHS / f"V1_1_SM_daily_risk_{display_label('VIN1_NF_SM')}_dashboard.png")
for p in (VIZ_FLEET_RISK, VIZ_VIN8_F, VIZ_VIN6_F,
          VIZ_DAILY_VIN1_F, VIZ_DAILY_VIN6_F, VIZ_DAILY_VIN1_NF):
    if not Path(p).exists():
        raise FileNotFoundError(f"Required image missing: {p}")


# ══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════
prs = make_prs()

# ─── SLIDE 1: TITLE ─────────────────────────────────────────
print("Building slide 1: Title")
s = blank_slide(prs)
add_rect(s, Inches(0), Inches(0), SW, SH, NAVY)
add_rect(s, Inches(0), Inches(5.7), SW, Inches(0.06), GOLD)

add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
         "DAIMLER STARTER MOTOR PREDICTIVE MAINTENANCE",
         font_size=34, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.5),
         "Calibrated Fleet Risk Tiers, Validated Early-Warning Alerts & a Measured Detection Horizon",
         font_size=20, color=GOLD)
add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.4),
         f"BharatBenz 5528T  |  Fleet of {N_VINS} VINs ({N_FAILED}F + {N_NF}NF)  |  V1.1_SM",
         font_size=14, color=LIGHT_GREY)
add_text(s, Inches(0.8), Inches(5.9), Inches(5), Inches(0.3),
         "June 2026  |  Version: V1.1_SM", font_size=12, color=LIGHT_GREY)
add_text(s, Inches(7), Inches(5.9), Inches(5), Inches(0.3),
         "Prepared by Data Science & Engineering Team", font_size=12,
         color=LIGHT_GREY, align=PP_ALIGN.RIGHT)
add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3),
         "BytEdge CONFIDENTIAL", font_size=10, bold=True, color=GOLD)


# ─── SLIDE 2: EXECUTIVE SUMMARY ─────────────────────────────
print("Building slide 2: Executive Summary")
s = blank_slide(prs)
add_header_bar(s, "EXECUTIVE SUMMARY",
               f"Starter motor failure prediction across {N_VINS} BharatBenz 5528T trucks")

y1, tw, th, gap, x0 = Inches(1.15), Inches(2.35), Inches(0.9), Inches(0.18), Inches(0.4)

add_kpi_tile(s, x0, y1, tw, th, "Nested-LOVO AUROC", f"{AUROC_V11*100:.1f}%", KPI_BLUE, "PASS (>89.3% V1)")
add_kpi_tile(s, x0+tw+gap, y1, tw, th, "Recall (Youden)", RECALL_YOUDEN, KPI_BLUE, f"SPEC {SPEC_YOUDEN}")
add_kpi_tile(s, x0+2*(tw+gap), y1, tw, th, "RED-Tier Recall", RECALL_RED, KPI_BLUE, f"SPEC {SPEC_RED}")
add_kpi_tile(s, x0+3*(tw+gap), y1, tw, th, "Calibration Slope", f"{CAL_SLOPE:.2f}", KPI_GREEN, "PASS [0.5, 2]")
add_kpi_tile(s, x0+4*(tw+gap), y1, tw, th, "Detection Horizon", f"{K_STAR} weeks", KPI_AMBER, "DECAY VERIFIED")

y2 = Inches(2.25)
add_kpi_tile(s, x0, y2, tw, th, "Failed Trucks", str(N_FAILED), KPI_RED, "POST-HOC")
add_kpi_tile(s, x0+tw+gap, y2, tw, th, "In-Service Trucks", str(N_NF), KPI_GREEN, f"{NF_GREEN} GRN / {NF_AMBER} AMB / {NF_RED} RED")
add_kpi_tile(s, x0+2*(tw+gap), y2, tw, th, "Features (Ridge)", str(N_FEATURES), KPI_BLUE, "ADMISSIBILITY-AUDITED")
add_kpi_tile(s, x0+3*(tw+gap), y2, tw, th, "A2 Cascade Alert", f"{A2_RECALL} F", KPI_GREEN, f"{A2_FP} NF, ~9.5 wk LEAD")
add_kpi_tile(s, x0+4*(tw+gap), y2, tw, th, "Permutation p", f"{PERMUTATION_P}", KPI_GREEN, "PASS (N=200 floor)")

add_multiline(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.8), [
    ("Objective: Predict starter motor failure risk and provide validated early warning from CAN bus voltage and crank telemetry.", True, NAVY),
    (f"Ridge classifier achieves nested-LOVO AUROC {AUROC_V11} (95% CI [{BOOTSTRAP_CI[0]}, {BOOTSTRAP_CI[1]}], permutation p={PERMUTATION_P}) with {N_FEATURES} audited features. Recall {RECALL_YOUDEN} at the Youden point; RED tier gives {RECALL_RED} at {SPEC_RED} specificity.", False, DARK_TEXT),
    (f"Honest-validation differentiator: V1's reported 0.9214 was restated to {AUROC_V1_REST} under the fully nested protocol before being beaten. V1.1's selection optimism is +{OPTIMISM_V11} (vs V1's hidden +{OPTIMISM_V1}). Probabilities are calibrated (slope {CAL_SLOPE}, Brier {BRIER}) and shippable.", False, DARK_TEXT),
    (f"Three validated alert channels (A2 battery-cascade ~9.5-week (66.5 d) lead with {A2_FP} NF false alarms; tier-gated persistence flag; A1 crank-burst corroborator) cover {COMBINED_RECALL} failed trucks. Detection horizon: k* = {K_STAR} weeks, decay-to-chance verified. ~4/14 silent/abrupt failures remain structurally invisible.", False, DARK_TEXT),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    f"Nested AUROC {AUROC_V11} beats the honestly-restated V1 baseline ({AUROC_V1_REST}) under a strictly harder protocol",
    f"V1's worst miss (VIN8_F_SM, P=0.303) is now caught at recalibrated 0.716, RED tier",
    f"New capability: validated early warning with a measured {K_STAR}-week horizon -- V1 had none (its trend battery fired on 90% of healthy trucks)",
    f"{NF_RED}/20 in-service trucks are RED (VIN5_NF, VIN20_NF) and should be prioritized for starter + battery circuit inspection",
], top=Inches(5.4))
add_footer(s)


# ─── SLIDE 3: BUSINESS CONTEXT ──────────────────────────────
print("Building slide 3: Business Context")
s = blank_slide(prs)
add_header_bar(s, "BUSINESS CONTEXT",
               "Why starter motor predictive maintenance matters for heavy-duty fleets")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.35),
         "THE CHALLENGE", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(1.6), [
    "A failed starter motor strands the truck wherever it last parked -- no crank, no trip",
    "No advance warning in the current time/km-based maintenance regime",
    "Starter failures are entangled with battery health -- the wrong part gets replaced",
    "Fleet telemetry (CAN bus) streams voltage and crank activity -- untapped",
    "14 of 34 study trucks failed within the observation window",
], font_size=11, bullet=True, color=DARK_TEXT)

add_text(s, Inches(6.8), Inches(1.15), Inches(6), Inches(0.35),
         "OBJECTIVES", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(6.8), Inches(1.55), Inches(6), Inches(1.6), [
    "Classify fleet into risk tiers (GREEN / AMBER / RED) with calibrated probabilities",
    "Ship only leakage-audited features under a fully nested validation protocol",
    "Provide validated early-warning alerts with measured lead times and FP burden",
    "Route battery-vs-starter triage so the cheapest correct intervention happens first",
    "Quantify what is achievable (and what is not) with 5-second telemetry",
], font_size=11, bullet=True, color=DARK_TEXT)

add_text(s, Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.35),
         "FOUR-LAYER APPROACH (methodology inherited from the Alternator V10.6.2 program)", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(1.4), [
    (f"Layer 1 -- Fleet Risk Model: {N_FEATURES}-feature nested-LOVO Ridge, calibrated probability + GREEN/AMBER/RED tier (AUROC {AUROC_V11})", True, NAVY),
    ("Layer 2 -- Early-Warning Alerts: A2 battery-cascade detector (~9.5-wk lead, 0/20 NF), tier-gated persistence flag, A1 crank-burst corroborator", True, NAVY),
    (f"Layer 3 -- Honest RUL: validity-horizon statement (flagged truck typically within ~{K_STAR} weeks of failure) instead of unsupportable day-precision dates", True, NAVY),
    ("Layer 4 -- Explainability: exact per-truck linear attributions, archetype assignment, raw-unit counterfactual with every alert", True, NAVY),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "Starter failure is a no-start stranding event; battery-vs-starter misdiagnosis is the single largest avoidable-cost lever (DICV A6)",
    "Four-layer architecture: calibrated risk tier, validated alerts, horizon statement, and an explanation card per truck",
    "Honest framing: ~4/14 failures are silent/abrupt and invisible to any voltage-based method at 5 s sampling",
])
add_footer(s)


# ─── SLIDE 4: PROBLEM STATEMENT / ARCHETYPES ────────────────
print("Building slide 4: Problem Statement")
s = blank_slide(prs)
add_header_bar(s, "PROBLEM STATEMENT",
               "Characterizing the starter motor failure landscape")

add_multiline(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5), [
    ("Key Question: Can we predict starter motor failure risk and warn early from CAN bus telemetry?", True, NAVY),
], font_size=14)

add_text(s, Inches(0.5), Inches(1.75), Inches(6.0), Inches(0.3),
         "FAILURE ARCHETYPES (DATA-DERIVED)", font_size=14, bold=True, color=NAVY)
s.shapes.add_picture(img_archetypes, Inches(0.4), Inches(2.1), Inches(6.2), Inches(2.8))

add_text(s, Inches(6.9), Inches(1.75), Inches(6), Inches(0.3),
         "WHAT EACH ARCHETYPE MEANS", font_size=14, bold=True, color=NAVY)
arch_data = [
    ['Archetype', 'Mechanism', 'Observable?'],
    ['A1 Solenoid', 'Intermittent engagement, crank retry bursts', 'Yes -- days-to-weeks'],
    ['A2 Battery', 'Rest-floor sag + deepening dips (cascade)', 'Yes -- best channel (~9 wk)'],
    ['A3 Volatility', 'Within-week voltage noise drifting up', 'Yes -- weeks'],
    ['A4 Silent/Abrupt', 'Windings, seizure, control loss, telemetry death', 'NO -- detection ceiling'],
]
add_table_shape(s, Inches(6.9), Inches(2.1), Inches(6.0), 5, 3, arch_data,
                col_widths=[1.4, 2.8, 1.8])

add_multiline(s, Inches(6.9), Inches(4.05), Inches(6.0), Inches(0.9), [
    "5/14 failed trucks go silent 32-142 days before the recorded failure date",
    "7/34 trucks have an SMA-dead telematics config -- crank features unobservable",
    "A4 (4/14) sets the honest recall ceiling at ~10-11/14 for lead-time alerting",
], font_size=9, bullet=True)

add_key_takeaways(s, [
    "Three of four archetypes (A1/A2/A3, ~10-11 trucks) are observable in voltage/crank telemetry -- A4 is not",
    "A2 battery cascade is the best-observed pathway and the prime confound: weak battery accelerates solenoid wear and vice versa",
    "Silent gaps and SMA-dead configs are telemetry-architecture problems, not modeling problems",
    "The physics audit predicted the recall ceiling (~10-11/14) before modeling -- the model lands exactly there at the RED tier",
], top=Inches(5.15))
add_footer(s)


# ─── SLIDE 5: DATA LANDSCAPE ────────────────────────────────
print("Building slide 5: Data Landscape")
s = blank_slide(prs)
add_header_bar(s, "DATA LANDSCAPE",
               "CAN bus telemetry streams and fleet composition")

s.shapes.add_picture(img_pipeline, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.8))

add_text(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(0.3),
         "CAN SIGNAL INVENTORY", font_size=14, bold=True, color=NAVY)
data = [
    ['Signal', 'Meaning', 'Range', 'Role'],
    ['VSI', 'Power Supply Voltage', '0-36V', 'Electrical health (primary)'],
    ['SMA', 'Starter Motor Active', '{0,1}', 'Crank-event catalog'],
    ['RPM', 'Engine Speed', '0-3500 rev/min', 'Crank success / rest detection'],
    ['CSP', 'Vehicle Speed', '0-100 km/h', 'Drive vs rest separation'],
    ['ANR', 'Engine Torque', '-400-1300 Nm', 'Load context'],
    ['GED', 'Alternator Excitation State', '{0,1,2,3}', 'Data-quality covariate only'],
]
add_table_shape(s, Inches(0.5), Inches(3.45), Inches(6.0), 7, 4, data,
                col_widths=[0.8, 2.0, 1.2, 2.0])

add_text(s, Inches(7.0), Inches(3.1), Inches(5.8), Inches(0.3),
         "FLEET COMPOSITION", font_size=14, bold=True, color=NAVY)
fleet_data = [
    ['Category', 'Count', 'Suffix', 'Status'],
    ['Failed', str(N_FAILED), '_F_SM', 'Post-hoc analysis'],
    ['In-Service', str(N_NF), '_NF_SM', 'Active monitoring'],
    ['Total', str(N_VINS), '--', f'~{N_ROWS/1e6:.0f}M CAN rows'],
]
add_table_shape(s, Inches(7.0), Inches(3.45), Inches(5.8), 4, 4, fleet_data,
                col_widths=[1.3, 0.8, 1.2, 2.0])

add_multiline(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(0.5), [
    f"{N_ROWS:,} raw 5-second rows -> {N_TRUCK_WEEKS:,} truck-weeks + {N_CRANK_EVENTS:,} gap-aware crank events",
    "5 s sampling / ~0.2 V resolution; no cranking current, temperature, or battery SoC signal",
    (VIN_FOOTNOTE, False, GREY_MED),
], font_size=9, bullet=True)

add_key_takeaways(s, [
    "VSI (power supply voltage) is again the winning signal family -- 3 of 4 model features derive from it; 1 from crank physics",
    "SM has no GED=2 channel: zero GED2 in all 14 failed trucks -- the ALT emergency layer does not transfer",
    "KT crank claims reconciled: durations only +3% (not +48%); whole-life >5% failed-crank threshold refuted -- only the last-90-day change discriminates (AUROC 0.74)",
], top=Inches(5.7), height=Inches(0.95))
add_footer(s)


# ─── SLIDE 6: HONEST VALIDATION & LEAKAGE AUDIT ─────────────
print("Building slide 6: Honest Validation & Leakage Audit")
s = blank_slide(prs)
add_header_bar(s, "HONEST VALIDATION & LEAKAGE AUDIT",
               "V1's headline was restated before being beaten -- the differentiator of this program")

s.shapes.add_picture(img_auroc_prog, Inches(0.4), Inches(1.15), Inches(6.6), Inches(2.6))

add_text(s, Inches(7.3), Inches(1.15), Inches(5.6), Inches(0.3),
         "WHY THE RESTATEMENT", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(7.3), Inches(1.5), Inches(5.6), Inches(2.2), [
    (f"V1 reported 0.9214 -- but feature screening + subset selection sat outside the LOVO loop (+{OPTIMISM_V1} optimism), and one recall point came from a post-hoc pooled threshold.", False, DARK_TEXT),
    (f"Nested restatement: V1's true out-of-sample AUROC is {AUROC_V1_REST} (recall 12/14).", False, DARK_TEXT),
    (f"V1's top-2 feature (vsi_dominant_freq) was a 1/n_weeks artifact -- it collapses 0.748 to 0.525 under a fixed-window control. Banned.", False, DARK_TEXT),
    (f"V1.1 redoes screening, subset search, threshold, and Platt recalibration inside every fold: optimism +{OPTIMISM_V11}.", False, DARK_TEXT),
    (f"Ablation: the nested protocol on V1-era features scores {AUROC_ABLATION} -- the V1.1 gain (+0.089) is feature engineering, not protocol arithmetic.", False, DARK_TEXT),
], font_size=9.5, bullet=True)

add_text(s, Inches(0.5), Inches(3.78), Inches(12.3), Inches(0.3),
         "LEAKAGE GATES (ALL ENFORCED PER REFIT)", font_size=14, bold=True, color=NAVY)
gate_data = [
    ['Gate', 'Hazard it controls', 'Result'],
    ['Leak ceilings measured', 'n_weeks alone classifies at AUROC 0.952; t_start at 0.893 (failed trucks stop transmitting because they failed)', 'All such features banned'],
    ['G1 fixed-L40 window control', 'Observation-length artifacts laundered through denominators', 'Drop 0.0000 -- matrices bit-identical'],
    ['G3 prequential time-locking', 'Epoch/length leak vs true failure signal', 'Screening-feature curve: holds to k=10 wks, chance at k=11 (0.536) -- failure-locked. NB: the frozen final-model curve (slide 16) decays slower: 0.704 at k=11'],
    ['G4 winner-subset stability', 'Fold-to-fold feature-selection instability', 'Strict modal criterion FAIL (14/34 < 17/34; a k=3/k=4 tie). Substantively stable: core pair vsi_withinwk + range_trend in 34/34 folds; only 3 distinct subsets total'],
    ['G6 banned-token scan', 'vsi_dominant_freq, calendar features, cumulative counts', '0 banned tokens in selected features'],
    ['Permutation test (full pipeline)', 'Selection noise masquerading as signal', f'p = {PERMUTATION_P} (0/200 shuffles reach {AUROC_V11})'],
]
add_table_shape(s, Inches(0.4), Inches(4.08), Inches(12.5), 7, 3, gate_data,
                col_widths=[2.6, 5.4, 4.5])

add_key_takeaways(s, [
    f"V1 0.9214 (reported) -> {AUROC_V1_REST} (restated, nested) -> {AUROC_V11} (V1.1 nested) -- measured under a strictly harder protocol; vs V1-as-reported the recall swap is VIN8_F gained, VIN9_F (nominally caught by V1 at P=0.4825) lost",
    "The leak ceilings in this data (AUROC 0.95 from data volume alone) exceed any honest model -- which is why the gates exist and are binding; G4 is the one strict FAIL, disclosed above",
], top=Inches(6.58), height=Inches(0.55))
add_footer(s)


# ─── SLIDE 7: SECTION DIVIDER ───────────────────────────────
print("Building slide 7: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 4, "MODELING APPROACH",
                  "Four-layer architecture: Risk Tiers -> Alerts -> Horizon -> Explanations")


# ─── SLIDE 8: MODELING APPROACH ──────────────────────────────
print("Building slide 8: Modeling Approach")
s = blank_slide(prs)
add_header_bar(s, "MODELING APPROACH",
               "Nested-LOVO Ridge + validated alert channels + horizon statement + explanations")

add_text(s, Inches(0.5), Inches(1.15), Inches(2.95), Inches(0.3),
         "LAYER 1: RISK MODEL", font_size=13, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.5), Inches(2.95), Inches(1.9), [
    "RidgeClassifier(alpha=1.0)",
    f"{N_FEATURES} audited features",
    f"Fully nested 34-fold LOVO: AUROC {AUROC_V11}",
    f"95% CI [{BOOTSTRAP_CI[0]}, {BOOTSTRAP_CI[1]}], p={PERMUTATION_P}",
    f"Per-fold Platt recalibration (slope {CAL_SLOPE})",
    "Tiers: GREEN <0.35 <= AMBER <0.55 <= RED",
    "Score weekly, review monthly",
], font_size=9, bullet=True)

add_text(s, Inches(3.7), Inches(1.15), Inches(2.95), Inches(0.3),
         "LAYER 2: ALERTS", font_size=13, bold=True, color=NAVY)
add_multiline(s, Inches(3.7), Inches(1.5), Inches(2.95), Inches(1.9), [
    f"A2 battery-cascade: {A2_RECALL} archetype, {A2_FP} NF, median lead {A2_MEDIAN_LEAD:.1f}d",
    f"Persistence flag: {PERS_RECALL}, {PERS_FP} NF -- tier-gated condition flag only",
    f"A1 crank-burst: corroborator only ({A1_FP_RATE} FP eps/truck-yr standalone)",
    f"Combined: {COMBINED_RECALL} failed fire >=1 channel; {NF_CLEAN} NF fully clean",
], font_size=9, bullet=True)

add_text(s, Inches(6.9), Inches(1.15), Inches(2.95), Inches(0.3),
         "LAYER 3: HONEST RUL", font_size=13, bold=True, color=NAVY)
add_multiline(s, Inches(6.9), Inches(1.5), Inches(2.95), Inches(1.9), [
    "No day-precision RUL, no survival layer",
    f"Validity-horizon statement: flagged truck typically within ~{K_STAR} weeks",
    "Tier -> maintenance window: RED 2-4 wks; AMBER next service",
    f"Proof of ceiling: hazard RUL MAE {SURV_RUL_MAE}d vs {CONST_RUL_MAE}d for a constant",
], font_size=9, bullet=True)

add_text(s, Inches(10.1), Inches(1.15), Inches(2.9), Inches(0.3),
         "LAYER 4: EXPLAINABILITY", font_size=13, bold=True, color=NAVY)
add_multiline(s, Inches(10.1), Inches(1.5), Inches(2.9), Inches(1.9), [
    "Exact linear attribution (coef x z = SHAP for a linear model)",
    "34 per-VIN explanation cards",
    "Archetype assignment per truck",
    "Raw-unit counterfactual per flagged truck",
    "Governance-grade model card",
], font_size=9, bullet=True)

add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.3),
         "DECISION ENGINE: HOW THE LAYERS COMBINE", font_size=14, bold=True, color=NAVY)
dec_data = [
    ['Condition', 'Action', 'Timing', 'Confidence'],
    ['A2 battery-cascade fires', 'BATTERY-FIRST INSPECTION (DICV A6 routing)', 'Weeks (~9.5 wk lead)', f'High ({A2_FP} NF false alarms)'],
    ['Tier RED', 'STARTER + BATTERY CIRCUIT INSPECTION', '2-4 weeks', f'High ({SPEC_RED} specificity)'],
    ['Tier AMBER + persistence flag active', 'BUNDLE INTO NEXT SCHEDULED SERVICE', 'Months', 'Medium (tier-gated flag)'],
    ['A1 crank burst on AMBER/RED truck', 'CORROBORATE -- ADVANCE THE INSPECTION', 'Days-weeks', 'Medium (rescued VIN1_F)'],
    ['Tier GREEN, no channel', 'NORMAL OPERATION', f'Score valid ~2.3 months (70 d)', 'High (10/20 NF fully clean)'],
]
add_table_shape(s, Inches(0.5), Inches(3.95), Inches(12.3), 6, 4, dec_data,
                col_widths=[3.4, 4.0, 1.9, 3.0])

add_key_takeaways(s, [
    "Four layers, each shipping only LOVO/physics-validated numbers -- alert channels are gated behind tiers, never standalone pagers",
    "A2 battery-cascade is the only short-fuse (~1-3 month) signal and the highest-confidence channel (0/20 NF false alarms)",
    "Alert channels recover 3 of the 4 Layer-1 GREEN-tier misses (VIN1/3/4_F) -- redundant coverage by design",
], top=Inches(6.15), height=Inches(0.85))
add_footer(s)


# ─── SLIDE 9: FEATURE ENGINEERING ────────────────────────────
print("Building slide 9: Feature Engineering")
s = blank_slide(prs)
add_header_bar(s, "FEATURE ENGINEERING",
               f"{N_FEATURES} winners from a 10-candidate audited pool, selected inside every fold")

s.shapes.add_picture(img_coefs, Inches(0.4), Inches(1.15), Inches(6.2), Inches(2.6))

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.3),
         "FEATURE DESCRIPTIONS & PHYSICS", font_size=14, bold=True, color=NAVY)
feat_data = [
    ['Feature', 'Physics', 'Coef'],
    ['vsi_withinwk_std_ratio_30d_w', 'Within-week voltage noise vs own baseline -- volatility drift (univariate AUROC 0.921)', '+0.886'],
    ['rest_vsi_p05_delta90', 'Engine-off rest-voltage floor sag -- battery cascade (step-aware re-baseline)', '-0.270'],
    ['vsi_range_trend', 'Weekly voltage-range slope -- suppressor (physics says widening = risk; flipped by r=+0.82 collinearity)', '-0.414'],
    ['dip_depth_last90_delta', 'Crank dip widening vs own baseline -- heavier cranking load', '+0.141'],
]
add_table_shape(s, Inches(7.0), Inches(1.55), Inches(5.8), 5, 3, feat_data,
                col_widths=[2.1, 2.9, 0.8])

add_multiline(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.0), [
    ("Selection: V1-faithful screening + exhaustive k=3-6 subset search redone inside each of 34 folds. Core pair selected 34/34 folds; only 3 distinct subsets ever chosen -- though gate G4's strict modal criterion FAILS (14/34 < 17/34, a k=3/k=4 tie; disclosed on the gates slide).", True, NAVY),
    "All features window-anchored (last 40 masked weeks / last-90-day events) and per-VIN baselined -- the L40 control matrix is bit-identical to production (zero-drop by construction).",
    "Two physics candidates (retry-burst, extended-crank-tail) were admissible but fleet-weak -- rejected by in-fold screening 34/34; they live on inside the A1 alert channel instead.",
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "The workhorse is the within-week voltage noise ratio (coef +0.886, univariate AUROC 0.921) -- volatility drift against the truck's own baseline",
    "vsi_range_trend is honestly flagged as a statistical suppressor: its negative weight is a collinearity correction, not physics",
    "Battery-step awareness matters: rest-VSI baselines re-anchored after detected battery replacements (VIN8_F + 5 NF trucks)",
])
add_footer(s)


# ─── SLIDE 10: RESULTS OVERVIEW ─────────────────────────────
print("Building slide 10: Results Overview")
s = blank_slide(prs)
add_header_bar(s, "RESULTS OVERVIEW",
               "Nested 34-fold LOVO performance, operating points and calibration")

s.shapes.add_picture(img_ridge, Inches(0.4), Inches(1.1), Inches(6.2), Inches(2.5))
s.shapes.add_picture(img_confusion, Inches(7.5), Inches(1.1), Inches(4.5), Inches(2.7))

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.3),
         "PERFORMANCE SCORECARD (ALL NESTED-LOVO OUT-OF-FOLD)", font_size=14, bold=True, color=NAVY)
score_data = [
    ['Metric', 'Value', 'Context'],
    ['AUROC', f'{AUROC_V11}', f'Bootstrap 95% CI [{BOOTSTRAP_CI[0]}, {BOOTSTRAP_CI[1]}]; permutation p = {PERMUTATION_P} (N=200 floor)'],
    ['Recall @ per-fold Youden', RECALL_YOUDEN, f'TP {TP}, FP {FP}, FN {FN}, TN {TN} -- recall-greedy operating point'],
    ['RED-tier operating point', f'{RECALL_RED} @ {SPEC_RED}', 'Matched-specificity alternative; choose per maintenance economics'],
    ['F1 / MCC @ Youden', f'{F1_SCORE} / {MCC}', 'Recall-weighted; V1 restated was 12/14 @ 18/20'],
    ['Calibration', f'slope {CAL_SLOPE}, Brier {BRIER}', f'CITL {CITL}; slope in [0.5, 2] -- probabilities shippable (V1 slope was 4.72)'],
    ['Selection optimism', f'+{OPTIMISM_V11}', f'Non-nested 0.9357 vs nested {AUROC_V11}; V1 hid +{OPTIMISM_V1}'],
    ['Jackknife stability', f'{JACKKNIFE[0]}-{JACKKNIFE[1]}', 'AUROC range 0.024, std 0.007 across leave-one-out reruns'],
]
add_table_shape(s, Inches(0.5), Inches(4.15), Inches(12.3), 8, 3, score_data,
                col_widths=[2.7, 1.9, 7.7])

add_key_takeaways(s, [
    f"Two pre-registered operating points: Youden ({RECALL_YOUDEN} recall @ {SPEC_YOUDEN}) and RED tier ({RECALL_RED} @ {SPEC_RED}) -- both honest, neither dominates V1 on every cell",
    "The single miss (VIN9_F_SM, prob 0.401 vs threshold 0.406) is the structurally invisible A4 + SMA-dead + 142-day-silent-gap truck",
    "Calibrated probabilities are a new deliverable: V1 scores were rank-only (slope 4.72); V1.1 ships slope 0.86",
], top=Inches(6.9), height=Inches(0.55))
add_footer(s)


# ─── SLIDE 11: SECTION DIVIDER ──────────────────────────────
print("Building slide 11: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 7, "DETAILED VIN ANALYSIS",
                  "Failed trucks, in-service fleet, alert channels and case studies")


# ─── SLIDE 12: FAILED VIN ANALYSIS ──────────────────────────
print("Building slide 12: Failed VIN Analysis")
s = blank_slide(prs)
add_header_bar(s, "FAILED TRUCK ANALYSIS -- 14 VINs",
               "Out-of-fold probability, tier, archetype and first-firing alert channel")

failed_data = [
    ['VIN', 'OOF P', 'Recal P', 'Tier', 'Archetype', 'First Channel', 'Lead (d)'],
    ['VIN1_F',  '0.406', '0.260', 'GREEN', 'A1 solenoid',      'A1 crank burst', '160'],
    ['VIN2_F',  '0.598', '0.904', 'RED',   'A2 battery',       'persistence',    '77'],
    ['VIN3_F',  '0.438', '0.338', 'GREEN', 'A2 battery',       'persistence',    '168'],
    ['VIN4_F',  '0.429', '0.339', 'GREEN', 'A4 silent/abrupt', 'persistence',    '28'],
    ['VIN5_F',  '0.756', '0.992', 'RED',   'A4 silent/abrupt', 'persistence',    '392'],
    ['VIN6_F',  '0.825', '0.998', 'RED',   'A2 battery',       'persistence',    '168'],
    ['VIN7_F',  '0.612', '0.906', 'RED',   'A3 volatility',    'persistence',    '266'],
    ['VIN8_F',  '0.521', '0.716', 'RED',   'A4 silent/abrupt', 'persistence',    '98'],
    ['VIN9_F',  '0.401', '0.224', 'GREEN', 'A4 silent/abrupt', 'NONE',           '--'],
    ['VIN10_F', '0.782', '0.995', 'RED',   'A1 solenoid',      'A1 crank burst', '160'],
    ['VIN11_F', '0.651', '0.958', 'RED',   'A3 volatility',    'persistence',    '266'],
    ['VIN12_F', '0.658', '0.955', 'RED',   'A3 volatility',    'A1 crank burst', '128'],
    ['VIN13_F', '0.496', '0.654', 'RED',   'A2 battery',       'persistence',    '301'],
    ['VIN14_F', '0.798', '0.998', 'RED',   'A1+A2 mixed',      'persistence',    '245'],
]
add_table_shape(s, Inches(0.4), Inches(1.1), Inches(12.5), 15, 7, failed_data,
                col_widths=[1.1, 0.9, 0.9, 0.9, 2.2, 1.6, 0.9])

add_key_takeaways(s, [
    "10/14 failed trucks score RED on recalibrated probability; 13/14 are caught at the per-fold Youden point",
    "VIN8_F -- V1's worst miss (P=0.303) -- is now RED at 0.716 via the within-week volatility and battery-floor features",
    "VIN9_F is the only full miss on every layer: A4 silent + SMA-dead + 142-day silent gap -- the irreducible blind spot of this dataset",
    "Three GREEN-tier failures (VIN1/3/4_F) are recovered by the alert channels -- the layers are redundant by design",
], top=Inches(6.45), height=Inches(0.95))
add_footer(s)


# ─── SLIDE 13: IN-SERVICE FLEET ─────────────────────────────
print("Building slide 13: In-Service Fleet")
s = blank_slide(prs)
add_header_bar(s, "IN-SERVICE FLEET STATUS -- 20 VINs",
               "Recalibrated out-of-fold risk across the full fleet")

s.shapes.add_picture(VIZ_FLEET_RISK, Inches(0.4), Inches(1.15), Inches(8.0), Inches(4.0))
s.shapes.add_picture(img_nf_tiers, Inches(8.6), Inches(1.15), Inches(4.4), Inches(2.7))

add_multiline(s, Inches(8.6), Inches(4.0), Inches(4.4), Inches(1.2), [
    ("RED: VIN5_NF (0.96), VIN20_NF (0.62) -- inspect starter + battery circuit within 2-4 weeks", False, DARK_TEXT),
    ("AMBER: VIN2_NF (0.45), VIN10_NF (0.43) -- bundle into next scheduled service", False, DARK_TEXT),
    ("Persistence-flag NF trucks (VIN2/5/8/15_NF) are tracked: future failures or first evidence of rule drift -- either is informative", False, DARK_TEXT),
], font_size=9, bullet=True)

add_key_takeaways(s, [
    f"{NF_GREEN}/20 in-service trucks are GREEN; {NF_AMBER} AMBER and {NF_RED} RED warrant action now",
    "NF false alarms may be right-censored degraders: VIN5_NF shows genuinely elevated electrical noise; VIN20_NF is tier-only evidence (SMA-dead truck, no alert channel fired)",
    "Re-score weekly (the 10-week horizon gives only ~2 reads per month otherwise); review tiers monthly",
], top=Inches(5.55))
add_footer(s)


# ─── SLIDE 14: VIN CASE STUDIES ──────────────────────────────
print("Building slide 14: VIN Case Studies")
s = blank_slide(prs)
add_header_bar(s, "VIN CASE STUDIES",
               "The recovered miss and the textbook battery cascade")

add_text(s, Inches(0.4), Inches(1.15), Inches(6.0), Inches(0.25),
         "VIN8_F -- V1's MISS, NOW CAUGHT (RED 0.716)", font_size=11, bold=True, color=GREEN_PASS)
s.shapes.add_picture(VIZ_VIN8_F, Inches(0.4), Inches(1.45), Inches(6.2), Inches(3.0))

add_text(s, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.25),
         "VIN6_F -- A2 BATTERY CASCADE (RED 0.998)", font_size=11, bold=True, color=NAVY)
s.shapes.add_picture(VIZ_VIN6_F, Inches(6.8), Inches(1.45), Inches(6.2), Inches(3.0))

add_multiline(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.0), [
    ("VIN8_F: V1 scored it 0.303 (missed). V1.1's within-week volatility ratio plus the battery-step-aware rest-VSI delta lift it to OOF 0.521 (recal 0.716, RED) -- and the persistence flag fired 98 days before end of telemetry. SMA-dead config and a 37-day silent gap are honestly flagged on its card.", False, DARK_TEXT),
    ("VIN6_F: the textbook A2 battery cascade -- rest floor sagging -3.15 V, crank dips +4.12 V deeper, within-week noise 2.07x baseline. RED at 0.998; the A2 triple detector fired 70 days before failure and routes the work order battery-first.", False, DARK_TEXT),
    ("Two visual views ship per truck: the weekly risk-probability trajectory (shown here) and a daily RUL-style degradation view (see the Daily-Risk RUL Views slides that follow the Detection Horizon).", False, GREY_MED),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "VIN8_F shows the V1.1 gain concretely: new physics features recover a previously-missed failure vs the restated baseline (transparency: vs V1-as-reported, V1.1 trades VIN9_F -- nominally caught at 0.4825 -- for VIN8_F)",
    "VIN6_F shows the triage value: battery-first routing on a cascade signature is the cheapest correct intervention",
], top=Inches(6.3), height=Inches(0.65))
add_footer(s)


# ─── SLIDE 15: ALERT CHANNELS ───────────────────────────────
print("Building slide 15: Alert Channels")
s = blank_slide(prs)
add_header_bar(s, "ALERT CHANNELS & VALIDATED LEADS",
               "Three LOVO/physics-validated channels -- V1 had zero")

s.shapes.add_picture(img_leads, Inches(0.4), Inches(1.15), Inches(7.4), Inches(2.7))

add_text(s, Inches(8.1), Inches(1.15), Inches(4.9), Inches(0.3),
         "SHIP / DON'T-SHIP VERDICTS", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(8.1), Inches(1.5), Inches(4.9), Inches(2.4), [
    (f"A2 battery cascade -- SHIP: {A2_RECALL} battery archetype, {A2_FP} NF, median lead {A2_MEDIAN_LEAD:.1f}d; battery replacements provably don't fire", True, GREEN_PASS),
    (f"Persistence flag -- SHIP tier-gated only: {PERS_RECALL} recall, {PERS_FP} NF at end-state; as a raw walking alarm it visits all 20 NF trucks (31% of weeks) -- never a first-crossing pager", True, KPI_AMBER),
    (f"A1 crank burst -- corroborator only: 4/12 applicable failed fire, but {A1_FP_RATE} FP episodes/truck-year standalone; rescued GREEN-tier VIN1_F", True, KPI_AMBER),
], font_size=9, bullet=True)

add_text(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.3),
         "COMBINED POLICY (TIER + PERSISTENCE + A1 + A2)", font_size=14, bold=True, color=NAVY)
pol_data = [
    ['Quantity', 'Value', 'Honest note'],
    ['Failed trucks firing >=1 channel', f'{COMBINED_RECALL}', 'Persistence first on 10, A1 first on 3; A2 corroborates 4 and is the only ~1-3-month-fuse signal'],
    ['Median first-fire lead', f'{PERS_MEDIAN_LEAD} d', 'Long leads = condition flag, not failure-imminent timer (min 28 d, VIN4_F)'],
    ['NF trucks fully clean', NF_CLEAN, '6 NF show 1 channel, 2 show 2, 2 show 3 -- repeat offenders may be right-censored degraders'],
    ['GREEN-tier failures recovered', '3 of 4', 'VIN1_F (A1 burst), VIN3_F and VIN4_F (persistence); only VIN9_F fires nothing'],
]
add_table_shape(s, Inches(0.5), Inches(4.4), Inches(12.3), 5, 3, pol_data,
                col_widths=[3.0, 1.6, 7.7])

add_key_takeaways(s, [
    "A2 is the deployment star: zero NF false alarms, ~9.5-week (66.5 d) median lead, and a built-in battery-vs-starter triage decision",
    "The persistence rule is honest about its weakness: recall holds out-of-fold (13/14) but NF FP doubled to 4/20 -- shipped tier-gated only",
], top=Inches(6.45), height=Inches(0.6))
add_footer(s)


# ─── SLIDE 16: DETECTION HORIZON ────────────────────────────
print("Building slide 16: Detection Horizon")
s = blank_slide(prs)
add_header_bar(s, "DETECTION HORIZON -- 10 WEEKS",
               "Prequential walk-back of the frozen 4-feature model (X4)")

s.shapes.add_picture(img_horizon, Inches(0.4), Inches(1.15), Inches(12.5), Inches(3.4))

add_multiline(s, Inches(0.5), Inches(4.75), Inches(12.3), Inches(1.0), [
    (f"Method: per VIN, truncate all data k weeks before its own end of history, recompute the 4 frozen features causally, re-run LOVO -- no re-screening. k=0 reconciles exactly to the X2 figure ({HORIZON_K0}).", False, DARK_TEXT),
    (f"Result: AUROC {HORIZON_K0} at k=0 decays to {HORIZON_K10} at k={K_STAR} and {HORIZON_K11} at k=11; far-tail (k=23-26) mean {HORIZON_TAIL} with every CI spanning 0.5.", False, DARK_TEXT),
    ("Why this matters twice: (1) operations gets a validity statement -- a flagged truck is typically within ~10 weeks of failure, a clean score is good for ~2.3 months (70 d); (2) the decay-to-chance is the leak disambiguation -- an observation-length artifact would NOT decay with distance to failure.", False, DARK_TEXT),
    ("Curve identity: this is the FROZEN FINAL-MODEL curve (0.704 at k=11, chance past ~k=20). The G3 gate's 'chance at k=11' (slide 6) is the discovery-phase screening-feature curve (0.536 at k=11) -- two different objects, not a contradiction.", False, GREY_MED),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    f"k* = {K_STAR} weeks: claim 'risk score valid ~2.3 months (70 d) out', never quote failure dates",
    "Decay is gradual (slow-burn degradation signal), reaching chance past ~k=20 -- honest about the k=13-16 hover at 0.62-0.77",
    "This curve is the single strongest evidence that the model reads failure physics, not telemetry bookkeeping",
], top=Inches(6.35), height=Inches(0.75))
add_footer(s)


# ─── DAILY-RISK RUL VIEWS (3 slides) ────────────────────────
# Daily-resolution per-truck dashboards, verified by the 2026-06-10 data audit.
HONEST_RUL_NOTE = ("Honest note: the daily RUL curve is a fleet-Weibull-anchored "
                   "illustration -- the validated deliverable is the calibrated risk "
                   "tier + <=10-week horizon statement, never a failure date.")

def add_daily_risk_slide(title, subtitle, img_path, panel_title, bullets):
    """ALT-style slide: full-aspect daily-risk dashboard left, takeaways right."""
    s = blank_slide(prs)
    add_header_bar(s, title, subtitle)
    # 3198x2065 px (aspect 1.549): width-only embed preserves the ratio (h = 5.52")
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.18), width=Inches(8.55))
    add_text(s, Inches(9.05), Inches(1.18), Inches(3.95), Inches(0.3),
             panel_title, font_size=13, bold=True, color=NAVY)
    add_multiline(s, Inches(9.05), Inches(1.55), Inches(3.95), Inches(3.9),
                  bullets, font_size=9.5, bullet=True)
    add_rounded_rect(s, Inches(9.05), Inches(5.65), Inches(3.95), Inches(1.15),
                     BG_LIGHT, border=KPI_AMBER)
    add_text(s, Inches(9.2), Inches(5.75), Inches(3.65), Inches(0.95),
             HONEST_RUL_NOTE, font_size=8.5, color=GREY_MED)
    add_footer(s)
    return s

print("Building slide 17: Daily-Risk RUL View -- VIN1_F (how to read)")
add_daily_risk_slide(
    "DAILY-RISK RUL VIEWS -- HOW TO READ (VIN1_F)",
    "A1 solenoid archetype -- early warning, a real telemetry gap, and a terminal silent gap on one canvas",
    VIZ_DAILY_VIN1_F,
    "HOW TO READ THE DAILY VIEW",
    [
        ("Curve breaks = no telemetry received. VIN1_F has a real 34-day mid-history gap (2025-07-31 -> 2025-09-02): the line is masked, never interpolated.", False, DARK_TEXT),
        ("Hatched band = terminal silent gap: 72 days with no data before the recorded failure (2025-11-26). Going quiet is itself information.", False, DARK_TEXT),
        ("Dotted segment = fleet-Weibull-anchored projection beyond the last received data -- an illustration layer, not a measurement.", False, DARK_TEXT),
        ("Early-warning annotation marks the A1 crank-burst corroboration; background colors are the action tiers (GREEN -> RED).", False, DARK_TEXT),
        ("Bottom strip: daily VSI mean/range against the truck's own setpoint -- the raw physics behind the score.", False, DARK_TEXT),
    ])

print("Building slide 18: Daily-Risk RUL View -- VIN6_F")
add_daily_risk_slide(
    "DAILY-RISK RUL VIEW -- VIN6_F (A2 BATTERY CASCADE)",
    "The textbook battery cascade: full GREEN -> RED zone progression before failure",
    VIZ_DAILY_VIN6_F,
    "TAKEAWAYS",
    [
        ("Full tier progression GREEN -> AMBER -> RED across the history; recalibrated Ridge risk reaches 0.998 (RED) at end of telemetry.", False, DARK_TEXT),
        ("Cascade physics visible in the VSI strip: rest-voltage floor sagging -3.15 V, crank dips +4.12 V deeper, within-week noise 2.07x baseline.", False, DARK_TEXT),
        ("The A2 battery-cascade alert fired 70 days before failure -- it routes the work order battery-first (DICV A6), the cheapest correct intervention.", False, DARK_TEXT),
        ("This is the trajectory shape the daily view is designed to surface: monotone degradation a weekly tier review can act on early.", False, DARK_TEXT),
    ])

print("Building slide 19: Daily-Risk RUL View -- VIN1_NF (healthy contrast)")
add_daily_risk_slide(
    "DAILY-RISK RUL VIEW -- VIN1_NF (HEALTHY IN-SERVICE)",
    "The healthy contrast: what a GREEN truck's daily view looks like",
    VIZ_DAILY_VIN1_NF,
    "TAKEAWAYS",
    [
        ("Healthy contrast: the truck stays GREEN across its entire history; recalibrated risk 0.07 at last score, status ACTIVE.", False, DARK_TEXT),
        ("The dotted forecast segment is the same fleet-Weibull-anchored projection drawn on every truck -- on a healthy truck it is a planning aid, not a failure prediction.", False, DARK_TEXT),
        ("Flat VSI strip with no rest-floor sag or dip widening -- the physics agrees with the score.", False, DARK_TEXT),
        ("Operating rule: a clean GREEN score is valid ~2.3 months (70 d); re-score weekly, review tiers monthly.", False, DARK_TEXT),
    ])


# ─── SLIDE 20: SECTION DIVIDER ──────────────────────────────
print("Building slide 20: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 8, "OPERATIONAL ASSESSMENT",
                  "Explainability, closed dead ends, limitations and the path forward")


# ─── SLIDE 21: EXPLAINABILITY ───────────────────────────────
print("Building slide 21: Explainability")
s = blank_slide(prs)
add_header_bar(s, "EXPLAINABILITY",
               "Exact attributions, archetypes and counterfactuals -- one card per truck")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.3),
         "HOW IT WORKS", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.4), [
    "Linear model -> exact decomposition: contribution = coefficient x z-scored feature (equals SHAP for a linear model; no approximation, no library)",
    "34 per-VIN explanation cards: ranked drivers in physical language, archetype assignment, data-quality caveats (SMA-dead, silent gaps, battery steps)",
    "Raw-unit counterfactual per flagged truck -- a ceteris paribus statement, not a repair prescription",
    "The suppressor (vsi_range_trend) is flagged on every card: physical value and model use stated separately",
], font_size=10, bullet=True)

add_text(s, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.3),
         "EXAMPLE CARD -- VIN6_F (RED, P 0.998, A2 BATTERY CASCADE)", font_size=12, bold=True, color=NAVY)
add_rounded_rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.9), BG_LIGHT, border=KPI_BLUE)
add_multiline(s, Inches(7.0), Inches(1.65), Inches(5.6), Inches(2.6), [
    ("Drivers (exact linear attribution):", True, NAVY),
    "Within-week noise 2.07x own baseline -> +1.202 (toward failure)",
    "Rest-voltage floor -3.15 V vs baseline -> +0.957 (battery floor sagging)",
    "Crank dips +4.12 V deeper (last 90 d) -> +0.579 (cascade signature)",
    "Range trend +0.200 V/wk -> -1.084 (suppressor term, flagged)",
    ("Counterfactual: drops RED -> AMBER if within-week noise ratio falls from 2.07x to 1.01x, all else equal.", True, DARK_TEXT),
    ("Routing: A2 cascade -> battery-first inspection (DICV A6).", True, GREEN_PASS),
], font_size=9, bullet=True)

add_key_takeaways(s, [
    "Every alert ships with its reasons: a depot engineer sees which physical channel moved and by how much in volts",
    "Counterfactuals make tiers actionable: 'returns to GREEN if the rest floor recovers 0.4 V' is a checkable statement",
    "Governance: model card with banned-feature registry, gate results, horizon curve and known-miss class (A4)",
], top=Inches(4.85))
add_footer(s)


# ─── SLIDE 22: CLOSED NEGATIVES ─────────────────────────────
print("Building slide 22: Closed Negatives")
s = blank_slide(prs)
add_header_bar(s, "CLOSED DEAD ENDS -- WITH NUMBERS",
               "Every shelved method now has evidence attached, not vibes")

neg_data = [
    ['Method family', 'What was tried', 'Result', 'Verdict'],
    ['Survival / hazard RUL',
     'Discrete-time hazard, Cox, Weibull AFT',
     f'Truck ranking 0.586 vs 0.893 (classifier); RUL MAE {SURV_RUL_MAE} d vs {CONST_RUL_MAE} d for a constant',
     'CLOSED -- a constant beats every survival formulation at 14 events'],
    ['Day-precision RUL',
     'All of the above + calibration analysis',
     'Calibrated weekly hazard (~0.005/wk) and day-precision are mathematically incompatible at 14 events',
     'CLOSED -- ship the 10-week horizon statement instead'],
    ['Deep sequence models',
     'LSTM/BiLSTM/TCN/Transformer/TFT/Informer/PatchTST/TimeXer, DeepSurv/DeepHit/DSM, VAE/contrastive/Siamese (sized)',
     '235x-6,275x over the parameter budget at 14 events; a 43-parameter LSTM\'s seed variance exceeds any signal difference',
     'CLOSED until n_failed >= 30-50'],
    ['Honest simple probes',
     'PCA, trend coefficients, distance methods',
     'All saturate at the same ~0.89-0.93 single degree of freedom the engineered features already capture',
     'CLOSED -- no untapped signal in the weekly aggregates'],
    ['GED emergency channel (ALT transfer)',
     'GED=2 storm detection on SM fleet',
     'Zero GED2 events in all 14 failed SM trucks',
     'CLOSED -- does not transfer; GED is a data-quality covariate only'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), 6, 4, neg_data,
                col_widths=[2.2, 3.6, 4.2, 2.5])

add_key_takeaways(s, [
    f"Survival modeling made things measurably worse: RUL MAE {SURV_RUL_MAE} d vs {CONST_RUL_MAE} d for a constant -- the proof that 'no RUL' is the honest deliverable",
    "Deep learning is closed by arithmetic, not opinion: every requested architecture is hundreds-to-thousands of times over the events-per-parameter budget",
    "One shelved future path: self-supervised crank-encoder pretraining on the 106M raw rows, defensible at n_failed >= 30-50",
], top=Inches(4.45))
add_footer(s)


# ─── SLIDE 23: CURRENT LIMITATIONS ──────────────────────────
print("Building slide 23: Current Limitations")
s = blank_slide(prs)
add_header_bar(s, "CURRENT LIMITATIONS",
               "Honest assessment of what the data and models cannot do")

limit_data = [
    ['Limitation', 'Impact', 'Root Cause', 'Mitigation'],
    ['A4 silent/abrupt failures are invisible',
     '~4/14 failures (29%) fire nothing on any layer',
     'Windings, seizure, control loss, telemetry death -- no electrical precursor at 5 s',
     'Transmission-health monitoring: a truck going quiet is itself a maintenance trigger'],
    ['No operating point dominates V1 everywhere',
     'Youden: 13/14 @ 15/20; RED tier: 10/14 @ 18/20',
     'Recall-specificity trade at n=34',
     'Two pre-registered points; operations choose per maintenance economics'],
    ['5 s / 0.2 V sampling destroys brush-wear channel',
     'The one genuine 60-120 d precursor is unreadable',
     'Sub-sample dip shape and sub-second duration growth',
     'High-frequency crank logging (>=1 Hz during SMA=1) -- the single biggest unlock'],
    ['n=34 trucks, 14 events',
     f'Wide CI [{BOOTSTRAP_CI[0]}, {BOOTSTRAP_CI[1]}]; thresholds decided by a handful of trucks',
     'Real-world failure-data scarcity',
     'Keep collecting failures; archetypes stay suggestive, not inferential'],
    ['OOF scores correlate with leak axes',
     'Spearman -0.64 (n_weeks), +0.51 (t_start) -- above tripwire',
     'Label-mediated: failed trucks transmit less because they failed',
     'Zero-drop L40 control + decay-to-chance curve; a larger fleet is the only definitive cure'],
    ['Persistence-flag NF false positives',
     '4/20 NF trucks flagged at end of history',
     'Right-censored degraders -- or rule drift',
     'Track VIN2/5/8/15_NF; either outcome is informative'],
    ['Battery-vs-starter ambiguity caps precision',
     'A2 routes battery-first but cannot prove the component',
     'No cranking current or battery SoC/SoH signal',
     'Instrument current or SoC -- ends the ambiguity'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), 8, 4, limit_data,
                col_widths=[3.0, 3.0, 3.3, 3.2])

add_key_takeaways(s, [
    "These are physics and data-size constraints, stated with numbers -- the honesty is itself a deliverable",
    "The 10-week horizon is the ceiling of this telemetry; reliable improvement requires new signals, not new models",
], top=Inches(6.3), height=Inches(0.6))
add_footer(s)


# ─── SLIDE 24: DEPLOYMENT, GOVERNANCE & ROADMAP ─────────────
print("Building slide 24: Deployment, Governance & Roadmap")
s = blank_slide(prs)
add_header_bar(s, "DEPLOYMENT, GOVERNANCE & ROADMAP",
               "Operating cadence, refit rules, and ranked data asks")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.9), Inches(0.3),
         "DEPLOYMENT & GOVERNANCE", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(5.9), Inches(2.6), [
    "Score the fleet weekly (features need 30-90-day windows; the 10-week horizon is built on weekly scoring); review tiers monthly",
    "RED -> starter + battery circuit inspection within 2-4 weeks; AMBER -> bundle into next scheduled service; GREEN -> normal operation",
    "Refit only when new failure labels arrive, always under the full nested protocol with all gates re-run and the permutation test",
    "Banned-feature registry is binding (vsi_dominant_freq, calendar features, observation-length features)",
    "If calibration slope leaves [0.5, 2] on a refit: ship tiers only, not probabilities",
    "Alert channels remain tier-gated -- no standalone pagers",
], font_size=10, bullet=True)

add_text(s, Inches(6.9), Inches(1.15), Inches(6.0), Inches(0.3),
         "RANKED DATA ASKS (VALUE PER RUPEE)", font_size=14, bold=True, color=NAVY)
ask_data = [
    ['#', 'Ask', 'What it unlocks'],
    ['1', 'High-frequency crank logging (>=1 Hz during SMA=1)', 'Revives the 60-120 d brush-wear channel and true dip physics'],
    ['2', 'Cranking current or battery SoC/SoH signal', 'Ends battery-vs-starter ambiguity; sharpens alert precision'],
    ['3', 'Maintenance / parts-replacement records', 'Turns data-derived archetypes into supervised labels'],
    ['4', 'Keep collecting failures (target n_failed >= 30-50)', 'Unlocks SSL crank-encoder pretraining and meaningful survival modeling'],
    ['5', 'Ambient temperature', 'Cold-start conditioning of crank features'],
]
add_table_shape(s, Inches(6.9), Inches(1.55), Inches(6.0), 6, 3, ask_data,
                col_widths=[0.4, 2.8, 2.8])

add_key_takeaways(s, [
    "V1.1 is the honest ceiling of this dataset -- improvement beyond it requires new signals, not new models",
    "Stop doing: chasing longer lead times in this data; evaluating deep/sequence/survival models below n_failed = 30",
    "VIN9_F-class failures need transmission-health monitoring, not a better classifier",
], top=Inches(5.45))
add_footer(s)


# ─── SLIDE 25: CONCLUSIONS ──────────────────────────────────
print("Building slide 25: Conclusions")
s = blank_slide(prs)
add_header_bar(s, "CONCLUSIONS & RECOMMENDATIONS",
               "Summary of findings and priority actions")

add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.3),
         "SUMMARY OF FINDINGS", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.0), [
    (f"Nested-LOVO Ridge (AUROC {AUROC_V11}, CI [{BOOTSTRAP_CI[0]}, {BOOTSTRAP_CI[1]}], p={PERMUTATION_P}) beats the honestly-restated V1 baseline ({AUROC_V1_REST}) under a strictly harder protocol", False, DARK_TEXT),
    (f"Calibrated probabilities ship for the first time (slope {CAL_SLOPE}, Brier {BRIER}); recall {RECALL_YOUDEN} including V1's worst miss (VIN8_F)", False, DARK_TEXT),
    (f"Three validated alert channels give {COMBINED_RECALL} early-warning coverage with a measured {K_STAR}-week horizon -- V1 had none", False, DARK_TEXT),
    ("A2 battery-cascade detector (0/20 NF false alarms, ~9.5-week lead) builds battery-vs-starter triage into the work order", False, DARK_TEXT),
    (f"Day-precision RUL is closed with proof (hazard MAE {SURV_RUL_MAE} d vs {CONST_RUL_MAE} d for a constant); the deliverable is the validity-horizon statement", False, DARK_TEXT),
    (f"~4/14 silent/abrupt failures are structurally invisible -- the honest recall ceiling for tier alerting is ~10-11/14", False, DARK_TEXT),
], font_size=11, bullet=True)

add_text(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.3),
         "PRIORITY RECOMMENDATIONS", font_size=16, bold=True, color=NAVY)
rec_data = [
    ['#', 'Recommendation', 'Rationale'],
    ['1', 'Inspect VIN5_NF_SM and VIN20_NF_SM (RED tier)', 'Recalibrated risk 0.96 / 0.62; starter + battery circuit within 2-4 weeks'],
    ['2', 'Stand up weekly fleet scoring with tier review monthly', 'The 10-week horizon requires weekly reads to be actionable'],
    ['3', 'Enable A2 battery-cascade alert with battery-first routing', '0/20 NF false alarms; the single largest avoidable-cost lever (DICV A6)'],
    ['4', 'Track persistence-flag NF trucks (VIN2/5/8/15_NF)', 'Future failures or first evidence of rule drift -- either is informative'],
    ['5', 'Procure high-frequency crank logging + battery SoC/current', 'The two ranked data asks that raise the ceiling'],
    ['6', 'Keep the banned-feature registry and nested protocol binding', 'Leak ceilings (AUROC 0.95 from data volume alone) exceed any honest model'],
    ['7', 'Communicate the honest picture to stakeholders', 'Tiers + horizon + triage, not failure dates; A4 blind spot stated in every deployment doc'],
]
add_table_shape(s, Inches(0.4), Inches(4.0), Inches(12.5), 8, 3, rec_data,
                col_widths=[0.5, 5.0, 7.0])

add_footer(s)


# ─── SLIDE 26: APPENDIX ─────────────────────────────────────
print("Building slide 26: Appendix")
s = blank_slide(prs)
add_header_bar(s, "APPENDIX & TECHNICAL REFERENCE",
               "Key terms, data sources, and pipeline specifications")

terms = [
    ['Term', 'Definition'],
    ['AUROC', 'Area Under ROC Curve -- measures classifier discrimination (1.0 = perfect)'],
    ['Nested LOVO', 'Leave-One-Vehicle-Out CV with feature screening, subset selection, threshold and recalibration redone inside every fold'],
    ['Ridge Classifier', 'L2-regularized linear classifier for binary failure-pattern prediction'],
    ['Platt Recalibration', 'Logistic mapping of decision values to calibrated probabilities, fitted per fold on inner out-of-fold values'],
    ['Youden’s J', 'Threshold maximizing Sensitivity + Specificity - 1, chosen per fold on inner-OOF predictions (pre-registered)'],
    ['VSI', 'Power Supply Voltage (CAN signal): 24V system; per-truck regulation setpoints span 27.6-28.2 V'],
    ['SMA', 'Starter Motor Active (CAN signal): 0=Not operated, 1=Operated -- basis of the crank-event catalog'],
    ['Archetype (A1-A4)', 'Data-derived failure pathway: A1 solenoid intermittency, A2 battery cascade, A3 volatility drift, A4 silent/abrupt'],
    ['Persistence flag', 'Causal within-week VSI-std ratio above training-fold NF p90 envelope in >=4 of last 12 weeks'],
    ['A2 triple detector', 'Rest-VSI step down + drive-VSI step up + crank-dip widening, all causal -- battery-cascade signature'],
    ['Detection horizon k*', 'Largest k (weeks of data removed before end of history) with sustained AUROC >= 0.75'],
    ['L40 control', 'Full rerun with every VIN clipped to its last 40 masked weeks -- zero drop proves no length artifact'],
    ['SMA-dead cohort', 'Telematics config with <1% SMA coverage (7 trucks) -- crank features NaN + fold-internal median impute'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), len(terms), 2, terms,
                col_widths=[2.5, 10.0])

add_multiline(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5), [
    "Data sources: STARTER MOTOR/V1.1/results/ (model spec, gates, nested predictions, horizon curve, alert validation), V1.1/reports/ (model card, comparison, alerts+horizon, experiment results), STARTER MOTOR/reports/V1_SM_final_report.md",
    "All metrics are nested-LOVO out-of-fold unless stated. See docs/column_dictionary.md for signal definitions. SM VINs are independent of ALT VINs (different physical trucks).",
], font_size=8, bullet=True, color=GREY_MED)
add_footer(s)


# ── SAVE ───────────────────────────────────────────────────────
out_path = str(OUT_DIR / "SM_Predictive_Maintenance_V1.1.pptx")
prs.save(out_path)
print(f"\nPresentation saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")

shutil.rmtree(TMPDIR, ignore_errors=True)
print("Done.")
